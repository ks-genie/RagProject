# -*- coding: utf-8 -*-
"""
RAG Pipeline Manager — Streamlit UI (Phase 9)

이 파일은 RAG(Retrieval-Augmented Generation) 파이프라인의
웹 관리 화면을 Streamlit으로 구현한 진입점입니다.

실행:
    streamlit run app.py
"""

import os
import threading
import time
from collections import Counter, defaultdict
from datetime import datetime
from io import BytesIO
from pathlib import Path

import numpy as np
import streamlit as st

from src.config import config as app_config
from src.database import DatabaseManager, Status, db
from src.log_setup import setup_logging, get_log_files
from src.ocr_postprocess import fix_easyocr_korean
from src.page_db import page_db
from src.pipeline import pipeline
from src.preprocessor import normalize_hard_linebreaks
from src.strategy_selector import STRATEGY_TESSERACT_CONFIG

# ---------------------------------------------------------------------------
# 로깅 초기화 — 앱 시작 시 1회, logs/log_yyyymmdd_hhmmss.log 생성
# setup_logging()은 날짜/시각이 포함된 로그 파일 경로를 반환합니다.
# ---------------------------------------------------------------------------
_LOG_PATH = setup_logging()

# ---------------------------------------------------------------------------
# 앱 버전 / 수정 시간
# os.path.getmtime()으로 이 파일의 마지막 수정 타임스탬프를 가져옵니다.
# ---------------------------------------------------------------------------
APP_VERSION  = "1.66"
_mtime       = os.path.getmtime(__file__)  # 현재 파일(app.py)의 수정 시각(초 단위 타임스탬프)
APP_MODIFIED = datetime.fromtimestamp(_mtime).strftime("%Y-%m-%d %H:%M:%S")

# ---------------------------------------------------------------------------
# 페이지 설정
# Streamlit 앱의 탭 제목, 아이콘, 레이아웃 등을 설정합니다.
# layout="wide"로 설정하면 화면 전체 너비를 사용합니다.
# ---------------------------------------------------------------------------
st.set_page_config(
    page_title="RAG Pipeline Manager",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------
# 전역 CSS — 문서 행 compact 스타일
# unsafe_allow_html=True 옵션을 사용해 HTML/CSS를 직접 주입합니다.
# Streamlit 기본 여백이 너무 넓기 때문에 UI를 더 촘촘하게 만들기 위한 조정입니다.
# ---------------------------------------------------------------------------
st.markdown("""
<style>
/* ── 페이지 상단 여백 축소 ── */
.block-container {
    padding-top: 0.75rem !important;
}

/* ── 탭 패널 내부만 행 간격 압축 (타이틀 영역 영향 없음) ── */
[data-testid="stTabPanel"] > div > [data-testid="stVerticalBlock"],
[role="tabpanel"] > div > [data-testid="stVerticalBlock"] {
    gap: 0 !important;
    row-gap: 0 !important;
}
[data-testid="stTabPanel"] .element-container,
[role="tabpanel"] .element-container {
    margin-top:    0 !important;
    margin-bottom: 0 !important;
}

/* ── 문서 행 컬럼 수직 정렬 ── */
[data-testid="stTabPanel"] [data-testid="stHorizontalBlock"],
[role="tabpanel"] [data-testid="stHorizontalBlock"] {
    align-items: center !important;
}
/* 컬럼 내부 stVerticalBlock도 세로 중앙 정렬 */
[data-testid="stTabPanel"] [data-testid="stHorizontalBlock"] [data-testid="stVerticalBlock"],
[role="tabpanel"] [data-testid="stHorizontalBlock"] [data-testid="stVerticalBlock"] {
    justify-content: center !important;
}
/* 컬럼 내부 상하 패딩 제거 */
[data-testid="stTabPanel"] [data-testid="stHorizontalBlock"] [data-testid="column"] > div,
[role="tabpanel"] [data-testid="stHorizontalBlock"] [data-testid="column"] > div {
    padding-top: 0 !important;
    padding-bottom: 0 !important;
}

/* ── 문서 행(bordered container) 내부 compact ── */

/* wrapper 자체 마진 */
[data-testid="stVerticalBlockBorderWrapper"] {
    margin-top:    0 !important;
    margin-bottom: 0 !important;
}
[data-testid="stVerticalBlockBorderWrapper"] > div {
    padding-top:    1px !important;
    padding-bottom: 1px !important;
    padding-left:   6px !important;
    padding-right:  6px !important;
}

/* bordered 내부 stVerticalBlock gap 0 */
[data-testid="stVerticalBlockBorderWrapper"] [data-testid="stVerticalBlock"] {
    gap: 0 !important;
}

/* bordered 내부 element-container 완전 압축 */
[data-testid="stVerticalBlockBorderWrapper"] .element-container {
    padding-top:   0 !important;
    padding-bottom:0 !important;
    min-height:    0 !important;
    margin-top:    0 !important;
    margin-bottom: 0 !important;
}

/* 내부 horizontal block */
[data-testid="stVerticalBlockBorderWrapper"] [data-testid="stHorizontalBlock"] {
    align-items: center !important;
    gap: 4px !important;
}

/* column 내부 패딩 제거 */
[data-testid="stVerticalBlockBorderWrapper"] [data-testid="column"] > div {
    padding: 0 !important;
    min-height: 0 !important;
}

/* 토글 스위치 컬럼 — 상단 정렬 */
[data-testid="stVerticalBlockBorderWrapper"] [data-testid="column"]:first-child {
    align-self: flex-start !important;
    padding-top: 1px !important;
}

/* 마크다운 여백·행높이 압축 */
[data-testid="stVerticalBlockBorderWrapper"] .stMarkdown p {
    margin: 0 !important;
    line-height: 1.25 !important;
    font-size: 0.875em;
}

/* caption 여백·행높이 압축 */
[data-testid="stVerticalBlockBorderWrapper"] [data-testid="stCaptionContainer"] p {
    margin: 0 !important;
    line-height: 1.1 !important;
}

/* 액션 버튼(재처리 등) — compact 유지 */
[data-testid="stVerticalBlockBorderWrapper"] .stButton > button {
    height:     22px !important;
    min-height: 22px !important;
    padding:    0 4px !important;
    font-size:  0.75em !important;
    line-height: 1 !important;
}

</style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------------------------
# 상태 색상 / 이모지
# 문서 처리 상태(Status)별로 화면에 표시할 이모지와 배경 색상을 정의합니다.
# Status 값은 src/database.py 의 Status 클래스에 정의되어 있습니다.
# ---------------------------------------------------------------------------
STATUS_EMOJI = {
    Status.NEW:              "🆕",
    Status.ANALYZING:        "🔍",
    Status.TEXT_EXTRACTED:   "📝",
    Status.OCR_DONE:         "📝",
    Status.TEXT_READY:       "✅",
    Status.UPLOAD_REQUESTED: "⬆️",
    Status.UPLOADED:         "☁️",
    Status.PARSING_DONE:     "⚙️",
    Status.CHUNKING_DONE:    "🔧",
    Status.EMBEDDING_DONE:   "🧩",
    Status.INDEXED:          "🎯",
    Status.REVIEW_REQUIRED:  "⚠️",
    Status.FAILED:           "❌",
}

# 상태별 배지 배경색 (Bootstrap 계열 색상 코드)
STATUS_COLOR = {
    Status.NEW:              "#6c757d",  # 회색 — 아직 처리 시작 전
    Status.ANALYZING:        "#0d6efd",  # 파란색 — 분석 중
    Status.TEXT_EXTRACTED:   "#0dcaf0",  # 하늘색 — 텍스트 추출 완료
    Status.OCR_DONE:         "#0dcaf0",  # 하늘색 — OCR 완료
    Status.TEXT_READY:       "#198754",  # 초록색 — 텍스트 준비 완료
    Status.UPLOAD_REQUESTED: "#ffc107",  # 노란색 — 업로드 요청됨
    Status.UPLOADED:         "#20c997",  # 청록색 — 업로드 완료
    Status.PARSING_DONE:     "#20c997",  # 청록색 — 파싱 완료
    Status.CHUNKING_DONE:    "#20c997",  # 청록색 — 청킹 완료
    Status.EMBEDDING_DONE:   "#20c997",  # 청록색 — 임베딩 완료
    Status.INDEXED:          "#198754",  # 초록색 — 인덱싱 완료 (최종 성공 상태)
    Status.REVIEW_REQUIRED:  "#fd7e14",  # 주황색 — 검토 필요
    Status.FAILED:           "#dc3545",  # 빨간색 — 처리 실패
}


def status_badge(status: str) -> str:
    """문서 상태를 색깔 있는 HTML 배지(badge) 문자열로 반환합니다.

    매개변수:
        status: 표시할 문서 처리 상태 문자열 (예: "INDEXED", "FAILED")

    반환값:
        HTML <span> 태그로 감싼 색깔 배지 문자열. st.markdown(..., unsafe_allow_html=True)로 표시해야 합니다.
    """
    color = STATUS_COLOR.get(status, "#6c757d")  # 알 수 없는 상태는 회색으로 처리
    emoji = STATUS_EMOJI.get(status, "•")
    return (
        f'<span style="background:{color};color:white;padding:2px 8px;'
        f'border-radius:12px;font-size:0.8em;font-weight:bold;">'
        f'{emoji} {status}</span>'
    )


# ---------------------------------------------------------------------------
# 세션 상태 초기화
# Streamlit은 매 상호작용마다 스크립트 전체를 재실행합니다.
# st.session_state를 이용하면 새로고침 사이에도 값을 유지할 수 있습니다.
# 키가 없을 때만 초기값을 설정합니다 (이미 있으면 덮어쓰지 않습니다).
# ---------------------------------------------------------------------------
if "pipeline_running" not in st.session_state:
    st.session_state.pipeline_running = False   # 파이프라인 실행 중 여부
if "pipeline_log" not in st.session_state:
    st.session_state.pipeline_log = ""          # 파이프라인 실행 결과 로그 텍스트
if "last_refresh" not in st.session_state:
    st.session_state.last_refresh = time.time() # 마지막 데이터 새로고침 시각(초 단위)


# ---------------------------------------------------------------------------
# 데이터 로드 (TTL 캐시)
# @st.cache_data(ttl=30): 30초 동안 결과를 캐시합니다.
# 같은 인수로 30초 이내에 다시 호출하면 DB를 조회하지 않고 캐시된 값을 반환합니다.
# ---------------------------------------------------------------------------
@st.cache_data(ttl=30)
def load_documents():
    """DB에서 전체 문서 목록을 불러옵니다.

    반환값:
        dict 리스트. 각 dict는 documents 테이블의 한 행(id, file_name, status 등)입니다.
    """
    return [dict(r) for r in db.get_all_documents()]


@st.cache_data(ttl=30)
def load_all_logs() -> dict:
    """전체 처리 로그를 단일 쿼리로 로드합니다.

    반환값:
        {doc_id: [log_dict, ...]} 형태의 딕셔너리.
        doc_id별로 관련 로그 목록이 묶여 있습니다.
    """
    return db.get_all_logs_grouped()


@st.cache_data(ttl=30)
def load_asset_counts() -> dict:
    """전체 자산(표·그림) 건수를 단일 쿼리로 로드합니다.

    반환값:
        {doc_id: {"name": 파일명, "counts": {"TABLE": N, "FIGURE": M}}} 형태의 딕셔너리.
        자산이 하나도 없는 문서는 결과에서 제외됩니다.
    """
    from src.asset_db import asset_db as _adb
    all_counts = _adb.count_all_assets()
    # 문서 ID → 파일명 매핑 테이블 생성
    docs_map   = {d["id"]: d["file_name"] for d in load_documents()}
    return {
        doc_id: {"name": docs_map.get(doc_id, f"doc_{doc_id}"), "counts": counts}
        for doc_id, counts in all_counts.items()
        if counts  # 자산이 있는 문서만 포함
    }


def refresh_data():
    """캐시를 전부 비워 다음 접근 시 DB에서 새로 불러오게 합니다.

    Streamlit의 @st.cache_data는 .clear()로 해당 함수의 캐시를 초기화할 수 있습니다.
    """
    load_documents.clear()
    load_all_logs.clear()
    load_asset_counts.clear()
    st.session_state.last_refresh = time.time()  # 새로고침 시각 업데이트


# ---------------------------------------------------------------------------
# PDF 뷰어 — 캐시 헬퍼
# ---------------------------------------------------------------------------

def _legacy_words_to_text_with_paragraphs(words: list[dict]) -> str:
    """단어 목록으로 문단 구조를 보존한 텍스트를 조립합니다. (구 버전 — 현재는 OcrEngine으로 대체)

    동작 방식:
    1) 단어를 (top, x0) 기준으로 정렬한 뒤, y 좌표 차이가 LINE_TOL(5pt) 이내인 단어들을 같은 줄로 묶습니다.
    2) 줄 간격 중앙값 × 1.3 초과 → 빈 줄(문단 구분), 이하 → 줄바꿈으로 처리합니다.
    DIGITAL_EXTRACT / DIGITAL_EXTRACT_MULTI 양쪽에서 공통 사용합니다.

    매개변수:
        words: pdfplumber extract_words()가 반환하는 단어 딕셔너리 목록.
               각 dict에는 "top", "x0", "text" 키가 포함됩니다.

    반환값:
        문단 구조가 보존된 텍스트 문자열.
    """
    if not words:
        return ""

    # top(y좌표) 기준으로 먼저 정렬하고, 같은 줄에서는 x0(x좌표) 기준으로 정렬
    ws = sorted(words, key=lambda w: (float(w["top"]), float(w["x0"])))

    LINE_TOL = 5  # y 좌표 차이가 이 픽셀 이하면 같은 줄로 취급
    lines: list[list[dict]] = []
    cur: list[dict] = [ws[0]]
    cur_top = float(ws[0]["top"])
    for w in ws[1:]:
        w_top = float(w["top"])
        if abs(w_top - cur_top) <= LINE_TOL:
            cur.append(w)  # 같은 줄에 추가
        else:
            lines.append(cur)
            cur = [w]
            cur_top = w_top
    lines.append(cur)  # 마지막 줄 추가

    # 각 줄을 (y좌표, 텍스트) 튜플로 변환
    line_data: list[tuple[float, str]] = []
    for ln in lines:
        ln_s = sorted(ln, key=lambda w: float(w["x0"]))  # x 좌표 순으로 단어 배열
        line_data.append((float(ln_s[0]["top"]),
                          " ".join(w["text"] for w in ln_s)))
    line_data.sort(key=lambda x: x[0])

    tops  = [t for t, _ in line_data]
    texts = [t for _, t in line_data]

    if len(tops) < 2:
        return "\n".join(texts)

    # 인접 줄들의 간격 목록을 구하고 중앙값으로 문단 경계 임계값 결정
    gaps = [tops[i + 1] - tops[i] for i in range(len(tops) - 1)]
    median_gap = sorted(gaps)[len(gaps) // 2]
    threshold  = median_gap * 1.3  # 이 값보다 간격이 크면 문단 경계로 판단

    parts: list[str] = [texts[0]]
    for i, gap in enumerate(gaps):
        parts.append(("\n\n" if gap > threshold else "\n") + texts[i + 1])

    return "".join(parts)


def _words_to_text_with_paragraphs(words: list[dict]) -> str:
    """단어 목록을 문단 구조가 보존된 텍스트로 변환합니다. (OcrEngine 위임 버전)

    매개변수:
        words: pdfplumber extract_words() 반환값.

    반환값:
        문단 구조가 보존된 텍스트 문자열.
    """
    from src.ocr_engine import OcrEngine

    # 실제 변환 로직은 OcrEngine 클래스 메서드에 위임합니다.
    return OcrEngine._words_to_text_with_paragraphs(words, line_height=5)


# ---------------------------------------------------------------------------
# 헤더·푸터 분리 헬퍼
# PDF 페이지의 상단 8%는 헤더, 하단 7%는 푸터로 간주합니다.
# ---------------------------------------------------------------------------

_HEADER_PCT = 0.08   # 페이지 상단 8% 영역 → 헤더로 분류
_FOOTER_PCT = 0.93   # 페이지 하단 7% 영역 (93% 이후) → 푸터로 분류


def _split_header_footer(
    page, words: list[dict]
) -> tuple[list[dict], list[dict], list[dict]]:
    """단어 목록을 헤더 / 본문 / 푸터 영역으로 분리합니다.

    헤더: top < page_height × 0.08
    푸터: top > page_height × 0.93

    매개변수:
        page: pdfplumber Page 객체 (page.height 를 사용합니다).
        words: extract_words() 반환값.

    반환값:
        (header_words, main_words, footer_words) 튜플.
    """
    h = float(page.height)
    h_th = h * _HEADER_PCT  # 헤더 경계 y 좌표
    f_th = h * _FOOTER_PCT  # 푸터 경계 y 좌표
    header = [w for w in words if float(w["top"]) < h_th]
    footer = [w for w in words if float(w["top"]) > f_th]
    main   = [w for w in words if h_th <= float(w["top"]) <= f_th]
    return header, main, footer


def _words_one_line(words: list[dict]) -> str:
    """단어 목록을 정렬해 한 줄 문자열로 이어붙입니다. (헤더·푸터 표시용)

    여러 줄이 있을 경우 줄 구분자로 " | "를 사용합니다.

    매개변수:
        words: pdfplumber extract_words() 반환값.

    반환값:
        단어들을 한 줄로 이어붙인 문자열.
    """
    if not words:
        return ""
    ws = sorted(words, key=lambda w: (float(w["top"]), float(w["x0"])))
    LINE_TOL = 5
    lines: list[list[dict]] = []
    cur = [ws[0]]; cur_top = float(ws[0]["top"])
    for w in ws[1:]:
        if abs(float(w["top"]) - cur_top) <= LINE_TOL:
            cur.append(w)
        else:
            lines.append(cur); cur = [w]; cur_top = float(w["top"])
    lines.append(cur)
    # 각 줄을 공백으로 묶고, 줄 사이는 " | "로 구분
    return " | ".join(
        " ".join(w["text"] for w in sorted(ln, key=lambda w: float(w["x0"])))
        for ln in lines
    )


def _connect_columns(left_text: str, right_text: str) -> tuple[str, str]:
    """좌→우 컬럼 경계에서 이어지는 문장·문단을 연결합니다.

    판단 기준 (영문 기준):
      - 좌 컬럼 마지막 단락이 문장 종결 기호(.?!;)로 끝나지 않음
      - 우 컬럼 첫 단락이 소문자로 시작함
    두 조건 모두 충족되면 하이픈 처리 포함하여 이어붙입니다.

    매개변수:
        left_text: 좌측 컬럼 텍스트.
        right_text: 우측 컬럼 텍스트.

    반환값:
        연결 처리된 (new_left_text, new_right_text) 튜플.
    """
    SENTENCE_END = frozenset('.?!;')
    left_paras  = [p for p in left_text.strip().split('\n\n')  if p.strip()]
    right_paras = [p for p in right_text.strip().split('\n\n') if p.strip()]
    if not left_paras or not right_paras:
        return left_text, right_text

    last_left   = left_paras[-1].rstrip()
    first_right = right_paras[0].lstrip()
    if not last_left or not first_right:
        return left_text, right_text

    last_char  = last_left[-1]
    first_char = first_right[0]

    # 좌측 마지막 줄이 종결 부호로 끝나지 않고, 우측 첫 줄이 소문자로 시작하면 연결
    if last_char not in SENTENCE_END and first_char.islower():
        if last_left.endswith('-'):
            connector = ""; last_left = last_left[:-1]  # 하이픈으로 끝나면 하이픈 제거 후 붙임
        else:
            connector = " "
        merged    = last_left + connector + first_right
        new_left  = '\n\n'.join(left_paras[:-1] + [merged])
        new_right = '\n\n'.join(right_paras[1:])
        return new_left, new_right

    return left_text, right_text


def _extract_academic_header(page, gap_x: float | None = None) -> str | None:
    """학술 논문 페이지에서 제목·저자를 폰트 크기로 감지해 레이블 분리합니다.

    gap_x 가 주어지면 좌측 컬럼 chars 만으로 상태 기계를 실행합니다.
    이 경우 적응형 임계값(비하드탑 영역 최대 폰트 크기 기준)을 사용하며,
    제목·저자가 좌측 컬럼 전용인지 전체 폭인지를 판별해 본문 필터를 조정합니다.

    동작 방식:
    ① page.chars 로 폰트 크기 기반 줄 분류 (상태 기계)
      - gap_x 없음: title_thresh = body_size × 1.5
      - gap_x 있음: HARD_TOP(4%) 위 chars 제외 + 적응형 title_thresh
    ② 실제 텍스트 조립은 extract_words() 사용 (공백 처리 정확)
    감지 실패 시 None 반환 → 호출자가 일반 추출로 폴백합니다.

    매개변수:
        page: pdfplumber Page 객체.
        gap_x: 다단 레이아웃에서 컬럼 사이 갭의 x 좌표. 단단 레이아웃이면 None.

    반환값:
        "[헤더]\\n...\\n\\n[제목]\\n...\\n\\n[저자]\\n..." 형식의 문자열.
        감지 실패 시 None.
    """
    from collections import Counter

    # ── 기본 설정 ──────────────────────────────────────────────────────────
    # page.chars: 페이지의 모든 개별 문자 정보 (위치, 폰트 크기 포함)
    chars = [c for c in page.chars if (c.get("text") or "").strip()]
    if not chars:
        return None

    page_h     = float(page.height)
    page_w     = float(page.width)
    page_mid_x = page_w / 2  # 페이지 중앙 x 좌표 (소속 기관 감지에 사용)
    LINE_TOL   = 5

    # 전체 chars 기준 body_size (최빈 폰트 크기 = 본문 폰트 크기)
    size_freq = Counter(
        round(float(c.get("size") or 0), 1)
        for c in chars if float(c.get("size") or 0) > 0
    )
    if len(size_freq) < 2:
        return None  # 폰트 종류가 1개 이하면 제목/본문 구분 불가
    body_size = size_freq.most_common(1)[0][0]  # 가장 많이 등장한 폰트 크기

    # ── 상태 기계용 chars 및 임계값 결정 ────────────────────────────────
    if gap_x is not None:
        # 좌측 컬럼 chars 만 사용 (다단 레이아웃)
        analysis_chars = [c for c in chars if float(c.get("x0", 0)) < gap_x]
        # 페이지 절대 상단 4% 는 세션 레이블 등 페이지 헤더로 간주하고 제외
        HARD_TOP = page_h * 0.04
        nh_chars = [c for c in analysis_chars if float(c.get("top", 0)) >= HARD_TOP]
        if not nh_chars:
            return None
        nh_sizes = [round(float(c.get("size") or 0), 1)
                    for c in nh_chars if float(c.get("size") or 0) > 0]
        nh_max   = max(nh_sizes) if nh_sizes else body_size * 1.5
        # 비하드탑 영역 최대 폰트의 95% → 제목 임계값 (너무 작은 글자는 제목에서 제외)
        title_thresh = nh_max * 0.95 if nh_max > body_size * 1.1 else body_size * 1.5
    else:
        # 단단 레이아웃: 전체 chars 대상, 제목은 본문 1.5배 이상 크기
        analysis_chars = chars
        HARD_TOP       = 0.0
        title_thresh   = body_size * 1.5

    auth_thresh = body_size * 1.05  # 저자 폰트는 본문보다 약간 큰 크기

    # ── chars → 줄 묶기 ──────────────────────────────────────────────────
    sc = sorted(analysis_chars, key=lambda c: (float(c["top"]), float(c["x0"])))
    if not sc:
        return None
    raw: list[list] = [[sc[0]]]
    for c in sc[1:]:
        # 이전 줄 마지막 문자와 y 좌표 차이가 LINE_TOL 이하면 같은 줄
        if abs(float(c["top"]) - float(raw[-1][-1]["top"])) <= LINE_TOL:
            raw[-1].append(c)
        else:
            raw.append([c])

    # 각 줄의 평균 폰트 크기와 y 좌표를 계산
    char_lines = []
    for ln in raw:
        avg_sz = sum(float(c.get("size") or 0) for c in ln) / len(ln)
        char_lines.append({"top": float(ln[0]["top"]), "size": avg_sz})
    char_lines.sort(key=lambda x: x["top"])

    # ── 상태 기계: pre → title → author → body ───────────────────────────
    # 페이지 상단에서 내려오면서 폰트 크기에 따라 상태를 전환합니다.
    title_tops:  list[float] = []  # 제목 줄들의 y 좌표
    author_tops: list[float] = []  # 저자 줄들의 y 좌표
    state = "pre"  # 초기 상태: 제목 이전

    for li in char_lines:
        sz, top = li["size"], li["top"]
        if top < HARD_TOP:
            continue  # 페이지 절대 상단(세션 레이블 등) 제외
        if state == "pre":
            if sz >= title_thresh:
                state = "title"; title_tops.append(top)
        elif state == "title":
            if sz >= title_thresh:
                title_tops.append(top)      # 제목 줄 계속
            elif sz >= auth_thresh:
                state = "author"; author_tops.append(top)  # 저자 줄 시작
            # sz < auth_thresh → 공백·수퍼스크립트, title 상태 유지
        elif state == "author":
            if sz >= auth_thresh:
                author_tops.append(top)     # 저자 줄 계속
            elif sz >= body_size:
                state = "body"              # 본문 시작 → 탐색 종료
            # sz < body_size → 수퍼스크립트(각주 번호 등), 무시

    if not title_tops:
        return None  # 제목을 찾지 못하면 None 반환

    # ── 제목이 전체 폭인지 좌측 컬럼 전용인지 판별 ───────────────────────
    # gap_x 가 있을 때만 필요: 우측 컬럼에서 제목 y-범위 chars 의 평균 폰트를 확인
    full_width_title  = True
    full_width_author = True
    if gap_x is not None:
        title_top_min = min(title_tops) - LINE_TOL
        title_top_max = max(title_tops) + LINE_TOL
        # 우측 컬럼에서 제목 y 범위에 있는 문자들
        right_tc = [c for c in chars
                    if float(c.get("x0", 0)) >= gap_x
                    and title_top_min <= float(c.get("top", 0)) <= title_top_max
                    and float(c.get("size") or 0) > 0]
        if right_tc:
            right_avg = sum(float(c.get("size", 0)) for c in right_tc) / len(right_tc)
            # 우측에도 큰 폰트가 있으면 전체 폭 제목으로 판단
            full_width_title = right_avg >= title_thresh * 0.9
        else:
            full_width_title = False   # 우측 컬럼에 제목 y범위 chars 없음 → 좌측 전용
        full_width_author = full_width_title  # 제목이 좌측 전용이면 저자도 동일

    # ── extract_words()로 텍스트 조립 ────────────────────────────────────
    all_words  = page.extract_words()
    footer_thr = page_h * _FOOTER_PCT
    title_min  = min(title_tops)
    col_limit  = None if full_width_title else gap_x  # 전체 폭이면 제한 없음

    def words_in_range(min_t: float, max_t: float,
                       max_x: float | None = None) -> list[dict]:
        """지정한 y 범위와 x 범위 안에 있는 단어만 필터링합니다."""
        return [w for w in all_words
                if min_t - LINE_TOL <= float(w["top"]) <= max_t + LINE_TOL
                and (max_x is None or float(w["x0"]) < max_x)]

    # 제목 단어 수집
    t_words = words_in_range(title_min, max(title_tops), col_limit)
    if len(t_words) < 3:
        return None  # 단어가 3개 미만이면 제목 감지 실패로 처리

    # 제목 단어를 줄별로 묶어 텍스트 조립
    t_by_line: dict[int, list] = {}
    for w in t_words:
        key = round(float(w["top"]))
        t_by_line.setdefault(key, []).append(w)
    title_line_strs = [
        " ".join(w["text"] for w in sorted(vs, key=lambda w: float(w["x0"])))
        for _, vs in sorted(t_by_line.items())
    ]
    title_text = title_line_strs[0]
    for tx in title_line_strs[1:]:
        # 하이픈으로 끝나면 공백 없이 이어붙임 (단어 분리 하이픈 처리)
        title_text = (title_text + tx) if title_text.endswith("-") \
                     else (title_text + " " + tx)

    # 저자 단어 수집 및 조립
    a_col_limit = None if full_width_author else gap_x
    author_text = ""
    if author_tops:
        a_words = words_in_range(min(author_tops), max(author_tops), a_col_limit)
        a_by_line: dict[int, list] = {}
        for w in sorted(a_words, key=lambda w: float(w["top"])):
            w_top = float(w["top"])
            placed = False
            for key in a_by_line:
                if abs(w_top - key) <= LINE_TOL:
                    a_by_line[key].append(w); placed = True; break
            if not placed:
                a_by_line[round(w_top)] = [w]
        author_text = " ".join(
            " ".join(w["text"] for w in sorted(vs, key=lambda w: float(w["x0"])))
            for _, vs in sorted(a_by_line.items())
        ).strip()

    # ── 소속 기관 감지 ────────────────────────────────────────────────────
    # 페이지 하단 좌측에 얇은 구분선이 있으면 소속 기관 블록의 시작으로 간주
    aff_start_top: float | None = None
    for seg in page.lines:
        seg_w  = float(seg.get("width",  0))
        seg_h  = float(seg.get("height", 1))
        seg_y  = float(seg.get("top",    0))
        seg_x0 = float(seg.get("x0",     0))
        # 얇고(높이 < 2pt) 적당히 긴(30~페이지 절반 미만) 가로선이 하단 60% 이후에 있으면
        if (seg_h < 2 and seg_w > 30 and seg_w < page_mid_x
                and seg_x0 < page_mid_x and seg_y > page_h * 0.60):
            if aff_start_top is None or seg_y < aff_start_top:
                aff_start_top = seg_y

    if aff_start_top is None:
        # 구분선이 없으면 하단 좌측의 소폰트 문자들로 소속 블록 탐색
        AFF_FONT_MAX = body_size - 1.5
        left_small = [
            c for c in chars
            if float(c.get("x0", 0)) < page_mid_x
            and float(c.get("top", 0)) > page_h * 0.70
            and 0 < float(c.get("size", 0)) < AFF_FONT_MAX
        ]
        # 서브스크립트(1~2자) 오탐 방지 — 실제 소속 블록은 여러 줄에 걸쳐 있음
        if len(left_small) >= 8:
            aff_start_top = min(float(c["top"]) for c in left_small)

    aff_words: list[dict] = []
    if aff_start_top is not None:
        # 소속 기관 영역의 단어 수집 (왼쪽 절반 영역만)
        aff_words = [w for w in all_words
                     if float(w["top"]) >= aff_start_top - LINE_TOL
                     and float(w["x0"]) < page_mid_x]

    # ── in_skip: 본문에서 제외할 영역 판별 함수 ──────────────────────────
    # 제목·저자가 좌측 컬럼 전용이면 같은 y범위의 우측 컬럼(본문)은 제외하지 않음
    def in_skip(top: float, x0: float = 0.0) -> bool:
        """해당 y, x 좌표의 단어를 본문에서 제외해야 하면 True를 반환합니다."""
        if top < title_min - LINE_TOL:   return True   # 제목 위 헤더 영역
        if top > footer_thr:              return True   # 푸터 영역
        # 제목 영역
        if top <= max(title_tops) + LINE_TOL:
            if full_width_title:          return True
            if x0 < (gap_x or page_w):   return True   # 좌측 컬럼만 제외
        # 저자 영역
        if author_tops and (min(author_tops) - LINE_TOL
                            <= top <= max(author_tops) + LINE_TOL):
            if full_width_author:         return True
            if x0 < (gap_x or page_w):   return True   # 좌측 컬럼만 제외
        # 소속 영역 (항상 왼쪽 컬럼)
        if (aff_start_top is not None
                and top >= aff_start_top - LINE_TOL
                and x0 < page_mid_x):     return True
        return False

    # 본문 단어: 제목·저자·소속·헤더·푸터 영역을 제외한 단어들
    body_words = [w for w in all_words
                  if not in_skip(float(w["top"]), float(w["x0"]))]
    if body_words and gap_x is not None:
        # 다단 레이아웃이면 OcrEngine의 다단 조합 로직을 사용
        from src.ocr_engine import OcrEngine
        body_start_y = OcrEngine._find_body_start_y(body_words, gap_x)
        body_text = OcrEngine._compose_multicol_text(
            body_words, gap_x, body_start_y
        )
    else:
        body_text = _words_to_text_with_paragraphs(body_words) if body_words else ""

    # [헤더]: 제목 시작 y 위쪽 단어 (전체 폭)
    hdr_words = [w for w in all_words if float(w["top"]) < title_min - LINE_TOL]
    ftr_words = [w for w in all_words if float(w["top"]) > footer_thr]

    # 각 섹션을 레이블과 함께 조립
    parts: list[str] = []
    if hdr_words:
        hdr_str = _words_one_line(hdr_words)
        if hdr_str.strip():
            parts.append(f"[헤더]\n{hdr_str}")
    parts.append(f"[제목]\n{title_text}")
    if author_text:
        parts.append(f"[저자]\n{author_text}")
    if aff_words:
        aff_text = _words_to_text_with_paragraphs(aff_words)
        if aff_text.strip():
            parts.append(f"[소속]\n{aff_text}")
    if body_text.strip():
        parts.append(f"[본문]\n{body_text}")
    if ftr_words:
        ftr_str = _words_one_line(ftr_words)
        if ftr_str.strip():
            parts.append(f"[푸터]\n{ftr_str}")
    return "\n\n".join(parts)


def _plumber_page_to_text(page) -> str:
    """pdfplumber 페이지에서 문단 구조를 보존해 텍스트를 추출합니다.

    추출 시도 순서:
    1) 학술 논문 제목/저자 감지 (_extract_academic_header) → 성공 시 즉시 반환
    2) layout=True (pdfminer 레이아웃 분석) — 단락 사이 빈 줄 자동 삽입
    3) y좌표 기반 문단 감지 (_words_to_text_with_paragraphs)

    매개변수:
        page: pdfplumber Page 객체.

    반환값:
        추출된 텍스트 문자열.
    """
    import re

    # 1단계: 학술 논문 헤더 감지 시도
    header = _extract_academic_header(page)
    if header:
        return header

    # 2단계: pdfminer 레이아웃 분석 기반 추출
    try:
        text = page.extract_text(layout=True)
        if text and text.strip():
            text = re.sub(r"\n{3,}", "\n\n", text)   # 3개 이상 연속 줄바꿈 → 2개로 축소
            text = re.sub(r"[ \t]{2,}", " ", text)   # 연속 공백 → 1개 공백
            return text.strip()
    except Exception:
        pass

    # 3단계: 단어 위치 기반 폴백 추출
    words = page.extract_words()
    if not words:
        return "(추출된 텍스트 없음)"
    hdr_w, main_w, ftr_w = _split_header_footer(page, words)
    parts: list[str] = []
    if hdr_w:
        hdr = _words_one_line(hdr_w)
        if hdr.strip():
            parts.append(f"[헤더]\n{hdr}")
    parts.append(_words_to_text_with_paragraphs(main_w) or "(추출된 텍스트 없음)")
    if ftr_w:
        ftr = _words_one_line(ftr_w)
        if ftr.strip():
            parts.append(f"[푸터]\n{ftr}")
    return "\n\n".join(parts)


def _strip_inline_captions(text: str) -> str:
    """뷰어 표시용: 페이지 텍스트 내 Fig./Table 캡션을 제거하고 분리된 문장을 연결합니다.

    캡션이 본문 문장 사이에 끼어 있으면 앞뒤 문장을 자연스럽게 연결합니다.

    매개변수:
        text: 처리할 텍스트 문자열.

    반환값:
        캡션이 제거된 텍스트 문자열.
    """
    import re

    _SENT_END = frozenset(".?!;:。！？")
    # 캡션 패턴: "그림 1.", "Figure 1:", "Fig. 2-3.", "Table 4." 등
    _CAPTION_RE = re.compile(
        r"^(?:그림|표|Figure|Fig\.?|Table)\s*\d+(?:[-\.]\d+)*(?![.-]\d)\s*[\.:]",
        re.IGNORECASE,
    )
    lines = text.splitlines()
    result: list[str] = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not _CAPTION_RE.match(stripped):
            result.append(lines[i])  # 캡션이 아닌 줄은 그대로 유지
            i += 1
            continue

        # 캡션 라인 발견 — 이전/다음 비어있지 않은 줄 탐색
        prev_idx = len(result) - 1
        while prev_idx >= 0 and not result[prev_idx].strip():
            prev_idx -= 1

        next_i = i + 1
        while next_i < len(lines) and not lines[next_i].strip():
            next_i += 1

        ptext = result[prev_idx].strip() if prev_idx >= 0 else ""
        ntext = lines[next_i].strip() if next_i < len(lines) else ""

        if (ptext and ptext[-1] not in _SENT_END
                and ntext and ntext[0].islower()):
            # 이전 줄 끝 + 다음 줄 머리 연결 (캡션 완전 제거)
            result[prev_idx] = result[prev_idx].rstrip() + " " + ntext
            while len(result) > prev_idx + 1:
                result.pop()
            i = next_i + 1
        else:
            # 앞뒤 연결이 불가능한 경우 — 캡션만 제거하고 다음 줄로 이동
            i = next_i

    return "\n".join(result)


@st.cache_data(ttl=600)
def get_pdf_page_count(file_path: str) -> int:
    """파일의 총 페이지 수(또는 슬라이드 수)를 반환합니다.

    매개변수:
        file_path: PDF, DOCX, PPTX 파일의 경로 문자열.

    반환값:
        정수형 페이지 수. 파일을 열 수 없으면 0 반환.
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".docx":
        return 1  # Word 문서는 전체를 단일 페이지로 처리
    if ext == ".pptx":
        try:
            from pptx import Presentation
            return len(Presentation(file_path).slides)
        except Exception:
            return 0
    # PDF는 pypdfium2로 페이지 수 확인
    try:
        import pypdfium2 as pdfium
        return len(pdfium.PdfDocument(file_path))
    except Exception:
        return 0


@st.cache_data(ttl=600)
def get_pdf_page_image_bytes(file_path: str, page_num: int, scale: float = 1.5) -> bytes | None:
    """PDF 특정 페이지를 PNG 이미지 바이트로 렌더링합니다.

    매개변수:
        file_path: PDF 파일 경로.
        page_num: 렌더링할 페이지 번호 (0-indexed).
        scale: 렌더링 배율. 기본값 1.5 (150 DPI).

    반환값:
        PNG 이미지의 bytes 객체. 실패 시 None.
    """
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(file_path)
        if page_num >= len(pdf):
            return None
        bitmap = pdf[page_num].render(scale=scale)
        buf = BytesIO()
        bitmap.to_pil().save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return None


@st.cache_resource
def _get_easyocr_reader(strategy: str):
    """EasyOCR Reader를 앱 시작 시 한 번만 로딩합니다. (모델 재로딩 방지)

    @st.cache_resource는 앱 전체 수명 동안 객체를 한 번만 생성합니다.
    EasyOCR 모델은 수백 MB 크기이므로 재생성을 피하는 것이 중요합니다.

    매개변수:
        strategy: OCR 전략 문자열. "OCR_KOR_ENG"이면 한/영, 그 외엔 한/영/일 지원.

    반환값:
        easyocr.Reader 인스턴스.
    """
    import easyocr
    langs = ["ko", "en"] if strategy == "OCR_KOR_ENG" else ["ko", "en", "ja"]
    return easyocr.Reader(langs, gpu=True, verbose=False)


@st.cache_data(ttl=600)
def extract_text_page(file_path: str, page_num: int, strategy: str, doc_id: int = 0) -> str:
    """페이지 단위 텍스트를 추출합니다.

    doc_id가 0이 아니면 page_contents DB를 먼저 조회하고 body_text가 있으면 반환합니다.
    DB에 없으면 strategy에 따라 PDF/DOCX/PPTX에서 직접 추출합니다.

    지원 전략:
        DOCX_EXTRACT       : Word 문서 텍스트 추출
        PPTX_EXTRACT       : PowerPoint 슬라이드 텍스트 추출
        DIGITAL_EXTRACT    : pdfplumber를 이용한 단단 PDF 추출
        DIGITAL_EXTRACT_MULTI : 컬럼 인식 다단 PDF 추출 (OcrEngine 재사용)
        OCR_KOR_ENG, OCR_JPN: EasyOCR 기반 이미지 OCR
        나머지              : Tesseract OCR

    매개변수:
        file_path: 문서 파일 경로.
        page_num: 페이지 번호 (0-indexed).
        strategy: OCR/추출 전략 문자열.
        doc_id: 데이터베이스 문서 ID. 0이면 DB 조회 생략.

    반환값:
        추출된 텍스트 문자열. 실패 시 오류 메시지 문자열.
    """
    try:
        # DB-first: 파이프라인이 저장한 구조화 본문 텍스트 우선 사용
        if doc_id:
            row = page_db.get_page_content(doc_id, page_num)
            if row and row.get("body_text"):
                return row["body_text"]

        if strategy == "DOCX_EXTRACT":
            from src.ocr_engine import OcrEngine
            pages = OcrEngine()._extract_docx_pages(file_path)
            return pages[0] if pages else ""

        if strategy == "PPTX_EXTRACT":
            from pptx import Presentation
            prs = Presentation(file_path)
            if page_num >= len(prs.slides):
                return "(슬라이드 없음)"
            slide = prs.slides[page_num]
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
                elif shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text_frame.text.strip() for c in row.cells
                                 if c.text_frame.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            return "\n\n".join(texts)

        if strategy == "DIGITAL_EXTRACT":
            import pdfplumber
            from src.ocr_engine import OcrEngine
            with pdfplumber.open(file_path) as pdf:
                if page_num < len(pdf.pages):
                    # 워터마크 필터링 후 텍스트 추출, 캡션 제거
                    return _strip_inline_captions(
                        _plumber_page_to_text(
                            OcrEngine.filter_watermarks(pdf.pages[page_num])
                        )
                    )
            return "(페이지 없음)"

        if strategy == "DIGITAL_EXTRACT_MULTI":
            import pdfplumber
            from src.ocr_engine import OcrEngine
            with pdfplumber.open(file_path) as pdf:
                if page_num >= len(pdf.pages):
                    return "(페이지 없음)"
                page  = OcrEngine.filter_watermarks(pdf.pages[page_num])
                words = page.extract_words()
                if not words:
                    return "(추출된 텍스트 없음)"
                # 컬럼 갭을 먼저 탐지 → 학술 헤더 감지에 전달
                # (gap_x 가 있으면 좌측 컬럼 전용 제목/저자도 올바르게 처리)
                hdr_w, main_w, ftr_w = _split_header_footer(page, words)
                pw    = float(page.width)
                gap_x = OcrEngine._find_column_gap(main_w, pw) if main_w else None
                academic = _extract_academic_header(page, gap_x=gap_x)
                if academic:
                    return academic
                # 학술 헤더 미감지 → 일반 컬럼 분리 추출
                if not main_w:
                    return _strip_inline_captions(_plumber_page_to_text(page))
                if gap_x is None:
                    return _strip_inline_captions(_plumber_page_to_text(page))
                body_start_y = OcrEngine._find_body_start_y(main_w, gap_x)
                body_text = OcrEngine._compose_multicol_text(main_w, gap_x, body_start_y)
                parts: list[str] = []
                if hdr_w:
                    hdr = _words_one_line(hdr_w)
                    if hdr.strip():
                        parts.append(f"[헤더]\n{hdr}")
                parts.append(body_text)
                if ftr_w:
                    ftr = _words_one_line(ftr_w)
                    if ftr.strip():
                        parts.append(f"[푸터]\n{ftr}")
                return _strip_inline_captions("\n\n".join(parts))

        # OCR 전략: 해당 페이지 이미지를 pypdfium2로 렌더링 (scale=2.0 → 200 DPI)
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(file_path)
        if page_num >= len(pdf):
            return "(페이지 없음)"
        bitmap = pdf[page_num].render(scale=2.0)
        pil_img = bitmap.to_pil()

        # EasyOCR 전략 (OCR_KOR_ENG, OCR_JPN)
        # page_contents DB가 비어있는 기처리 문서도 올바른 엔진으로 표시
        if strategy in ("OCR_KOR_ENG", "OCR_JPN"):
            import numpy as np
            reader = _get_easyocr_reader(strategy)
            result = reader.readtext(np.array(pil_img), detail=0, paragraph=True)
            return fix_easyocr_korean("\n".join(result))

        # Tesseract 전략 (OCR_ENG, OCR_MULTI_COL, OCR_FORMULA 등)
        import pytesseract
        if app_config.ocr_tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = app_config.ocr_tesseract_cmd
        tcfg = STRATEGY_TESSERACT_CONFIG.get(strategy, {})
        lang = tcfg.get("lang") or "kor+eng"
        psm  = tcfg.get("psm", 3)   # PSM(Page Segmentation Mode): 페이지 레이아웃 분석 방식
        oem  = tcfg.get("oem", 1)   # OEM(OCR Engine Mode): 1=LSTM 기반 신경망 엔진
        custom_cfg = f"--psm {psm} --oem {oem}"
        return pytesseract.image_to_string(pil_img, lang=lang, config=custom_cfg)

    except Exception as exc:
        return f"텍스트 추출 실패: {exc}"


# ---------------------------------------------------------------------------
# OCR 텍스트 리플로우 — 단락 내 줄바꿈으로 끊어진 문장 연결
# ---------------------------------------------------------------------------
def _legacy_reflow_ocr_text(text: str) -> str:
    """단락 내 줄바꿈을 공백으로 이어 연속 문장을 한 줄로 합칩니다. (구 버전)

    OCR은 물리적 줄바꿈을 그대로 출력하므로, 의미상 같은 문장이 여러 줄로
    나뉘어져 있을 때 이를 하나로 이어주는 역할을 합니다.

    규칙:
    - 빈 줄(단락 경계)은 유지
    - 하이픈 줄바꿈(-\\n)은 공백 없이 연결  (예: multi-\\ncolumn → multicolumn)
    - 나머지 줄바꿈은 공백으로 연결

    매개변수:
        text: OCR로 얻은 원시 텍스트.

    반환값:
        줄바꿈이 정리된 텍스트 문자열.
    """
    import re
    paragraphs = re.split(r'(?:\r?\n){2,}', text.strip())  # 빈 줄 기준으로 단락 분리
    merged: list[str] = []
    for para in paragraphs:
        buf = ""
        for line in para.splitlines():
            line = line.rstrip()
            if not line:
                continue
            if not buf:
                buf = line
            elif buf.endswith("-"):
                buf = buf[:-1] + line          # 하이픈 이음: 공백 없이 연결
            else:
                buf = buf + " " + line         # 일반 이음: 공백으로 연결
        if buf:
            merged.append(buf)
    return "\n\n".join(merged)


def reflow_ocr_text(text: str) -> str:
    """OCR 텍스트의 불필요한 줄바꿈을 정리합니다. (현재 버전 — src.preprocessor 위임)

    매개변수:
        text: 정리할 텍스트 문자열.

    반환값:
        줄바꿈이 정리된 텍스트 문자열.
    """
    return normalize_hard_linebreaks(text)


# ---------------------------------------------------------------------------
# PDF 뷰어 — 스플릿 뷰 렌더링
# ---------------------------------------------------------------------------

# 반투명 네비게이션 바 CSS (최초 1회 주입)
# 뷰어 상단의 페이지 이동 버튼들을 스타일링합니다.
_NAV_CSS = """
<style>
/* 뷰어 네비게이션 바 — 반투명 글래스 효과 */
.pdf-nav-bar {
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 6px;
    background: rgba(14, 17, 23, 0.55);
    backdrop-filter: blur(10px);
    -webkit-backdrop-filter: blur(10px);
    border-radius: 10px;
    padding: 4px 12px;
    margin: 0 0 6px 0;
    border: 1px solid rgba(255,255,255,0.12);
}
.pdf-nav-bar .nav-info {
    color: rgba(255,255,255,0.90);
    font-size: 0.85em;
    padding: 0 12px;
    white-space: nowrap;
}
/* number_input 위 여백 제거 */
.pdf-nav-bar + div [data-testid="stNumberInput"] {
    margin-top: 0;
}

"""


def render_pdf_viewer(doc: dict, key_prefix: str = ""):
    """원본 PDF 페이지(좌) ↔ OCR 추출 텍스트(우)를 나란히 표시합니다.

    네비게이션 바는 문서 영역 위에 반투명 오버레이 스타일로 표시하여
    문서 표시 영역을 최대화합니다.

    매개변수:
        doc: documents 테이블의 한 행을 dict로 변환한 객체.
             "file_path", "id", "ocr_strategy" 키를 사용합니다.
        key_prefix: 여러 문서를 동시에 표시할 때 세션 상태 키 충돌을 방지하기 위한 접두사.
    """
    file_path = doc["file_path"]
    doc_id    = doc["id"]
    strategy  = doc.get("ocr_strategy") or "DIGITAL_EXTRACT"

    if not Path(file_path).exists():
        st.warning(f"파일을 찾을 수 없습니다: `{file_path}`")
        return

    page_count = get_pdf_page_count(file_path)
    if page_count == 0:
        st.warning("PDF를 열 수 없습니다.")
        return

    # ── 페이지 번호 세션 상태 초기화 ────────────────────────────────────
    # pg_key: 이 문서 뷰어의 현재 페이지 번호를 저장하는 세션 상태 키
    pg_key = f"{key_prefix}pg_{doc_id}"
    if pg_key not in st.session_state:
        st.session_state[pg_key] = 1  # 1-indexed (첫 페이지 = 1)

    # 버튼 클릭을 위젯 렌더 전에 먼저 처리
    # Streamlit은 버튼 클릭 시 스크립트를 재실행하므로, 이 시점에 상태를 업데이트해야 합니다.
    if st.session_state.get(f"{pg_key}_first"):
        st.session_state[pg_key] = 1
    if st.session_state.get(f"{pg_key}_prev"):
        st.session_state[pg_key] = max(1, st.session_state[pg_key] - 1)
    if st.session_state.get(f"{pg_key}_next"):
        st.session_state[pg_key] = min(page_count, st.session_state[pg_key] + 1)
    if st.session_state.get(f"{pg_key}_last"):
        st.session_state[pg_key] = page_count

    current = st.session_state[pg_key]
    page_num = current - 1  # 내부 처리는 0-indexed

    # ── 반투명 네비게이션 바 CSS 주입 ────────────────────────────────────
    st.markdown(_NAV_CSS, unsafe_allow_html=True)

    # ── 네비게이션 바 (문서 위, 반투명 글래스) ────────────────────────────
    # 레이아웃: [⏮] [◀] [─ 페이지 직접 입력 / 전체 페이지 │ OCR전략 ─] [▶] [⏭]
    nc1, nc2, nc3, nc4, nc5, nc6, nc7 = st.columns([1, 1, 2, 1, 2, 1, 1])

    with nc1:
        st.button("⏮", key=f"{pg_key}_first", disabled=(current == 1),
                  use_container_width=True, help="처음 페이지")
    with nc2:
        st.button("◀", key=f"{pg_key}_prev", disabled=(current == 1),
                  use_container_width=True, help="이전 페이지")
    with nc3:
        # number_input: 사용자가 직접 페이지 번호를 입력할 수 있는 위젯
        st.number_input(
            "페이지",
            min_value=1, max_value=page_count, step=1,
            key=pg_key,
            label_visibility="collapsed",
        )
    with nc4:
        st.markdown(
            f"<div style='text-align:center;padding:6px 0;font-size:0.9em;'>"
            f"/ {page_count}</div>",
            unsafe_allow_html=True,
        )
    with nc5:
        st.markdown(
            f"<div style='text-align:center;padding:6px 0;font-size:0.82em;"
            f"color:#888;'>전략: <code>{strategy}</code></div>",
            unsafe_allow_html=True,
        )
    with nc6:
        st.button("▶", key=f"{pg_key}_next", disabled=(current == page_count),
                  use_container_width=True, help="다음 페이지")
    with nc7:
        st.button("⏭", key=f"{pg_key}_last", disabled=(current == page_count),
                  use_container_width=True, help="마지막 페이지")

    # ── 문서 본문 — 전체 너비를 PDF(좌)·텍스트(우)에 할당 ─────────────────
    col_pdf, col_text = st.columns(2)

    with col_pdf:
        ext = Path(file_path).suffix.lower()
        if ext == ".docx":
            st.caption("원본 문서 (Word)")
            st.info("Word 문서는 원본 미리보기를 지원하지 않습니다.\n우측에서 추출된 텍스트를 확인하세요.")
        elif ext == ".pptx":
            st.caption(f"원본 슬라이드  {current} / {page_count}")
            st.info(f"PowerPoint 슬라이드 {current} — 원본 미리보기를 지원하지 않습니다.\n우측에서 추출된 텍스트를 확인하세요.")
        else:
            st.caption(f"원본 PDF  p.{current} / {page_count}")
            img_bytes = get_pdf_page_image_bytes(file_path, page_num)
            if img_bytes:
                st.image(img_bytes, use_container_width=True)
            else:
                st.error("페이지 렌더링 실패")

    with col_text:
        st.caption(f"추출 텍스트  p.{current}")
        with st.spinner("텍스트 추출 중..."):
            # 텍스트 추출 후 리플로우(줄바꿈 정리) 적용
            text = reflow_ocr_text(extract_text_page(file_path, page_num, strategy, doc_id))
        st.text_area(
            label="",
            value=text,
            height=680,
            key=f"{key_prefix}txtarea_{doc_id}_{page_num}",
        )


# ---------------------------------------------------------------------------
# 파이프라인 실행 (백그라운드 스레드)
# ---------------------------------------------------------------------------
def _run_pipeline_bg():
    """파이프라인을 실행하고 결과를 세션 상태에 기록합니다. (백그라운드 스레드 대상 함수)

    Streamlit은 싱글 스레드 모델이므로, 파이프라인처럼 오래 걸리는 작업은
    별도 스레드에서 실행하고 결과를 session_state에 저장합니다.
    """
    st.session_state.pipeline_running = True
    st.session_state.pipeline_log = "파이프라인 실행 중...\n"
    try:
        stats = pipeline.run_once()  # 파이프라인 1회 실행 → 통계 반환
        st.session_state.pipeline_log = (
            f"완료!\n"
            f"  신규 탐지: {stats.scanned}건\n"
            f"  처리 완료: {stats.processed}건\n"
            f"  INDEXED:  {stats.indexed}건\n"
            f"  REVIEW:   {stats.review}건\n"
            f"  FAILED:   {stats.failed}건\n"
        )
    except Exception as exc:
        st.session_state.pipeline_log = f"오류 발생: {exc}"
    finally:
        st.session_state.pipeline_running = False  # 실행 완료(성공/실패 모두) 후 플래그 해제


def run_pipeline():
    """파이프라인을 데몬 스레드로 실행합니다.

    daemon=True: 메인 스레드(Streamlit 서버)가 종료되면 이 스레드도 함께 종료됩니다.
    """
    threading.Thread(target=_run_pipeline_bg, daemon=True).start()


# ---------------------------------------------------------------------------
# 단일 문서 재처리 헬퍼
# ---------------------------------------------------------------------------
def do_retry_and_process(doc_id: int):
    """실패하거나 검토가 필요한 문서를 재처리합니다.

    DB 상태를 FAILED/REVIEW_REQUIRED → NEW로 초기화한 뒤 파이프라인을 즉시 실행합니다.

    매개변수:
        doc_id: 재처리할 문서의 DB ID.
    """
    db.reset_for_retry(doc_id)
    refresh_data()
    with st.spinner(f"문서 {doc_id} 재처리 중..."):
        ok = pipeline.process_document(doc_id)
    refresh_data()
    if ok:
        st.success(f"문서 {doc_id} 재처리 완료 → INDEXED")
    else:
        st.error(f"문서 {doc_id} 재처리 실패. 로그를 확인하세요.")


def do_force_reprocess(doc_id: int):
    """INDEXED 상태의 문서를 강제로 처음부터 재처리합니다.

    이미 성공적으로 인덱싱된 문서도 강제로 재분석·재임베딩합니다.

    매개변수:
        doc_id: 강제 재처리할 문서의 DB ID.
    """
    db.force_reset_for_reprocess(doc_id)
    refresh_data()
    with st.spinner(f"문서 {doc_id} 강제 재처리 중..."):
        ok = pipeline.process_document(doc_id)
    refresh_data()
    if ok:
        st.success(f"문서 {doc_id} 강제 재처리 완료 → INDEXED")
    else:
        st.error(f"문서 {doc_id} 강제 재처리 실패. 로그를 확인하세요.")


def do_reset(doc_id: int):
    """문서 상태를 NEW로 초기화합니다. (파이프라인 재실행 없음)

    매개변수:
        doc_id: 초기화할 문서의 DB ID.
    """
    db.reset_for_retry(doc_id)
    refresh_data()
    st.success(f"문서 {doc_id} → NEW 초기화 완료.")


# ---------------------------------------------------------------------------
# 문서 행 렌더링 (재사용)
# ---------------------------------------------------------------------------
def render_document_row(doc: dict, key_prefix: str = ""):
    """문서 목록에서 한 행(문서)을 렌더링합니다.

    접혀 있을 때: 파일명, 상태, 소스, 레이아웃, 언어, 수식 여부를 한 줄로 표시합니다.
    펼쳐지면: 상세 정보(경로, OCR 품질, 오류 메시지, 처리 로그)와 액션 버튼을 표시합니다.

    매개변수:
        doc: documents 테이블의 한 행 dict.
        key_prefix: 위젯 키 충돌 방지용 접두사 (탭마다 다르게 설정).
    """
    status_icon = STATUS_EMOJI.get(doc["status"], "•")
    layout_val  = doc.get("layout_type") or ""
    expand_key  = f"{key_prefix}expand_{doc['id']}"  # 토글 위젯의 고유 키

    with st.container(border=False):
        # ── [토글] 파일명+상태 | 소스 | 레이아웃 | 언어 | 수식 — 항상 표시 ─
        c_btn, c_name, c_src, c_layout, c_lang, c_formula = st.columns(
            [0.4, 8.6, 1.1, 0.9, 1.4, 0.9]
        )
        with c_btn:
            # 토글 스위치: 클릭하면 상세 정보가 펼쳐짐
            is_expanded = st.toggle(
                label="",
                key=expand_key,
                label_visibility="collapsed",
            )
        c_name.markdown(f"{status_icon}  **{doc['file_name']}**  —  {doc['status']}")
        if doc.get("source_type"):
            c_src.caption(doc["source_type"])
            c_layout.caption("다단" if "MULTI" in layout_val else "단일")
            c_lang.caption(doc.get("detected_languages") or "—")
            c_formula.caption("수식O" if doc.get("has_formula") else "수식X")
        else:
            # 분석이 아직 완료되지 않아 분류 정보가 없는 경우
            c_src.caption("—")
            c_layout.caption("—")
            c_lang.caption("—")
            c_formula.caption("—")

        # ── 상세 정보 / 재처리 (토글 펼침) ───────────────────────────────
        if is_expanded:
            st.divider()
            col_info, col_action = st.columns([3, 1])

            with col_info:
                st.markdown(
                    f"**ID**: {doc['id']}  |  **상태**: "
                    + status_badge(doc["status"]),
                    unsafe_allow_html=True,
                )
                st.markdown(f"**경로**: `{doc['file_path']}`")

                if doc.get("source_type"):
                    st.markdown(
                        f"소스: **{doc['source_type']}**  |  "
                        f"레이아웃: **{doc.get('layout_type', '-')}**  |  "
                        f"언어: **{doc.get('detected_languages', '-')}**  |  "
                        f"수식: **{'있음' if doc.get('has_formula') else '없음'}**  |  "
                        f"OCR전략: **{doc.get('ocr_strategy', '-')}**"
                    )

                if doc.get("ocr_quality_score") is not None:
                    st.markdown(f"OCR 품질 점수: **{doc['ocr_quality_score']:.3f}**")

                if doc.get("error_message"):
                    st.error(f"오류: {doc['error_message']}")

                st.caption(
                    f"등록: {doc.get('created_at', '')}  |  "
                    f"갱신: {doc.get('updated_at', '')}  |  "
                    f"재시도 횟수: {doc.get('retry_count', 0)}회"
                )

                with st.expander("처리 로그"):
                    logs = load_all_logs().get(doc["id"], [])
                    if logs:
                        for log in logs:
                            icon = "✅" if log["result"] == "SUCCESS" else "❌"
                            st.markdown(
                                f"`{log['logged_at']}` {icon} "
                                f"**{log['step']}** [{log['result']}] "
                                f"{log.get('message', '')}"
                            )
                    else:
                        st.caption("로그 없음")

            with col_action:
                view_key = f"{key_prefix}view_{doc['id']}"
                st.checkbox("📖 원본/OCR 비교", key=view_key)

                # 재처리 가능한 상태(FAILED, REVIEW_REQUIRED 등)일 때만 버튼 표시
                if doc["status"] in Status.RETRYABLE:
                    k = f"{key_prefix}proc_{doc['id']}"
                    if st.button("🔄 즉시 재처리", key=k, use_container_width=True):
                        do_retry_and_process(doc["id"])
                        st.rerun()  # 상태 변경 후 화면 즉시 갱신

                    k2 = f"{key_prefix}reset_{doc['id']}"
                    if st.button("➡️ NEW 초기화", key=k2, use_container_width=True):
                        do_reset(doc["id"])
                        st.rerun()

                # 이미 인덱싱된 문서는 강제 재처리 버튼만 표시
                if doc["status"] == Status.INDEXED:
                    k3 = f"{key_prefix}force_{doc['id']}"
                    if st.button("🔁 강제 재처리", key=k3, use_container_width=True):
                        do_force_reprocess(doc["id"])
                        st.rerun()

            # 원본/OCR 비교 체크박스가 체크된 경우 PDF 뷰어 표시
            if st.session_state.get(f"{key_prefix}view_{doc['id']}", False):
                st.divider()
                render_pdf_viewer(doc, key_prefix=key_prefix)


# ===========================================================================
# 사이드바
# with st.sidebar 블록 안의 모든 위젯은 화면 왼쪽 사이드바에 표시됩니다.
# ===========================================================================
with st.sidebar:
    st.title("📄 RAG Pipeline")
    # 버전 및 수정 시각을 작은 글씨로 표시
    st.markdown(
        f'<div style="margin-top:-8px;margin-bottom:4px;font-size:0.75em;color:#999;">'
        f'<b style="color:#7eb8f7;">v{APP_VERSION}</b> &nbsp; {APP_MODIFIED}'
        f'</div>',
        unsafe_allow_html=True,
    )
    st.divider()

    st.subheader("파이프라인 실행")
    btn_disabled = st.session_state.pipeline_running  # 실행 중일 때는 버튼 비활성화
    if st.button(
        "▶ 파이프라인 실행",
        type="primary",
        disabled=btn_disabled,
        use_container_width=True,
    ):
        run_pipeline()
        st.rerun()

    if st.session_state.pipeline_running:
        st.info("파이프라인 실행 중...")

    if st.session_state.pipeline_log:
        st.code(st.session_state.pipeline_log, language=None)

    st.divider()
    if st.button("🔄 새로고침", use_container_width=True):
        # (250623) 쿨다운 2초 — 마지막 새로고침 후 2초 내 재클릭은 무시
        if time.time() - st.session_state.last_refresh > 2.0:
            refresh_data()
            st.rerun()

    st.caption(
        f"마지막 갱신: "
        f"{time.strftime('%H:%M:%S', time.localtime(st.session_state.last_refresh))}"
    )

    st.divider()
    st.subheader("전체 재OCR")
    st.caption("모든 문서를 NEW로 초기화하고 OCR부터 재처리합니다.")
    if st.button("🔁 전체 재OCR", use_container_width=True, type="secondary"):
        count = db.reset_all_for_reocr()
        refresh_data()
        st.success(f"{count}건 초기화 완료.")
        st.info("▶ 파이프라인 실행 버튼을 눌러 재처리를 시작하세요.")
        st.rerun()


# ===========================================================================
# 메인 타이틀 + 요약 (한 줄)
# ===========================================================================
# 전체 문서 목록 로드 및 상태별 집계
docs = load_documents()
status_counts = Counter(d["status"] for d in docs)  # 상태별 문서 수 집계
_idx  = status_counts.get(Status.INDEXED, 0)
_rev  = status_counts.get(Status.REVIEW_REQUIRED, 0)
_fail = status_counts.get(Status.FAILED, 0)
_new  = status_counts.get(Status.NEW, 0)

st.title("RAG Pipeline Manager")
# 전체 현황을 제목 오른쪽에 한 줄로 요약 표시
st.markdown(
    f'<div style="display:flex;justify-content:space-between;align-items:center;">'
    f'<span style="font-size:0.85em;color:#888;">'
    f'전체 {len(docs)}건&nbsp;&nbsp;|&nbsp;&nbsp;🎯 INDEXED {_idx}'
    f'&nbsp;&nbsp;|&nbsp;&nbsp;⚠️ 검토 {_rev}'
    f'&nbsp;&nbsp;|&nbsp;&nbsp;❌ 실패 {_fail}'
    f'&nbsp;&nbsp;|&nbsp;&nbsp;🆕 대기 {_new}'
    f'</span>'
    f'<span>'
    f'<span style="background:#1e3a5f;color:#7eb8f7;padding:2px 8px;'
    f'border-radius:5px;font-size:0.8em;font-family:monospace;">v{APP_VERSION}</span>'
    f'&nbsp;<span style="font-size:0.72em;color:#999;">{APP_MODIFIED}</span>'
    f'</span>'
    f'</div>',
    unsafe_allow_html=True,
)

# ── 탭 ──────────────────────────────────────────────────────────────────────
# st.tabs()로 여러 탭을 생성합니다. 각 탭은 with 블록으로 내용을 채웁니다.
tab_all, tab_review, tab_failed, tab_dashboard, tab_assets, tab_diag = st.tabs(
    ["📋 전체 파일 목록", "⚠️ 검토 필요", "❌ 실패 목록", "📊 대시보드", "🖼️ 표·그림 자산", "🔬 PDF 분석 진단"]
)


# ===========================================================================
# 탭 1 — 전체 파일 목록
# ===========================================================================
with tab_all:
    # 검색창과 상태 필터를 나란히 배치
    fc1, fc2 = st.columns([2, 1])
    with fc1:
        search = st.text_input("파일명 검색", placeholder="🔍 파일명 검색...", key="search_all", label_visibility="collapsed")
    with fc2:
        all_statuses = sorted(set(d["status"] for d in docs))
        status_filter = st.multiselect("상태 필터", all_statuses, key="filter_all", label_visibility="collapsed")

    # 검색어와 상태 필터 적용
    filtered = docs
    if search:
        filtered = [d for d in filtered if search.lower() in d["file_name"].lower()]
    if status_filter:
        filtered = [d for d in filtered if d["status"] in status_filter]

    if not filtered:
        st.info("표시할 문서가 없습니다.")
    else:
        # 컬럼 헤더 출력 (실제 문서 행과 같은 비율의 컬럼 사용)
        _, h_name, h_src, h_layout, h_lang, h_formula = st.columns(
            [0.4, 8.6, 1.1, 0.9, 1.4, 0.9]
        )
        h_name.caption("파일명 / 진행상태")
        h_src.caption("소스")
        h_layout.caption("레이아웃")
        h_lang.caption("언어")
        h_formula.caption("수식")
        for doc in filtered:
            render_document_row(doc, key_prefix="all_")


# ===========================================================================
# 탭 2 — 검토 필요
# ===========================================================================
with tab_review:
    review_docs = [d for d in docs if d["status"] == Status.REVIEW_REQUIRED]
    st.subheader(f"검토 필요 문서 ({len(review_docs)}건)")

    if not review_docs:
        st.success("검토가 필요한 문서가 없습니다.")
    else:
        st.info(
            "아래 문서들은 OCR 품질 또는 텍스트 길이가 기준에 미치지 못했습니다. "
            "확인 후 재처리하세요."
        )
        for doc in review_docs:
            with st.container(border=False):
                rc1, rc2 = st.columns([4, 1])
                with rc1:
                    st.markdown(f"**{doc['file_name']}** (ID: {doc['id']})")
                    st.markdown(f"**사유**: {doc.get('error_message', '알 수 없음')}")
                    if doc.get("source_type"):
                        st.caption(
                            f"소스: {doc['source_type']}  |  "
                            f"OCR전략: {doc.get('ocr_strategy', '-')}  |  "
                            f"품질점수: {doc.get('ocr_quality_score', '-')}"
                        )
                    with st.expander("처리 로그"):
                        for log in load_all_logs().get(doc["id"], []):
                            icon = "✅" if log["result"] == "SUCCESS" else "❌"
                            st.markdown(
                                f"`{log['logged_at']}` {icon} "
                                f"**{log['step']}** [{log['result']}] "
                                f"{log.get('message', '')}"
                            )
                with rc2:
                    if st.button("🔄 즉시 재처리", key=f"rev_proc_{doc['id']}",
                                 use_container_width=True):
                        do_retry_and_process(doc["id"])
                        st.rerun()
                    if st.button("➡️ NEW 초기화", key=f"rev_reset_{doc['id']}",
                                 use_container_width=True):
                        do_reset(doc["id"])
                        st.rerun()


# ===========================================================================
# 탭 3 — 실패 목록
# ===========================================================================
with tab_failed:
    failed_docs = [d for d in docs if d["status"] == Status.FAILED]
    st.subheader(f"실패 문서 ({len(failed_docs)}건)")

    if not failed_docs:
        st.success("실패한 문서가 없습니다.")
    else:
        for doc in failed_docs:
            with st.container(border=False):
                fc1, fc2 = st.columns([4, 1])
                with fc1:
                    st.markdown(f"**{doc['file_name']}** (ID: {doc['id']})")
                    st.markdown(
                        f"**실패 단계**: `{doc.get('failed_step', '?')}`  |  "
                        f"**재시도 횟수**: {doc.get('retry_count', 0)}회"
                    )
                    if doc.get("error_message"):
                        st.error(doc["error_message"])
                    with st.expander("처리 로그"):
                        for log in load_all_logs().get(doc["id"], []):
                            icon = "✅" if log["result"] == "SUCCESS" else "❌"
                            st.markdown(
                                f"`{log['logged_at']}` {icon} "
                                f"**{log['step']}** [{log['result']}] "
                                f"{log.get('message', '')}"
                            )
                with fc2:
                    if st.button("🔄 즉시 재처리", key=f"fail_proc_{doc['id']}",
                                 use_container_width=True):
                        do_retry_and_process(doc["id"])
                        st.rerun()
                    if st.button("➡️ NEW 초기화", key=f"fail_reset_{doc['id']}",
                                 use_container_width=True):
                        do_reset(doc["id"])
                        st.rerun()


# ===========================================================================
# 탭 4 — 대시보드
# ===========================================================================
with tab_dashboard:
    st.subheader("처리 현황 대시보드")

    if not docs:
        st.info("등록된 문서가 없습니다.")
    else:
        # 상태별 분포 테이블
        st.markdown("#### 상태별 분포")
        status_table = [
            {
                "상태": f"{STATUS_EMOJI.get(s, '')} {s}",
                "건수": c,
                "비율": f"{c / len(docs) * 100:.1f}%",
            }
            for s, c in sorted(status_counts.items(), key=lambda x: -x[1])  # 건수 내림차순 정렬
        ]
        st.table(status_table)

        st.divider()

        # 분류 통계 (소스 유형, 레이아웃, OCR 전략별)
        source_counts  = Counter(d["source_type"]  for d in docs if d.get("source_type"))
        layout_counts  = Counter(d["layout_type"]  for d in docs if d.get("layout_type"))
        strategy_counts = Counter(d["ocr_strategy"] for d in docs if d.get("ocr_strategy"))
        formula_count  = sum(1 for d in docs if d.get("has_formula"))
        analyzed_count = sum(1 for d in docs if d.get("source_type"))

        dc1, dc2 = st.columns(2)

        with dc1:
            st.markdown("#### PDF 분류 통계")
            if source_counts:
                st.markdown("**소스 유형**")
                for k, v in source_counts.items():
                    st.markdown(f"- {k}: **{v}건**")
            if layout_counts:
                st.markdown("**레이아웃 유형**")
                for k, v in layout_counts.items():
                    st.markdown(f"- {k}: **{v}건**")
            st.markdown(
                f"**수식 포함**: {formula_count}건 "
                f"(분류 완료 {analyzed_count}건 중)"
            )

        with dc2:
            st.markdown("#### OCR 전략 분포")
            if strategy_counts:
                max_v = max(strategy_counts.values())
                for k, v in sorted(strategy_counts.items(), key=lambda x: -x[1]):
                    # 막대 그래프를 ASCII 블록 문자로 간단히 시각화
                    bar = "█" * int(v / max_v * 10)
                    st.markdown(f"`{k}` **{v}건**  {bar}")
            else:
                st.caption("분류 데이터 없음")

        # OCR 품질 통계
        quality_scores = [
            d["ocr_quality_score"]
            for d in docs
            if d.get("ocr_quality_score") is not None
        ]
        if quality_scores:
            st.divider()
            st.markdown("#### OCR 품질 점수 통계")
            avg_q = sum(quality_scores) / len(quality_scores)
            qc1, qc2, qc3 = st.columns(3)
            qc1.metric("평균 품질", f"{avg_q:.3f}")
            qc2.metric("최솟값", f"{min(quality_scores):.3f}")
            qc3.metric("최댓값", f"{max(quality_scores):.3f}")

            # 구간별 분포 — 품질 점수를 4개 구간으로 나눠 건수 표시
            bins   = [(0,    0.5,  "0.0 ~ 0.5 (낮음)"),
                      (0.5,  0.7,  "0.5 ~ 0.7 (보통)"),
                      (0.7,  0.85, "0.7 ~ 0.85 (양호)"),
                      (0.85, 1.01, "0.85 ~ 1.0 (우수)")]
            for lo, hi, label in bins:
                cnt = sum(1 for q in quality_scores if lo <= q < hi)
                st.markdown(f"- {label}: **{cnt}건**")


# ===========================================================================
# 탭 5 — PDF 분석 진단
# _diag_analyze 함수는 탭 6 렌더링 전에 정의되어야 하므로 여기서 선언합니다.
# ===========================================================================
@st.cache_data(ttl=60)
def _diag_analyze(pdf_path: str):
    """PDF 레이아웃 진단 — 줄 커버리지 프로파일을 반환합니다.

    각 페이지(최대 5페이지)에서 단어의 x 좌표 분포를 분석하여
    다단 레이아웃 여부와 컬럼 갭 위치를 감지합니다.

    동작 방식:
    1) 단어를 x 좌표 bin으로 매핑해 각 bin의 줄 커버리지를 계산합니다.
    2) 중앙(35%~65%) 구간에서 임계값 이하로 연속으로 비어 있는 bin이 3개 이상이면
       MULTI_COLUMN으로 판정합니다.

    매개변수:
        pdf_path: 분석할 PDF 파일의 경로 문자열.

    반환값:
        페이지별 결과 딕셔너리 목록.
        각 dict에는 page, n_lines, n_words, detected("MULTI_COLUMN"/"SINGLE_COLUMN"),
        line_cov(200개 bin의 커버리지 배열) 등이 포함됩니다.
        건너뛴 페이지는 "skipped" 키에 사유가 있습니다.
    """
    import pdfplumber

    BINS   = 200   # 가로를 200개 구간으로 나눠 단어 분포를 분석
    LINE_H = 3     # y 좌표를 3pt 단위로 반올림해 같은 줄로 묶음
    results = []

    try:
        with pdfplumber.open(pdf_path) as pdf:
            for pi, page in enumerate(pdf.pages[:5]):  # 최대 5페이지만 분석
                words = page.extract_words()
                page_width = float(page.width)
                if len(words) < 6:
                    results.append({"page": pi + 1, "skipped": "단어 < 6개"})
                    continue

                # 단어를 y 좌표 그룹(줄)별로, x 좌표를 bin 번호로 매핑
                y_ranges = defaultdict(list)
                for w in words:
                    yk = round(float(w["top"]) / LINE_H) * LINE_H  # y 좌표 그룹화
                    x0b = max(0, int(float(w["x0"]) / page_width * BINS))
                    x1b = min(BINS - 1, int(float(w["x1"]) / page_width * BINS))
                    y_ranges[yk].append((x0b, x1b))

                n_lines = len(y_ranges)
                if n_lines < 10:
                    results.append({"page": pi + 1, "skipped": f"줄 수 {n_lines} < 10 (그림/참고문헌 페이지)"})
                    continue

                # 각 bin이 전체 줄 중 몇 %에서 단어로 덮여 있는지 계산
                line_cov = np.zeros(BINS)
                for ranges in y_ranges.values():
                    covered = np.zeros(BINS, dtype=bool)
                    for x0b, x1b in ranges:
                        covered[x0b: x1b + 1] = True
                    line_cov += covered
                line_cov /= n_lines  # 정규화: 0.0(항상 비어있음) ~ 1.0(항상 단어 있음)

                # 중앙 구간(35%~65%)에서 갭 탐색
                s, e = int(BINS * 0.35), int(BINS * 0.65)
                center_cov = line_cov[s:e]
                side_mean  = float((line_cov[:s].mean() + line_cov[e:].mean()) / 2)
                # 갭 임계값: 양쪽 사이드의 평균 커버리지의 30% (최대 0.25)
                gap_thresh = min(0.25, side_mean * 0.30)

                # 중앙에서 가장 긴 연속 갭(임계값 미만 bin) 탐지
                max_gap = curr = 0
                gap_start = gap_end = -1
                for i, v in enumerate(center_cov < gap_thresh):
                    curr = curr + 1 if v else 0
                    if curr > max_gap:
                        max_gap = curr
                        gap_end   = s + i
                        gap_start = gap_end - max_gap + 1

                # 연속 갭이 3bin 이상이면 다단 레이아웃으로 판정
                detected = "MULTI_COLUMN" if max_gap >= 3 else "SINGLE_COLUMN"
                results.append({
                    "page":       pi + 1,
                    "skipped":    None,
                    "n_lines":    n_lines,
                    "n_words":    len(words),
                    "side_mean":  round(side_mean, 3),
                    "gap_thresh": round(gap_thresh, 3),
                    "max_gap":    max_gap,
                    "gap_pct":    f"{gap_start/BINS*100:.1f}%~{(gap_end+1)/BINS*100:.1f}%" if max_gap >= 3 else "-",
                    "detected":   detected,
                    "line_cov":   line_cov.tolist(),  # Streamlit 차트용으로 리스트로 변환
                })
    except Exception as exc:
        results.append({"page": "?", "skipped": str(exc)})

    return results


# ===========================================================================
# 탭 5 — 표·그림 자산 뷰어
# ===========================================================================
with tab_assets:
    st.subheader("표·그림 자산 뷰어")
    st.caption("파이프라인이 추출한 표(TABLE)·그림(FIGURE) 이미지와 OCR 결과를 확인합니다.")

    from src.asset_db import asset_db as _asset_db

    asset_counts = load_asset_counts()

    if not asset_counts:
        st.info("아직 추출된 자산이 없습니다. 파이프라인을 실행하면 표·그림이 자동으로 추출됩니다.")
    else:
        # 문서 선택 드롭다운: "파일명 (TABLE:N / FIGURE:M)" 형식으로 표시
        doc_options = {
            f"{v['name']} (TABLE:{v['counts'].get('TABLE',0)} / FIGURE:{v['counts'].get('FIGURE',0)})": k
            for k, v in asset_counts.items()
        }
        selected_label = st.selectbox("문서 선택", list(doc_options.keys()), key="asset_doc_sel")
        selected_doc_id = doc_options[selected_label]

        assets = _asset_db.get_assets_for_document(selected_doc_id)

        # 유형 필터 (전체 / TABLE / FIGURE)
        filter_type = st.radio("유형 필터", ["전체", "TABLE", "FIGURE"], horizontal=True, key="asset_type_filter")
        if filter_type != "전체":
            assets = [a for a in assets if dict(a)["asset_type"] == filter_type]

        st.markdown(f"**{len(assets)}건** 표시 중")

        for asset in assets:
            a = dict(asset)
            ref_tag  = a["ref_tag"]    # 참조 태그 (예: "Fig.1", "Table.2")
            atype    = a["asset_type"] # "TABLE" 또는 "FIGURE"
            page_num = a["page_num"] + 1  # 1-indexed로 변환해 표시
            ocr_text = a.get("ocr_text") or ""
            caption  = a.get("caption") or ""

            with st.expander(f"{ref_tag}  —  p.{page_num}  |  {atype}", expanded=False):
                col_img, col_txt = st.columns([1, 1])

                with col_img:
                    img_bytes = _asset_db.get_asset_image(a["id"])
                    if img_bytes:
                        st.image(img_bytes, caption=caption or ref_tag, use_container_width=True)
                    else:
                        st.warning("이미지 없음")

                with col_txt:
                    st.markdown(f"**참조 태그:** `{ref_tag}`")
                    st.markdown(f"**페이지:** {page_num}")
                    st.markdown(f"**유형:** {atype}")
                    if caption:
                        st.markdown(f"**캡션:** {caption}")
                    bbox = (a.get("bbox_x0"), a.get("bbox_y0"), a.get("bbox_x1"), a.get("bbox_y1"))
                    if all(v is not None for v in bbox):
                        st.markdown(f"**위치 (pt):** x0={bbox[0]:.1f}, y0={bbox[1]:.1f}, x1={bbox[2]:.1f}, y1={bbox[3]:.1f}")
                    else:
                        st.markdown("**위치 (pt):** —")
                    st.markdown("**OCR 텍스트:**")
                    if ocr_text:
                        st.text_area("", ocr_text, height=200, key=f"ocr_{a['id']}", label_visibility="collapsed")
                    else:
                        st.caption("OCR 결과 없음")


# ===========================================================================
# 탭 6 — PDF 레이아웃 분석 진단
# ===========================================================================
with tab_diag:
    st.subheader("PDF 레이아웃 분석 진단")
    st.caption("실제 PDF 파일 경로를 입력하면 페이지별 줄 커버리지 프로파일과 다단 감지 결과를 표시합니다.")

    diag_path = st.text_input(
        "PDF 파일 경로",
        placeholder="예: data/pdf_watch/paper.pdf",
        key="diag_path",
    )

    if diag_path and Path(diag_path).exists():
        with st.spinner("분석 중..."):
            diag_results = _diag_analyze(diag_path)

        for r in diag_results:
            pg = r["page"]
            if r.get("skipped"):
                st.info(f"p{pg}: 건너뜀 — {r['skipped']}")
                continue

            # 다단이면 초록 원, 단단이면 파란 원으로 구분
            col_badge = "🟢" if r["detected"] == "MULTI_COLUMN" else "🔵"
            with st.expander(
                f"{col_badge} p{pg}  [{r['detected']}]  "
                f"갭={r['max_gap']}bins @ {r['gap_pct']}  "
                f"(줄:{r['n_lines']}, side_mean:{r['side_mean']})",
                expanded=True,
            ):
                # 커버리지 프로파일 차트 — x축: 페이지 가로 위치(%), y축: 줄 커버리지
                lc = r["line_cov"]
                BINS = len(lc)
                xs   = [i / BINS * 100 for i in range(BINS)]  # 0~100% 범위로 변환
                thresh = r["gap_thresh"]

                # pandas DataFrame으로 변환해 Streamlit 꺾은선 차트로 표시
                import pandas as pd
                df_cov = pd.DataFrame({
                    "x(%)":      xs,
                    "줄커버리지": lc,
                    "임계값":    [thresh] * BINS,  # 갭 판정 기준선
                })
                st.line_chart(df_cov.set_index("x(%)"), height=180)

                mc1, mc2, mc3, mc4 = st.columns(4)
                mc1.metric("텍스트 줄 수", r["n_lines"])
                mc2.metric("단어 수", r["n_words"])
                mc3.metric("side_mean", r["side_mean"])
                mc4.metric("max_gap", f"{r['max_gap']} bin")

    elif diag_path:
        st.error(f"파일을 찾을 수 없습니다: `{diag_path}`")

