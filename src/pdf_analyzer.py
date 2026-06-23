# -*- coding: utf-8 -*-
"""
PDF Analyzer — PDF 문서를 분류 분석하여 PDFProfile 반환

분석 항목:
  1. source_type  : DIGITAL (텍스트 레이어 존재) / SCANNED (이미지 기반)
  2. layout_type  : SINGLE_COLUMN / MULTI_COLUMN (2단 이상 논문 형식)
  3. languages    : kor / eng / jpn 포함 여부
  4. has_formula  : 수식 포함 여부
"""

import logging
import re
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path

import cv2
import numpy as np
import pdfplumber
import pypdfium2 as pdfium
from langdetect import detect_langs, LangDetectException

from src.config import config
from src.database import SourceType, LayoutType

# 이 모듈 전용 로거를 생성합니다. 로그 출력 시 모듈 이름이 함께 표시됩니다.
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 언어 감지 임계값 (해당 문자 비율 기준)
# ---------------------------------------------------------------------------
# 전체 텍스트 중 해당 언어 문자가 차지하는 비율이 이 값 이상이어야 해당 언어로 판정합니다.
LANG_THRESHOLD = 0.03          # 전체 텍스트 대비 3% 이상이면 해당 언어 포함으로 판정
# 전체 텍스트 중 수학 기호가 차지하는 비율이 이 값 이상이어야 수식 포함으로 판정합니다.
FORMULA_CHAR_THRESHOLD = 0.008 # 수학 기호 밀도 0.8% 이상이면 수식 포함으로 판정

# ---------------------------------------------------------------------------
# Unicode 문자 범위 정규식 패턴
# 각 언어 및 수학 기호에 해당하는 유니코드 범위를 미리 컴파일해 둡니다.
# 미리 컴파일(re.compile)하면 반복 호출 시 속도가 빠릅니다.
# ---------------------------------------------------------------------------
_KOR_RE  = re.compile(r'[가-힣ᄀ-ᇿ㄰-㆏]')  # 한글 (완성형+자모+호환자모)
_JPN_RE  = re.compile(r'[぀-ゟ゠-ヿ]')               # 일본어 히라가나/가타카나
_ENG_RE  = re.compile(r'[A-Za-z]')                                    # 영문 알파벳
_MATH_RE = re.compile(                                                 # 수학 기호 모음
    r'[∀-⋿'   # Mathematical Operators (∀, ∃, ∈ 등)
    r'⨀-⫿'    # Supplemental Math Operators (추가 수학 연산자)
    r'Ͱ-Ͽ'    # Greek Letters (α, β, γ 등 수식에 자주 쓰이는 그리스 문자)
    r'∫∞∑∏∂∇≠≤≥]'
    # ∫(적분), ∞(무한대), ∑(시그마), ∏(파이), ∂(편미분), ∇(나블라), ≠(부등호), ≤, ≥
)


# ---------------------------------------------------------------------------
# PDFProfile — 분석 결과를 담는 데이터 컨테이너
# ---------------------------------------------------------------------------
@dataclass
class PDFProfile:
    """PDF 또는 문서 파일 분석 결과를 저장하는 데이터 클래스.

    하나의 문서 파일을 분석한 뒤 얻은 정보를 한 곳에 모아 두는 구조체입니다.
    다운스트림(파이프라인의 다음 단계)에서 이 프로파일을 읽고
    적절한 처리 방식(OCR 여부, 청킹 방식 등)을 결정합니다.

    매개변수:
        source_type (str): 문서 소스 유형. "DIGITAL"(텍스트 레이어 있음) 또는 "SCANNED"(스캔 이미지).
        layout_type (str): 페이지 레이아웃. "SINGLE_COLUMN"(단단) 또는 "MULTI_COLUMN"(다단).
        detected_languages (list[str]): 감지된 언어 코드 목록. 예: ["kor", "eng"]
        has_formula (bool): 수식(수학 기호) 포함 여부.
        page_count (int): 문서의 총 페이지(슬라이드) 수.
        avg_chars_per_page (float): 페이지당 평균 글자 수.
    """
    source_type: str
    layout_type: str
    # field(default_factory=list)는 각 인스턴스가 독립적인 빈 리스트를 갖도록 합니다.
    # 기본값으로 []를 직접 쓰면 모든 인스턴스가 같은 리스트를 공유하는 버그가 생깁니다.
    detected_languages: list[str] = field(default_factory=list)
    has_formula: bool = False
    page_count: int = 0
    avg_chars_per_page: float = 0.0

    def __str__(self) -> str:
        """객체를 사람이 읽기 좋은 문자열로 변환합니다.

        반환값:
            str: "[소스유형] [레이아웃] 언어:xxx | 수식 있음/없음 | N페이지" 형식의 문자열.
        """
        # 감지된 언어가 없으면 "unknown"으로 표시합니다.
        langs = ",".join(self.detected_languages) or "unknown"
        formula = "수식 있음" if self.has_formula else "수식 없음"
        return (
            f"[{self.source_type}] [{self.layout_type}] "
            f"언어:{langs} | {formula} | {self.page_count}페이지"
        )


# ---------------------------------------------------------------------------
# PDFAnalyzer — 문서 분석의 핵심 클래스
# ---------------------------------------------------------------------------
class PDFAnalyzer:
    """PDF, DOCX, PPTX 파일을 분석하여 PDFProfile을 생성하는 분석기 클래스.

    파일 확장자에 따라 자동으로 적절한 분석 메서드를 선택합니다.
    분석 항목: 소스 유형(디지털/스캔), 레이아웃(단단/다단), 언어, 수식 포함 여부.
    """

    def analyze(self, file_path: str) -> PDFProfile:
        """파일 확장자를 보고 알맞은 분석 메서드를 호출하여 PDFProfile을 반환합니다.

        매개변수:
            file_path (str): 분석할 파일의 경로 (PDF, DOCX, PPTX 지원).

        반환값:
            PDFProfile: 문서 분석 결과 객체.
        """
        # Path(file_path).suffix는 파일 확장자를 추출합니다. 예: ".pdf", ".docx"
        # .lower()는 대소문자 차이를 무시하기 위해 소문자로 통일합니다.
        suffix = Path(file_path).suffix.lower()
        if suffix == ".docx":
            return self._analyze_docx(file_path)
        if suffix == ".pptx":
            return self._analyze_pptx(file_path)
        # 위 두 조건에 해당하지 않으면 PDF로 간주합니다.
        return self._analyze_pdf(file_path)

    def _analyze_pdf(self, pdf_path: str) -> PDFProfile:
        """PDF 파일을 4단계로 분석하여 PDFProfile을 반환합니다.

        분석 순서:
          1. 소스 유형 판별 (디지털 텍스트 vs 스캔 이미지)
          2. 언어 감지 (한국어, 영어, 일본어)
          3. 수식 포함 여부 탐지
          4. 페이지 레이아웃 분석 (단단 vs 다단)

        매개변수:
            pdf_path (str): 분석할 PDF 파일 경로.

        반환값:
            PDFProfile: 분석 결과가 담긴 프로파일 객체.
        """
        logger.info("PDF 분석 시작: %s", pdf_path)

        # ── 1. 소스 유형 및 텍스트 추출 ─────────────────────────────────────
        # pdfplumber로 텍스트를 추출해 봅니다.
        # 텍스트가 충분히 많으면 DIGITAL, 거의 없으면 스캔본(SCANNED)으로 판단합니다.
        source_type, full_text, avg_chars, page_count = self._detect_source_type(pdf_path)
        logger.debug("소스 유형: %s (평균 %.1f자/페이지)", source_type, avg_chars)

        # ── 2. 언어 감지 ────────────────────────────────────────────────────
        # 전체 텍스트가 매우 길 수 있으므로 앞 3000자만 샘플로 사용해 속도를 높입니다.
        sample_text = full_text[:3000] if full_text else ""
        languages = self._detect_languages(sample_text, fallback_path=pdf_path)
        logger.debug("감지된 언어: %s", languages)

        # ── 3. 수식 감지 ────────────────────────────────────────────────────
        # 우선 추출된 텍스트에서 수학 기호를 검색합니다.
        has_formula = self._detect_formula_digital(full_text)
        if source_type == SourceType.SCANNED and not has_formula:
            # 스캔본은 텍스트가 없으므로 이미지 자체에서 수식 특징(분수선 등)을 찾습니다.
            images = self._get_page_images(pdf_path, max_pages=3)
            has_formula = self._detect_formula_scanned(images)
        logger.debug("수식 포함: %s", has_formula)

        # ── 4. 레이아웃 분석 ─────────────────────────────────────────────────
        # 디지털 PDF는 단어 위치 데이터로, 스캔본은 이미지 분석으로 레이아웃을 판별합니다.
        if source_type == SourceType.DIGITAL:
            layout_type = self._detect_layout_digital(pdf_path)
        else:
            images = self._get_page_images(pdf_path, max_pages=3)
            layout_type = self._detect_layout_scanned(images)
        logger.debug("레이아웃: %s", layout_type)

        # 분석 결과를 PDFProfile 객체에 담아 반환합니다.
        profile = PDFProfile(
            source_type=source_type,
            layout_type=layout_type,
            detected_languages=languages,
            has_formula=has_formula,
            page_count=page_count,
            avg_chars_per_page=avg_chars,
        )
        logger.info("PDF 분석 완료: %s", profile)
        return profile

    def _analyze_docx(self, file_path: str) -> PDFProfile:
        """Word 문서(.docx)를 분석하여 PDFProfile을 반환합니다.

        Word 문서는 텍스트 레이어가 항상 존재하므로 소스 유형은 DOCX로 고정되고,
        레이아웃은 SINGLE_COLUMN으로 가정합니다.

        매개변수:
            file_path (str): 분석할 .docx 파일 경로.

        반환값:
            PDFProfile: 분석 결과가 담긴 프로파일 객체.
        """
        logger.info("DOCX 분석 시작: %s", file_path)
        try:
            from docx import Document  # python-docx 라이브러리로 Word 문서를 읽습니다.
            doc = Document(file_path)
            # 모든 단락(문단) 텍스트를 줄바꿈으로 이어 붙여 전체 텍스트를 만듭니다.
            # p.text.strip()이 빈 문자열이면 건너뜁니다(빈 줄 제외).
            full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            # Word에서 '섹션' 수를 페이지 수의 근사값으로 사용합니다.
            # 실제 페이지 수와 다를 수 있지만 Word API로 정확한 페이지 수를 구하기 어렵습니다.
            page_count = max(1, len(doc.sections))
        except Exception as e:
            logger.warning("DOCX 열기 실패 [%s]: %s", file_path, e)
            # 파일 열기에 실패하면 빈 텍스트와 페이지 수 1로 기본값을 설정합니다.
            full_text = ""
            page_count = 1

        languages   = self._detect_languages(full_text[:3000])
        has_formula = self._detect_formula_digital(full_text)
        profile = PDFProfile(
            source_type=SourceType.DOCX,
            layout_type=LayoutType.SINGLE_COLUMN,
            detected_languages=languages,
            has_formula=has_formula,
            page_count=page_count,
            avg_chars_per_page=len(full_text) / page_count,
        )
        logger.info("DOCX 분석 완료: %s", profile)
        return profile

    def _analyze_pptx(self, file_path: str) -> PDFProfile:
        """PowerPoint 파일(.pptx)을 분석하여 PDFProfile을 반환합니다.

        각 슬라이드의 모든 텍스트 박스(도형)에서 텍스트를 추출하여 분석합니다.
        소스 유형은 PPTX, 레이아웃은 SINGLE_COLUMN으로 고정합니다.

        매개변수:
            file_path (str): 분석할 .pptx 파일 경로.

        반환값:
            PDFProfile: 분석 결과가 담긴 프로파일 객체.
        """
        logger.info("PPTX 분석 시작: %s", file_path)
        try:
            from pptx import Presentation  # python-pptx 라이브러리로 PPT 파일을 읽습니다.
            prs = Presentation(file_path)
            texts = []
            # 각 슬라이드를 순회하며 텍스트를 담고 있는 도형(shape)을 찾습니다.
            for slide in prs.slides:
                for shape in slide.shapes:
                    # has_text_frame이 True인 도형만 텍스트를 포함합니다.
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            texts.append(t)
            full_text  = "\n".join(texts)
            # 슬라이드 수를 페이지 수로 사용합니다.
            page_count = max(1, len(prs.slides))
        except Exception as e:
            logger.warning("PPTX 열기 실패 [%s]: %s", file_path, e)
            full_text  = ""
            page_count = 1

        languages   = self._detect_languages(full_text[:3000])
        has_formula = self._detect_formula_digital(full_text)
        profile = PDFProfile(
            source_type=SourceType.PPTX,
            layout_type=LayoutType.SINGLE_COLUMN,
            detected_languages=languages,
            has_formula=has_formula,
            page_count=page_count,
            avg_chars_per_page=len(full_text) / page_count,
        )
        logger.info("PPTX 분석 완료: %s", profile)
        return profile

    # ------------------------------------------------------------------
    # 1. 소스 유형 판별
    # ------------------------------------------------------------------
    def _detect_source_type(
        self, pdf_path: str
    ) -> tuple[str, str, float, int]:
        """pdfplumber로 텍스트를 추출하고 페이지당 평균 글자 수로 소스 유형을 판별합니다.

        텍스트 레이어가 있는 디지털 PDF는 글자가 많이 추출되고,
        스캔된 이미지 PDF는 글자가 거의 추출되지 않는 점을 이용합니다.

        매개변수:
            pdf_path (str): 분석할 PDF 파일 경로.

        반환값:
            tuple:
                - source_type (str): "DIGITAL" 또는 "SCANNED"
                - full_text (str): 모든 페이지에서 추출한 텍스트를 합친 문자열
                - avg_chars_per_page (float): 페이지당 평균 글자 수
                - page_count (int): 총 페이지 수
        """
        texts: list[str] = []
        try:
            with pdfplumber.open(pdf_path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    # extract_text()는 페이지에서 텍스트를 추출합니다.
                    # 텍스트가 없으면 None을 반환하므로 "or """로 빈 문자열로 대체합니다.
                    t = page.extract_text() or ""
                    texts.append(t)
        except Exception as e:
            logger.warning("pdfplumber 열기 실패 [%s]: %s", pdf_path, e)
            # 파일 열기 자체가 실패하면 스캔본으로 간주하고 빈 결과를 반환합니다.
            return SourceType.SCANNED, "", 0.0, 0

        full_text = "\n".join(texts)
        total_chars = sum(len(t) for t in texts)
        # 페이지 수가 0이면 ZeroDivisionError가 발생하므로 조건부로 처리합니다.
        avg_chars = total_chars / page_count if page_count else 0.0

        # 설정값(config.ocr_min_text_length) 이상의 글자가 있으면 디지털 PDF로 판정합니다.
        source_type = (
            SourceType.DIGITAL
            if avg_chars >= config.ocr_min_text_length
            else SourceType.SCANNED
        )
        return source_type, full_text, avg_chars, page_count

    # ------------------------------------------------------------------
    # 2. 언어 감지
    # ------------------------------------------------------------------
    def _detect_languages(self, text: str, fallback_path: str = "") -> list[str]:
        """유니코드 문자 비율 분석과 langdetect 보정을 결합하여 언어 목록을 반환합니다.

        1차로 유니코드 범위별 문자 비율을 계산하고,
        2차로 langdetect 라이브러리 결과로 누락된 언어를 보완합니다.

        매개변수:
            text (str): 언어를 감지할 텍스트.
            fallback_path (str): 텍스트가 부족할 때 참고할 파일 경로 (파일명에서 언어 추측).

        반환값:
            list[str]: 감지된 언어 코드 목록. 예: ["kor", "eng"]
                       언어를 감지하지 못하면 ["eng"]를 기본값으로 반환합니다.
        """
        if not text or len(text) < 20:
            # 텍스트가 너무 짧으면(스캔 PDF 등) 파일명에 한글이 포함되어 있는지 확인합니다.
            # 파일명이 한글이면 한국어 문서일 가능성이 높습니다.
            if fallback_path and _KOR_RE.search(Path(fallback_path).name):
                return ["kor", "eng"]
            return ["eng"]  # 판단 근거가 없으면 영어로 기본 설정합니다.

        # ZeroDivisionError 방지를 위해 최솟값을 1로 보장합니다.
        total = max(len(text), 1)
        languages: list[str] = []

        # ── 1차 판정: 유니코드 범위별 문자 비율 계산 ──────────────────────────
        # 전체 텍스트에서 각 언어 문자가 차지하는 비율이 임계값(3%) 이상이면 해당 언어로 인정합니다.
        if len(_KOR_RE.findall(text)) / total >= LANG_THRESHOLD:
            languages.append("kor")
        if len(_JPN_RE.findall(text)) / total >= LANG_THRESHOLD:
            languages.append("jpn")
        if len(_ENG_RE.findall(text)) / total >= LANG_THRESHOLD:
            languages.append("eng")

        # ── 2차 보정: langdetect 라이브러리로 1차에서 놓친 언어 보완 ──────────
        # langdetect는 통계 기반이라 짧은 텍스트에서 불안정하므로 앞 1000자만 사용합니다.
        try:
            detected = detect_langs(text[:1000])
            for lang_prob in detected:
                code = lang_prob.lang  # 언어 코드 (예: "ko", "en", "ja")
                prob = lang_prob.prob  # 해당 언어일 확률 (0.0 ~ 1.0)
                # 확률이 30% 이상이고 아직 목록에 없는 언어만 추가합니다.
                if code == "ko" and prob > 0.3 and "kor" not in languages:
                    languages.append("kor")
                elif code == "ja" and prob > 0.3 and "jpn" not in languages:
                    languages.append("jpn")
                elif code == "en" and prob > 0.3 and "eng" not in languages:
                    languages.append("eng")
        except LangDetectException:
            # 텍스트가 너무 짧거나 판별 불가 상황에서 발생하는 예외를 조용히 무시합니다.
            pass

        # 최종적으로 아무 언어도 감지되지 않으면 영어를 기본값으로 반환합니다.
        return languages if languages else ["eng"]

    # ------------------------------------------------------------------
    # 3. 수식 감지
    # ------------------------------------------------------------------
    def _detect_formula_digital(self, text: str) -> bool:
        """텍스트에서 수학 기호의 밀도를 계산하여 수식 포함 여부를 판별합니다.

        매개변수:
            text (str): 수식을 탐지할 텍스트.

        반환값:
            bool: 수식이 포함되어 있으면 True, 없으면 False.
        """
        if not text:
            return False
        # _MATH_RE에 매칭되는 수학 기호 문자를 모두 찾아 개수를 셉니다.
        math_chars = len(_MATH_RE.findall(text))
        # 수학 기호 수를 전체 텍스트 길이로 나눠 밀도를 구합니다.
        density = math_chars / max(len(text), 1)
        # 밀도가 임계값(0.8%) 이상이면 수식이 있다고 판정합니다.
        return density >= FORMULA_CHAR_THRESHOLD

    def _detect_formula_scanned(self, images: list) -> bool:
        """스캔 이미지에서 수식 특징(분수선 등 수평 세선, 고밀도 소형 윤곽)을 탐지합니다.

        스캔본은 텍스트 추출이 불가능하므로 이미지 처리로 수식의 시각적 특징을 찾습니다.
        주요 단서: 분수의 가로 줄(분수선)은 매우 얇고 넓은 수평선으로 나타납니다.

        매개변수:
            images (list): PIL 이미지 객체의 리스트 (각 항목이 한 페이지).

        반환값:
            bool: 1페이지 이상에서 수식 특징이 감지되면 True, 아니면 False.
        """
        formula_page_count = 0
        for img_pil in images:
            # PIL 이미지를 NumPy 배열로 변환하고, 그레이스케일(흑백)로 바꿉니다.
            img = np.array(img_pil.convert("L"))
            # Otsu 이진화: 픽셀값을 자동으로 임계값을 정해 흑(글자)/백(배경)으로 분리합니다.
            # THRESH_BINARY_INV는 글자(어두운 부분)를 흰색(255)으로 반전시킵니다.
            _, binary = cv2.threshold(img, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

            # 이진화된 이미지에서 연결된 픽셀 덩어리(윤곽선)를 모두 찾습니다.
            contours, _ = cv2.findContours(binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            h_img, w_img = binary.shape
            formula_hints = 0
            for cnt in contours:
                # 각 윤곽선을 감싸는 최소 직사각형의 위치(x,y)와 크기(w,h)를 구합니다.
                x, y, w, h = cv2.boundingRect(cnt)
                aspect_ratio = w / max(h, 1)      # 가로:세로 비율 (클수록 납작한 모양)
                relative_w = w / w_img             # 페이지 전체 너비 대비 이 윤곽의 너비 비율

                # 분수선 판별 조건:
                # - h <= 3: 높이가 3픽셀 이하 (매우 얇은 선)
                # - aspect_ratio > 10: 가로가 세로의 10배 이상 (납작한 수평선 모양)
                # - relative_w > 0.02: 페이지 너비의 2% 이상 (너무 짧은 선은 노이즈로 제외)
                if h <= 3 and aspect_ratio > 10 and relative_w > 0.02:
                    formula_hints += 1

            # 한 페이지에서 분수선 후보가 3개 이상이면 수식이 있다고 판단합니다.
            if formula_hints >= 3:
                formula_page_count += 1

        # 1페이지라도 수식이 감지되면 True를 반환합니다.
        return formula_page_count >= 1

    # ------------------------------------------------------------------
    # 4. 레이아웃 분석
    # ------------------------------------------------------------------
    def _detect_layout_digital(self, pdf_path: str) -> str:
        """디지털 PDF에서 단어 위치 데이터를 이용해 단단/다단 레이아웃을 판별합니다.

        알고리즘 개요:
          - 페이지를 200개의 수직 슬라이스(bin)로 나눕니다.
          - 각 bin에 텍스트 줄이 걸쳐 있는 비율(커버리지)을 계산합니다.
          - 페이지 중앙(35~65%) 구간에서 커버리지가 낮은 공간(갭)이 있으면 다단으로 판정합니다.
          - 제목처럼 전폭(full-width)으로 걸쳐 있는 줄이 있어도, 다단 본문 줄에 의해
            갭이 희석되지 않도록 줄 단위 비율을 사용합니다.

        BINS=200 (0.5% 해상도)로 좁은 컬럼 간격(~9pt)도 탐지할 수 있습니다.

        매개변수:
            pdf_path (str): 분석할 PDF 파일 경로.

        반환값:
            str: "MULTI_COLUMN" 또는 "SINGLE_COLUMN"
        """
        multi_count = 0       # 다단으로 판정된 페이지 수
        analyzed_pages = 0    # 실제 분석한 페이지 수 (텍스트가 충분한 페이지만 카운트)
        BINS    = 200   # x축 분할 수 (0.5% 단위) — 좁은 갭 탐지에 충분한 해상도
        LINE_H  = 3     # 같은 줄로 묶는 y 좌표 허용 오차 (단위: pt)
        MIN_GAP = 3     # 컬럼 갭으로 인정하는 최소 너비 (bin 단위, 약 1.5% of 페이지)

        try:
            with pdfplumber.open(pdf_path) as pdf:
                # 처음 5페이지만 분석합니다. 문서 전체를 처리하면 너무 느리기 때문입니다.
                for page in pdf.pages[:5]:
                    words = page.extract_words()
                    # 단어가 6개 미만인 페이지는 그림/표 전용 페이지일 가능성이 높아 건너뜁니다.
                    if len(words) < 6:
                        continue
                    page_width = float(page.width)

                    # ── 줄별 x 범위 수집 ─────────────────────────────────
                    # 각 단어를 y 좌표 기준으로 같은 줄에 묶고,
                    # 각 줄에서 단어들이 차지하는 x 범위(bin 단위)를 기록합니다.
                    y_ranges: dict = defaultdict(list)
                    for w in words:
                        # y 좌표를 LINE_H(3pt) 단위로 반올림하여 같은 줄의 단어를 하나로 묶습니다.
                        yk = round(float(w["top"]) / LINE_H) * LINE_H
                        # 단어의 x 시작/끝 좌표를 0~199 범위의 bin 인덱스로 변환합니다.
                        x0b = max(0, int(float(w["x0"]) / page_width * BINS))
                        x1b = min(BINS - 1, int(float(w["x1"]) / page_width * BINS))
                        y_ranges[yk].append((x0b, x1b))

                    n_lines = len(y_ranges)
                    # 줄이 10개 미만이면 참고문헌, 그림 등 본문이 아닐 가능성이 높아 건너뜁니다.
                    if n_lines < 10:
                        continue
                    analyzed_pages += 1

                    # ── 줄 단위 커버리지 계산: 각 bin을 덮는 줄의 비율 ──────────
                    # line_cov[i]는 i번째 bin에 텍스트가 있는 줄이 전체 줄 중 몇 %인지를 나타냅니다.
                    line_cov = np.zeros(BINS)
                    for ranges in y_ranges.values():
                        # 이 줄에서 텍스트가 있는 bin을 표시합니다.
                        covered = np.zeros(BINS, dtype=bool)
                        for x0b, x1b in ranges:
                            covered[x0b: x1b + 1] = True
                        line_cov += covered  # 커버된 줄 수를 누적합니다.

                    line_cov /= n_lines  # 전체 줄 수로 나눠 0~1 비율로 정규화합니다.

                    # ── 중앙 35~65% 구간에서 커버리지가 낮은 공간(갭) 탐색 ────
                    # 2단 레이아웃의 컬럼 간격은 항상 페이지 중앙에 위치합니다.
                    s = int(BINS * 0.35)  # 중앙 구간 시작 bin 인덱스 (35%)
                    e = int(BINS * 0.65)  # 중앙 구간 끝 bin 인덱스 (65%)
                    center_cov = line_cov[s:e]
                    # 좌우 텍스트 영역의 평균 커버리지를 기준값으로 사용합니다.
                    side_mean  = float(
                        (line_cov[:s].mean() + line_cov[e:].mean()) / 2
                    )

                    # 좌우 텍스트 밀도가 너무 낮은 페이지는 신뢰할 수 없어 건너뜁니다.
                    if side_mean <= 0.10:
                        continue

                    # 갭 판정 임계값: 좌우 평균의 30% 이하이거나 절댓값 25% 이하 중 더 낮은 값을 사용합니다.
                    # 이렇게 하면 텍스트 밀도에 비례해 갭 기준이 자동으로 조정됩니다.
                    gap_thresh = min(0.25, side_mean * 0.30)
                    max_gap = curr = 0
                    # 중앙 구간에서 gap_thresh 미만인 bin이 연속으로 몇 개 이어지는지 셉니다.
                    for v in (center_cov < gap_thresh):
                        curr = int(curr + 1) if v else 0
                        if curr > max_gap:
                            max_gap = curr

                    logger.debug(
                        "레이아웃 디지털 p%d: n_lines=%d, side_mean=%.2f, "
                        "gap_thresh=%.2f, max_gap=%d (min=%d)",
                        analyzed_pages, n_lines, side_mean,
                        gap_thresh, max_gap, MIN_GAP,
                    )

                    # 연속 갭이 MIN_GAP(3 bin = 약 9pt) 이상이면 다단 페이지로 기록합니다.
                    if max_gap >= MIN_GAP:
                        multi_count += 1

        except Exception as e:
            logger.warning("레이아웃 분석 실패: %s", e)

        # 분석한 페이지의 2/3 이상이 다단이면 전체 문서를 다단으로 판정합니다.
        threshold = max(1, analyzed_pages * 2 // 3)
        return LayoutType.MULTI_COLUMN if multi_count >= threshold else LayoutType.SINGLE_COLUMN

    def _detect_layout_scanned(self, images: list) -> str:
        """스캔 이미지에서 수직 투영(vertical projection)으로 단단/다단 레이아웃을 판별합니다.

        알고리즘 개요:
          - 각 x열의 픽셀 합(수직 투영)을 구하면, 텍스트가 있는 열은 값이 높고
            컬럼 사이 빈 공간은 값이 낮게 나타납니다.
          - 스무딩(이동 평균)으로 글자 간 짧은 공백 노이즈를 제거합니다.
          - 페이지 중앙 35~65% 구간에 깊은 골(valley)이 있으면 다단으로 판정합니다.

        매개변수:
            images (list): PIL 이미지 객체의 리스트 (각 항목이 한 페이지).

        반환값:
            str: "MULTI_COLUMN" 또는 "SINGLE_COLUMN"
        """
        multi_count = 0   # 다단으로 판정된 페이지 수
        analyzed = 0      # 실제 분석 완료한 페이지 수
        for img_pil in images:
            # PIL 이미지를 그레이스케일 NumPy 배열로 변환합니다.
            img = np.array(img_pil.convert("L"))
            # Otsu 이진화로 글자(흰색=255)와 배경(검정=0)으로 분리합니다.
            _, binary = cv2.threshold(img, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

            # 수직 투영: 각 x열(열 방향)의 픽셀 값을 모두 더합니다.
            # 결과는 이미지 너비 길이의 1차원 배열입니다.
            col_sums = binary.sum(axis=0).astype(float)
            # 모든 열이 0이면(완전히 빈 페이지) 분석할 수 없으므로 건너뜁니다.
            if col_sums.max() == 0:
                continue

            analyzed += 1
            col_sums /= col_sums.max()  # 최댓값으로 나눠 0~1 범위로 정규화합니다.

            w = len(col_sums)
            # 글자 간 짧은 공백 노이즈 제거를 위해 이동 평균(moving average) 스무딩을 적용합니다.
            # 커널 크기를 페이지 너비의 4%로 설정합니다 (약 33px for 826px wide image).
            k = max(1, int(w * 0.04))
            # np.convolve: 신호(col_sums)와 균등 커널을 합성곱하여 이동 평균을 구합니다.
            smoothed = np.convolve(col_sums, np.ones(k) / k, mode="same")
            # 수치 안정성을 위해 최댓값이 0에 너무 가까우면 1e-9로 대체합니다.
            smoothed /= max(smoothed.max(), 1e-9)

            # 페이지 중앙 35~65% 구간에서 최솟값(골 깊이)을 찾습니다.
            s = int(w * 0.35)
            e = int(w * 0.65)
            valley_depth = float(smoothed[s:e].min())  # 중앙 구간의 최솟값 (골 깊이)
            # 좌우 텍스트 영역의 평균 투영값을 기준값으로 사용합니다.
            side_mean = float((smoothed[:s].mean() + smoothed[e:].mean()) / 2)

            # 다단 판정 조건:
            # - side_mean > 0.01: 좌우에 텍스트가 어느 정도 있어야 함 (빈 페이지 제외)
            # - valley_depth < side_mean * 0.50: 중앙 골 깊이가 좌우 평균의 50% 미만
            #   → 컬럼 사이 공간이 텍스트 영역보다 확연히 비어 있음을 의미합니다.
            if side_mean > 0.01 and valley_depth < side_mean * 0.50:
                multi_count += 1

        # 분석한 페이지의 2/3 이상이 다단이면 전체 문서를 다단으로 판정합니다.
        threshold = max(1, analyzed * 2 // 3)
        return LayoutType.MULTI_COLUMN if multi_count >= threshold else LayoutType.SINGLE_COLUMN

    # ------------------------------------------------------------------
    # 유틸 — 페이지 이미지 변환 (pypdfium2 사용, poppler 불필요)
    # ------------------------------------------------------------------
    def _get_page_images(self, pdf_path: str, max_pages: int = 3) -> list:
        """PDF의 각 페이지를 PIL 이미지 객체로 변환하여 리스트로 반환합니다.

        pypdfium2 라이브러리를 사용하므로 poppler 등 외부 프로그램 설치 없이 동작합니다.
        스캔본 분석(수식 탐지, 레이아웃 분석)에 사용됩니다.

        매개변수:
            pdf_path (str): 이미지로 변환할 PDF 파일 경로.
            max_pages (int): 변환할 최대 페이지 수. 기본값은 3.

        반환값:
            list: PIL.Image 객체의 리스트. 변환 실패 시 빈 리스트를 반환합니다.
        """
        images = []
        try:
            pdf = pdfium.PdfDocument(pdf_path)
            # 요청한 max_pages와 실제 페이지 수 중 작은 값만큼만 처리합니다.
            n = min(max_pages, len(pdf))
        except Exception as e:
            logger.warning("PDF 열기 실패: %s", e)
            return images
        for i in range(n):
            try:
                page = pdf[i]
                # scale=1.5는 약 108 DPI로 렌더링합니다 (기본 72 DPI × 1.5).
                # 해상도가 높을수록 이미지 분석 정확도가 높아지지만 메모리와 속도가 증가합니다.
                bitmap = page.render(scale=1.5)
                # pypdfium2의 bitmap 객체를 PIL Image로 변환합니다.
                images.append(bitmap.to_pil())
            except Exception as e:
                logger.warning("페이지 %d 이미지 변환 실패 (건너뜀): %s", i, e)
        return images


# ---------------------------------------------------------------------------
# 싱글턴 인스턴스 — 모듈 임포트 시 한 번만 생성되어 전역에서 재사용됩니다.
# 다른 모듈에서 `from src.pdf_analyzer import analyzer`로 바로 사용할 수 있습니다.
# ---------------------------------------------------------------------------
analyzer = PDFAnalyzer()
