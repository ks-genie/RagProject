# -*- coding: utf-8 -*-
"""
OCR Engine — 전략별 Tesseract 실행 및 품질 점수 산출

전략에 따라 Layout Splitter / Formula Handler와 연동하여
최적의 OCR 파이프라인을 실행.

주요 기능:
  - 디지털 PDF 텍스트 직접 추출 (단일/다단 컬럼)
  - 스캔 이미지 PDF Tesseract OCR (단일/다단/수식 포함)
  - EasyOCR 기반 한국어/영어/일본어 인식
  - LLM(Claude/Gemini)을 이용한 OCR 결과 후교정
  - 워터마크 필터링 (색상·회전 기반)
  - Office 문서(DOCX/PPTX) 텍스트 추출
"""

import logging
import os
import re
from collections import defaultdict

import cv2
import numpy as np
import pdfplumber
import pypdfium2 as pdfium
import pytesseract
from PIL import Image as PILImage

from src.config import config
from src.database import OcrStrategy, Status
from src.formula_handler import FormulaHandler
from src.layout_splitter import LayoutSplitter
from src.ocr_postprocess import fix_easyocr_korean
from src.strategy_selector import get_tesseract_config

# 모듈 전체에서 사용할 로거 설정
logger = logging.getLogger(__name__)

# 컬럼 병합 시 우 컬럼 앞에서 건너뛸 그림/표 캡션 패턴
# 예: "그림 1.", "표 2-3:", "Figure 1.2:", "Fig. 3:" 등의 형태를 감지
_COL_CAP_RE = re.compile(
    r"^(?:그림|표|Figure|Fig\.?|Table)\s*\d+(?:[-\.]\d+)*(?![.-]\d)\s*[\.:]",
    re.IGNORECASE,
)
# 좌 컬럼 끝 바로 다음에 오는 약어 정의 패턴: (HBM3), (DQS), (CLK) 등 대문자+숫자 약어
# 이 패턴이 나타나면 앞 문장과 이어지는 것으로 판단해 병합한다.
_COL_ABBREV_RE = re.compile(r"^\([A-Z][A-Z0-9]{1,10}\)")

# 유효 문자 패턴 (한글·영문·일문·숫자·공백·기본 구두점)
# OCR 품질 점수 계산 시, 아래 패턴에 매칭되는 문자만 '유효한 문자'로 간주한다.
_VALID_CHAR_RE = re.compile(
    r'[가-힣'    # 한글 (가~힣)
    r'぀-ヿ'     # 히라가나·가타카나 (일본어)
    r'A-Za-z0-9'         # 영문 대소문자 및 숫자
    r'\s.,!?;:()\-\'"·]' # 공백과 기본 구두점
)


class OcrEngine:
    """OCR 파이프라인 전체를 관리하는 핵심 클래스.

    전략(strategy) 값에 따라 적합한 OCR 방법을 선택하여 실행하고,
    추출된 텍스트와 품질 점수를 반환한다.

    사용 예:
        engine = OcrEngine()
        text, score = engine.run("document.pdf", "OCR_KOR_ENG")
    """

    def __init__(self):
        """OcrEngine 초기화: Tesseract 설정, 레이아웃 분리기, 수식 처리기를 준비한다."""
        self._setup_tesseract()
        # 페이지 레이아웃(단일/다단 컬럼)을 감지·분리하는 객체
        self.splitter = LayoutSplitter()
        # 수식 영역을 감지하고 분리하는 객체
        self.formula_handler = FormulaHandler()
        # run() 호출 후 페이지별 텍스트를 저장하는 리스트 (외부에서 참조 가능)
        self.last_page_texts: list[str] = []
        # EasyOCR Reader를 언어 조합별로 캐싱하는 딕셔너리 (lang_key → easyocr.Reader)
        # 모델 재로딩 비용이 크기 때문에 한 번 로딩한 모델을 재사용한다.
        self._easy_readers: dict[tuple, object] = {}

    def _setup_tesseract(self):
        """config에서 Tesseract 실행 파일 경로 및 tessdata 폴더를 환경에 적용한다.

        config에 값이 설정된 경우에만 동작하며, 미설정 시 시스템 기본값을 사용한다.
        """
        # Tesseract 실행 파일 경로를 pytesseract에 직접 지정
        if config.ocr_tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = config.ocr_tesseract_cmd
        # 학습 데이터(언어팩) 경로를 환경변수로 설정
        if config.ocr_tessdata_dir:
            os.environ["TESSDATA_PREFIX"] = config.ocr_tessdata_dir

    # ------------------------------------------------------------------
    # EasyOCR
    # ------------------------------------------------------------------
    def _get_easy_reader(self, langs: list[str]):
        """EasyOCR Reader를 언어 조합별로 캐싱하여 반환 (모델 재로딩 방지).

        같은 언어 조합의 Reader는 처음 한 번만 생성되며, 이후 요청 시
        캐싱된 객체를 그대로 반환한다. GPU가 사용 가능하면 자동으로 활성화한다.

        Args:
            langs: 인식할 언어 코드 목록. 예: ["ko", "en"], ["ko", "en", "ja"]

        Returns:
            easyocr.Reader: 초기화된 EasyOCR Reader 객체
        """
        import easyocr
        import torch
        # 언어 목록을 정렬하여 순서 차이 없이 동일한 키로 캐싱
        key = tuple(sorted(langs))
        if key not in self._easy_readers:
            # CUDA GPU 사용 가능 여부 자동 감지
            gpu = torch.cuda.is_available()
            logger.info("EasyOCR Reader 초기화: langs=%s, gpu=%s", langs, gpu)
            self._easy_readers[key] = easyocr.Reader(langs, gpu=gpu, verbose=False)
        return self._easy_readers[key]

    def _easyocr_pages(self, images: list[PILImage.Image], langs: list[str]) -> list[str]:
        """EasyOCR로 페이지별 텍스트를 추출하여 반환한다.

        paragraph=True 옵션으로 줄들을 단락 단위로 묶은 뒤,
        단락 목록을 빈 줄로 연결하여 하나의 페이지 텍스트를 만든다.
        추출 후 한국어 EasyOCR 특유의 오인식 패턴을 후처리로 교정한다.

        Args:
            images: 페이지별 PIL 이미지 목록
            langs:  인식할 언어 코드 목록. 예: ["ko", "en"]

        Returns:
            list[str]: 페이지별 추출 텍스트 목록 (페이지 수와 동일한 길이)
        """
        reader = self._get_easy_reader(langs)
        pages: list[str] = []
        for page_idx, img in enumerate(images):
            try:
                # PIL 이미지를 NumPy 배열로 변환하여 EasyOCR에 전달
                result = reader.readtext(np.array(img), detail=0, paragraph=True)
                # 단락 목록을 빈 줄(\n\n)로 구분하여 연결
                text = "\n\n".join(result) if result else ""
                # EasyOCR 한국어 오인식 패턴 후처리 (예: 'ㅇ' → 'O' 등)
                pages.append(fix_easyocr_korean(text))
            except Exception as e:
                logger.warning("EasyOCR 페이지 %d 실패 (빈 결과): %s", page_idx + 1, e)
                pages.append("")
        return pages

    def _llm_correct_pages(self, pages: list[str]) -> list[str]:
        """LLM(Claude 또는 Gemini)으로 EasyOCR 결과를 페이지별로 교정한다.

        API 키가 미설정되어 있거나 교정 호출이 실패하는 경우 원본 텍스트를 그대로 유지한다.

        Args:
            pages: 교정 전 페이지별 텍스트 목록

        Returns:
            list[str]: 교정 완료된 페이지별 텍스트 목록
        """
        from src.llm_corrector import correct_ocr_page
        # config에서 사용할 LLM 공급자 선택 (gemini 또는 claude)
        provider = config.llm_provider
        if provider == "gemini":
            api_key = config.gemini_api_key
            model   = config.gemini_model
            timeout = config.gemini_timeout
            max_tokens = 8192
        else:
            api_key = config.claude_api_key
            model   = config.claude_model
            timeout = config.claude_timeout
            max_tokens = config.claude_max_tokens

        # API 키가 없으면 교정을 건너뛰고 원본 반환
        if not api_key:
            logger.warning("LLM(%s) API 키 미설정 — 교정 생략", provider)
            return pages

        corrected: list[str] = []
        for i, text in enumerate(pages):
            # 각 페이지를 LLM에게 보내 교정 요청
            result = correct_ocr_page(
                text,
                provider=provider,
                api_key=api_key,
                model=model,
                max_tokens=max_tokens,
                timeout=timeout,
            )
            # 교정 결과가 원본과 다른 경우에만 로그 기록
            if result != text:
                logger.info(
                    "LLM(%s) 교정 완료: 페이지 %d (%d→%d자)",
                    provider, i + 1, len(text), len(result),
                )
            corrected.append(result)
        return corrected

    # ------------------------------------------------------------------
    # 공개 인터페이스
    # ------------------------------------------------------------------
    def run(self, pdf_path: str, strategy: str) -> tuple[str, float]:
        """OCR 실행의 진입점. 전략에 맞는 추출 방법을 선택하여 실행한다.

        strategy 값에 따라 디지털 추출, 스캔 OCR, EasyOCR 등
        서로 다른 파이프라인이 선택된다.
        스캔 전략은 모두 _ocr_adaptive를 통해
        페이지별 레이아웃을 감지하여 다단/단일을 자동 전환한다.

        Args:
            pdf_path: OCR을 수행할 파일 경로 (PDF, DOCX, PPTX 등)
            strategy: OcrStrategy 상수값. 예: "DIGITAL_EXTRACT", "OCR_KOR_ENG"

        Returns:
            tuple[str, float]:
                - str:   추출된 전체 텍스트 (페이지 구분: 빈 줄 두 개)
                - float: 품질 점수 0.0 ~ 1.0 (유효 문자 비율)
        """
        logger.info("OCR 시작: strategy=%s [%s]", strategy, pdf_path)

        try:
            if strategy == OcrStrategy.DIGITAL_EXTRACT:
                # 일반 디지털 PDF — 텍스트 레이어를 직접 추출 (단일 컬럼)
                pages = self._extract_digital_pages(pdf_path)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            elif strategy == OcrStrategy.DIGITAL_EXTRACT_MULTI:
                # 2단 레이아웃 디지털 PDF — 컬럼을 분리하여 좌→우 순서로 추출
                pages = self._extract_digital_multicol_pages(pdf_path)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            elif strategy == OcrStrategy.DOCX_EXTRACT:
                # Word 문서(.docx) 텍스트 추출
                pages = self._extract_docx_pages(pdf_path)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            elif strategy == OcrStrategy.PPTX_EXTRACT:
                # PowerPoint 문서(.pptx) 슬라이드별 텍스트 추출
                pages = self._extract_pptx_pages(pdf_path)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            elif strategy == OcrStrategy.OCR_KOR_ENG:
                # EasyOCR으로 한국어+영어 인식
                images = self._get_page_images(pdf_path)
                pages = self._easyocr_pages(images, ["ko", "en"])
                # LLM 교정이 활성화된 경우 추가로 교정 수행
                if config.llm_correction_enabled:
                    logger.info("LLM(%s) OCR 교정 시작: %d 페이지", config.llm_provider, len(pages))
                    pages = self._llm_correct_pages(pages)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            elif strategy == OcrStrategy.OCR_JPN:
                # EasyOCR으로 한국어+영어+일본어 인식
                images = self._get_page_images(pdf_path)
                pages = self._easyocr_pages(images, ["ko", "en", "ja"])
                # LLM 교정이 활성화된 경우 추가로 교정 수행
                if config.llm_correction_enabled:
                    logger.info("LLM(%s) OCR 교정 시작: %d 페이지", config.llm_provider, len(pages))
                    pages = self._llm_correct_pages(pages)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            else:
                # 그 외 전략: Tesseract 기반 적응형 OCR
                # 수식 포함 전략인지 여부 확인
                images = self._get_page_images(pdf_path)
                tess_cfg = get_tesseract_config(strategy)
                with_formula = strategy in (
                    OcrStrategy.OCR_FORMULA,
                    OcrStrategy.OCR_MULTI_FORMULA,
                )
                # 페이지별로 레이아웃을 자동 감지하여 OCR 수행
                pages = self._ocr_adaptive_pages(images, tess_cfg,
                                                 with_formula=with_formula)
                text = "\n\n".join(pages)
                score = self._compute_quality(text)
            # 마지막으로 처리한 페이지별 텍스트를 저장 (외부 참조용)
            self.last_page_texts = pages
        except Exception as e:
            logger.error("OCR 실패: %s", e)
            raise

        logger.info("OCR 완료: 글자 수=%d, 품질=%.3f", len(text), score)
        return text, score

    def meets_quality(self, text: str, score: float) -> bool:
        """OCR 결과가 최소 품질 기준을 충족하는지 확인한다.

        텍스트 길이가 최소 기준 이상이고, 품질 점수가 임계값 이상일 때만 True를 반환한다.

        Args:
            text:  OCR로 추출한 텍스트
            score: _compute_quality()로 계산한 품질 점수 (0.0 ~ 1.0)

        Returns:
            bool: 품질 기준 충족 시 True, 미충족 시 False
        """
        return (
            len(text) >= config.ocr_min_text_length
            and score >= config.ocr_quality_threshold
        )

    # ------------------------------------------------------------------
    # 워터마크 필터 (색상·회전 기반)
    # ------------------------------------------------------------------
    @staticmethod
    def _luminance(color) -> float | None:
        """PDF 색상값을 밝기(luminance) 값으로 변환한다.

        Grayscale, RGB, CMYK 등 다양한 색상 포맷을 지원하며,
        변환할 수 없는 포맷이면 None을 반환한다.

        Args:
            color: pdfplumber에서 추출한 색상값.
                   None, 단일 숫자, 또는 1·3·4 원소 배열/튜플.

        Returns:
            float | None: 밝기 값 (0.0=검정, 1.0=흰색), 변환 불가 시 None
        """
        if color is None:
            return None
        # 단일 숫자인 경우 그대로 밝기값으로 사용
        if isinstance(color, (int, float)):
            return float(color)
        if not isinstance(color, (list, tuple)):
            return None
        # 그레이스케일: 채널 1개
        if len(color) == 1:
            return float(color[0])
        # RGB: 인간 시각에 맞는 가중 평균 공식 사용
        if len(color) == 3:
            r, g, b = (float(v) for v in color)
            return 0.299 * r + 0.587 * g + 0.114 * b
        # CMYK: 먼저 RGB로 변환 후 밝기 계산
        if len(color) == 4:
            c, m, y, k = (float(v) for v in color)
            r = (1 - c) * (1 - k)
            g = (1 - m) * (1 - k)
            b = (1 - y) * (1 - k)
            return 0.299 * r + 0.587 * g + 0.114 * b
        return None

    @classmethod
    def filter_watermarks(
        cls,
        page,
        light_threshold: float = 0.80,
        stamp_size_pt: float = 30.0,
        stamp_lum_min: float = 0.15,
    ):
        """워터마크로 의심되는 문자를 제거한 pdfplumber 페이지 객체를 반환한다.

        세 가지 기준으로 워터마크를 판별하고 필터링한다:
          ① upright=False: 회전된 텍스트 (대각선 스탬프 등)
          ② 채움색 또는 획 색상의 밝기가 light_threshold(기본 0.80) 초과:
             반투명하거나 연한 배경 워터마크
          ③ 폰트 크기 > stamp_size_pt AND 밝기 > stamp_lum_min:
             CTM 행렬 회전으로 렌더링된 대형 스탬프의 개별 글자 (예: 'T', 'S')

        Args:
            page:             pdfplumber 페이지 객체
            light_threshold:  밝기 임계값. 이 값 초과 시 연한 색상으로 판단 (기본: 0.80)
            stamp_size_pt:    대형 스탬프 판별 기준 폰트 크기(pt) (기본: 30.0)
            stamp_lum_min:    대형 스탬프 판별 기준 최소 밝기 (기본: 0.15)

        Returns:
            pdfplumber 페이지 객체 (워터마크 문자가 필터링된 상태)
        """
        def _keep(obj):
            # 문자(char) 유형이 아닌 객체(선, 사각형 등)는 그대로 유지
            if obj.get("object_type") != "char":
                return True
            # ① 텍스트 행렬 회전 → 워터마크로 판단하여 제거
            if not obj.get("upright", True):
                return False
            lum_fill   = cls._luminance(obj.get("non_stroking_color"))
            lum_stroke = cls._luminance(obj.get("stroking_color"))
            # ② 연한 색상 (채움 또는 획)이면 워터마크로 판단
            if lum_fill   is not None and lum_fill   > light_threshold:
                return False
            if lum_stroke is not None and lum_stroke > light_threshold:
                return False
            # ③ 큰 폰트 + 비검정 색상 → CTM 회전 스탬프 낱글자로 판단
            size = float(obj.get("size") or 0)
            lum  = lum_fill if lum_fill is not None else lum_stroke
            if size > stamp_size_pt and lum is not None and lum > stamp_lum_min:
                return False
            return True
        # pdfplumber의 filter()에 판별 함수를 전달하여 조건에 맞는 객체만 유지
        return page.filter(_keep)

    # ------------------------------------------------------------------
    # Office 문서 추출 — DOCX / PPTX
    # ------------------------------------------------------------------
    def _extract_docx_pages(self, file_path: str) -> list[str]:
        """python-docx 라이브러리로 Word 문서(.docx)에서 텍스트를 추출한다.

        문서 전체를 단일 페이지로 반환하며, 단락, 표, 텍스트박스를 순서대로 처리한다.
        python-docx가 열지 못하는 파일(HTML로 저장된 .docx 등)은
        HTML 태그 제거 방식으로 폴백 추출을 시도한다.

        Args:
            file_path: Word 문서 파일 경로

        Returns:
            list[str]: 추출된 텍스트를 담은 단일 원소 리스트
        """
        from docx import Document
        try:
            doc = Document(file_path)
            parts: list[str] = []
            # 일반 단락 텍스트 추출
            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
            # 표(table) 텍스트 추출 — 각 행의 셀을 '|'로 구분
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            # 텍스트박스 등 python-docx가 누락하는 요소: XML w:t 태그를 직접 순회
            # python-docx API가 지원하지 않는 일부 요소를 보완하기 위한 방법
            if not parts:
                ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                for elem in doc.element.body.iter(f'{{{ns}}}t'):
                    t = (elem.text or '').strip()
                    if t:
                        parts.append(t)
            return ["\n\n".join(parts)] if parts else [""]
        except Exception as first_err:
            # ZIP(DOCX) 형식이 아닌 경우 — Word HTML export 등 비표준 포맷 처리
            logger.warning("DOCX ZIP 열기 실패, HTML 폴백 시도 [%s]: %s", file_path, first_err)
            return self._extract_html_as_docx_pages(file_path)

    def _extract_html_as_docx_pages(self, file_path: str) -> list[str]:
        """확장자가 .docx이지만 실제로는 HTML 형식인 파일에서 텍스트를 추출한다.

        Word에서 'HTML로 저장(mso- CSS 속성 포함)'한 파일만 처리하며,
        일반 웹페이지가 .docx로 저장된 경우에는 거부하여 빈 결과를 반환한다.
        HTML 태그와 엔티티를 제거하여 순수 텍스트만 추출한다.

        Args:
            file_path: 실제로는 HTML인 .docx 파일 경로

        Returns:
            list[str]: 추출된 텍스트를 담은 단일 원소 리스트, 실패 시 [""]
        """
        import re
        try:
            with open(file_path, encoding='utf-8', errors='replace') as f:
                raw = f.read()

            # Word HTML export 여부 판별 — mso- CSS 속성, Microsoft Word 언급, WordDocument 태그 중 하나라도 있으면 Word HTML로 판단
            is_word_html = bool(
                re.search(r'mso-', raw, re.IGNORECASE) or
                re.search(r'Microsoft\s+Word', raw, re.IGNORECASE) or
                re.search(r'WordDocument', raw)
            )
            if not is_word_html:
                logger.warning(
                    "웹페이지 HTML로 판단됨 — 문서 내용 추출 불가 (Word HTML export가 아님): [%s]",
                    file_path,
                )
                return [""]

            # CSS 스타일 시트와 JavaScript 블록을 먼저 제거 (텍스트 노이즈 방지)
            raw = re.sub(r'<style[^>]*>.*?</style>', ' ', raw,
                         flags=re.DOTALL | re.IGNORECASE)
            raw = re.sub(r'<script[^>]*>.*?</script>', ' ', raw,
                         flags=re.DOTALL | re.IGNORECASE)
            # 남은 HTML 태그 제거
            text = re.sub(r'<[^>]+>', ' ', raw)
            # 기본 HTML 엔티티를 실제 문자로 변환
            for ent, ch in [('&nbsp;', ' '), ('&lt;', '<'), ('&gt;', '>'),
                            ('&amp;', '&'), ('&quot;', '"')]:
                text = text.replace(ent, ch)
            # 연속 공백·탭을 단일 공백으로 정규화
            text = re.sub(r'[ \t]+', ' ', text)
            # 3개 이상 연속 줄바꿈을 2개로 통일
            text = re.sub(r'\n{3,}', '\n\n', text).strip()
            if text:
                logger.info("Word HTML 폴백 추출 완료: %d자 [%s]", len(text), file_path)
                return [text]
        except Exception as e:
            logger.error("HTML 폴백 추출 실패 [%s]: %s", file_path, e)
        return [""]

    def _extract_pptx_pages(self, file_path: str) -> list[str]:
        """python-pptx로 PowerPoint 슬라이드별 텍스트를 추출한다.

        각 슬라이드에서 텍스트 프레임과 표(table)를 추출하며,
        슬라이드 하나를 페이지 하나로 취급한다.

        Args:
            file_path: PowerPoint 파일 경로

        Returns:
            list[str]: 슬라이드별 추출 텍스트 목록
        """
        from pptx import Presentation
        try:
            prs = Presentation(file_path)
            pages: list[str] = []
            for slide in prs.slides:
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # 텍스트 프레임을 가진 도형(텍스트박스, 제목 등)에서 텍스트 추출
                        t = shape.text_frame.text.strip()
                        if t:
                            slide_texts.append(t)
                    elif shape.has_table:
                        # 표(table) 도형에서 셀 텍스트를 '|'로 구분하여 추출
                        for row in shape.table.rows:
                            cells = [c.text_frame.text.strip() for c in row.cells
                                     if c.text_frame.text.strip()]
                            if cells:
                                slide_texts.append(" | ".join(cells))
                pages.append("\n\n".join(slide_texts))
            return pages if pages else [""]
        except Exception as e:
            logger.error("PPTX 추출 실패 [%s]: %s", file_path, e)
            return [""]

    # ------------------------------------------------------------------
    # Digital 추출 — 단일 컬럼
    # ------------------------------------------------------------------
    def _extract_digital_pages(self, pdf_path: str) -> list[str]:
        """디지털 PDF(텍스트 레이어 존재)에서 단일 컬럼 기준으로 페이지별 텍스트를 추출한다.

        워터마크를 필터링한 후, 단어 좌표 정보를 바탕으로 단락 구조를 재구성한다.

        Args:
            pdf_path: PDF 파일 경로

        Returns:
            list[str]: 페이지별 추출 텍스트 목록
        """
        pages = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                # 워터마크 문자 제거 후 단어 목록 추출
                page  = self.filter_watermarks(page)
                words = page.extract_words()
                pages.append(
                    # 단어 좌표를 이용해 단락 구조를 재구성
                    self._words_to_text_with_paragraphs(words) if words
                    # 단어가 없으면 pdfplumber 기본 텍스트 추출로 폴백
                    else (page.extract_text() or "")
                )
        return pages

    def _extract_digital(self, pdf_path: str) -> tuple[str, float]:
        """디지털 PDF 단일 컬럼 추출 후 (전체 텍스트, 품질 점수) 튜플을 반환한다."""
        pages = self._extract_digital_pages(pdf_path)
        full_text = "\n\n".join(pages)
        return full_text, self._compute_quality(full_text)

    # ------------------------------------------------------------------
    # Digital 추출 — 다단 컬럼 인식
    # ------------------------------------------------------------------
    def _extract_digital_multicol_pages(self, pdf_path: str) -> list[str]:
        """2단 레이아웃 디지털 PDF를 컬럼별로 분리하여 순서대로 텍스트를 추출한다.

        각 페이지에서 컬럼 간격(gap)을 자동으로 찾아 좌 컬럼 → 우 컬럼 순으로 읽는다.
        컬럼이 감지되지 않으면 단일 컬럼으로 처리한다.

        Args:
            pdf_path: PDF 파일 경로

        Returns:
            list[str]: 페이지별 추출 텍스트 목록
        """
        pages = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page  = self.filter_watermarks(page)
                words = page.extract_words()
                if not words:
                    pages.append("")
                    continue
                page_width = float(page.width)
                # 컬럼 사이 빈 공간(gap)의 x 좌표 중심을 탐지
                gap_x = self._find_column_gap(words, page_width)
                if gap_x is None:
                    # 컬럼 구분이 없으면 단일 컬럼으로 처리
                    pages.append(self._words_to_text_with_paragraphs(words))
                else:
                    # 본문 시작 y 좌표 탐지 (헤더/제목 영역 이후를 본문으로 간주)
                    body_start_y = self._find_body_start_y(
                        words, gap_x, page_height=float(page.height)
                    )
                    pages.append(self._compose_multicol_text(words, gap_x, body_start_y))
        return pages

    def _extract_digital_multicol(self, pdf_path: str) -> tuple[str, float]:
        """다단 컬럼 디지털 PDF 추출 후 (전체 텍스트, 품질 점수) 튜플을 반환한다."""
        pages = self._extract_digital_multicol_pages(pdf_path)
        full_text = "\n\n".join(pages)
        return full_text, self._compute_quality(full_text)

    @staticmethod
    def _find_column_gap(words: list, page_width: float) -> float | None:
        """줄 단위 x 커버리지 분석으로 2단 컬럼 사이의 빈 공간 중심 x 좌표를 찾는다.

        페이지를 200개 구간(bins)으로 나누어 각 줄이 각 구간을 덮는 비율을 계산한다.
        페이지 중앙 30~70% 구간에서 커버리지가 낮은 연속 구간이 컬럼 갭 후보다.
        3개 이상의 줄에서 연속 3개 이상의 빈 구간이 발견될 때만 갭으로 인정한다.

        Args:
            words:      pdfplumber에서 추출한 단어 목록
            page_width: 페이지 너비(pt)

        Returns:
            float | None: 컬럼 갭 중심 x 좌표(pt), 갭이 없으면 None
        """
        BINS   = 200  # 페이지를 몇 개 구간으로 나눌지 (해상도)
        LINE_H = 3    # 줄 그룹핑 간격 (픽셀/pt 단위)

        lines = OcrEngine._group_words_by_line(words, line_height=LINE_H)
        # 양쪽 컬럼 모두에 단어가 있는 줄을 우선 후보로 선택
        candidate_lines = [
            line_words
            for _, line_words in lines
            if OcrEngine._line_spans_both_columns(line_words, page_width * 0.5)
        ]
        # 후보 줄이 3개 미만이면 모든 줄을 대상으로 분석
        selected_lines = (
            candidate_lines
            if len(candidate_lines) >= 3
            else [line_words for _, line_words in lines]
        )

        n_lines = len(selected_lines)
        if n_lines < 3:
            return None

        # 각 구간별로 해당 구간을 덮는 줄의 비율을 계산
        line_cov = np.zeros(BINS)
        for line_words in selected_lines:
            covered = np.zeros(BINS, dtype=bool)
            for w in line_words:
                # 단어의 좌우 경계를 bin 인덱스로 변환
                x0b = max(0, int(float(w["x0"]) / page_width * BINS))
                x1b = min(BINS - 1, int(float(w["x1"]) / page_width * BINS))
                covered[x0b: x1b + 1] = True
            line_cov += covered
        # 비율(0~1)로 정규화
        line_cov /= n_lines

        # 페이지 중앙 30~70% 영역만 갭 탐색 범위로 사용
        s, e = int(BINS * 0.35), int(BINS * 0.65)
        center_cov = line_cov[s:e]
        # 좌우 사이드의 평균 커버리지 — 텍스트가 충분히 있어야 분석 의미 있음
        side_mean  = float((line_cov[:s].mean() + line_cov[e:].mean()) / 2)
        if side_mean <= 0.10:
            # 텍스트가 거의 없으면 다단 레이아웃으로 볼 수 없음
            return None

        # 갭 임계값: 사이드 평균 커버리지의 30% 또는 0.25 중 작은 값
        gap_thresh = min(0.25, side_mean * 0.30)
        max_gap = curr = best_start = best_end = 0
        # 슬라이딩 윈도우로 가장 긴 연속 빈 구간을 탐색
        for i, v in enumerate(center_cov < gap_thresh):
            curr = curr + 1 if v else 0
            if curr > max_gap:
                max_gap   = curr
                best_end   = s + i
                best_start = best_end - max_gap + 1

        # 3개 미만의 연속 빈 구간은 갭으로 보기 어려움
        if max_gap < 3:
            return None

        # 갭 구간의 중심 bin 인덱스를 실제 x 좌표(pt)로 변환
        gap_center_bin = (best_start + best_end + 1) / 2
        return gap_center_bin / BINS * page_width

    @staticmethod
    def _group_words_by_line(words: list, line_height: int = 3) -> list[tuple[float, list]]:
        """단어 목록을 y 좌표 기준으로 줄(행) 단위로 그룹화한다.

        top 좌표를 line_height 단위로 반올림하여 같은 줄의 단어를 묶는다.
        아래첨자/위첨자처럼 y 간격이 작고 x가 인접한 단어는 이전 줄에 합친다.

        Args:
            words:       pdfplumber에서 추출한 단어 목록
            line_height: 줄 그룹핑 반올림 단위 (기본: 3)

        Returns:
            list[tuple[float, list]]: (y 좌표, 해당 줄 단어 목록) 튜플의 정렬된 리스트
        """
        grouped: dict[float, list] = defaultdict(list)
        for w in words:
            # top 좌표를 line_height 단위로 반올림하여 같은 그룹에 모음
            yk = round(float(w["top"]) / line_height) * line_height
            grouped[yk].append(w)
        groups = sorted(grouped.items())
        if len(groups) < 2:
            return groups

        # 아래첨자/위첨자 병합:
        # y 간격이 line_height*3 미만이고 x 좌표가 이전 줄의 특정 단어에 인접하면 합침.
        # 전체 줄 범위가 아닌 개별 단어 기준 인접성을 확인하여 오병합 방지.
        merged: list[list] = [[groups[0][0], list(groups[0][1])]]
        for yk, grp_words in groups[1:]:
            prev_yk, prev_words = merged[-1]
            gap = yk - prev_yk
            if gap < line_height * 3:
                sub: list = []
                rest: list = []
                for cw in grp_words:
                    cx0 = float(cw.get("x0", 0))
                    # 아래첨자/위첨자 판별: 현재 단어의 왼쪽 끝이 이전 줄 특정 단어의 오른쪽 끝에 인접
                    is_sub = any(
                        -1.0 <= cx0 - float(pw.get("x1", pw.get("x0", 0))) <= 1.5
                        for pw in prev_words
                    )
                    if is_sub:
                        sub.append(cw)
                    else:
                        rest.append(cw)
                if sub:
                    # 첨자 단어를 이전 줄에 합침
                    merged[-1][1].extend(sub)
                if rest:
                    # 첨자가 아닌 나머지는 새 줄로 추가
                    merged.append([yk, rest])
            else:
                merged.append([yk, list(grp_words)])
        return [(yk, ws) for yk, ws in merged]

    @staticmethod
    def _join_line_words(words: list) -> str:
        """한 줄 내 단어들을 x 좌표 순서로 연결하여 텍스트를 만든다.

        단어 간격이 1.0pt 미만인 경우(아래첨자·위첨자 등 밀착된 경우)는 공백 없이 붙이고,
        그 이상이면 단어 사이에 공백을 삽입한다.

        Args:
            words: 한 줄에 해당하는 단어 목록

        Returns:
            str: 연결된 줄 텍스트
        """
        # x0 기준으로 왼쪽에서 오른쪽 순서로 정렬
        sorted_w = sorted(words, key=lambda w: float(w.get("x0", 0)))
        if not sorted_w:
            return ""
        result = sorted_w[0]["text"]
        for i in range(1, len(sorted_w)):
            # 이전 단어 끝(x1)과 현재 단어 시작(x0) 사이의 간격 계산
            gap = float(sorted_w[i].get("x0", 0)) - float(sorted_w[i - 1].get("x1", sorted_w[i - 1].get("x0", 0)))
            # 1.0pt 미만 간격: 밀착(첨자 등) → 공백 없이 연결
            # 1.0pt 이상: 일반 단어 간격 → 공백 삽입
            result += sorted_w[i]["text"] if gap < 1.0 else (" " + sorted_w[i]["text"])
        return result

    @staticmethod
    def _line_spans_both_columns(line_words: list, split_x: float) -> bool:
        """한 줄의 단어들이 좌우 컬럼 양쪽에 모두 걸쳐 있는지 확인한다.

        단어 중심 x 좌표가 split_x보다 작으면 좌 컬럼,
        크거나 같으면 우 컬럼으로 분류하여 양쪽에 단어가 있는지 판단한다.

        Args:
            line_words: 한 줄에 해당하는 단어 목록
            split_x:    좌/우 컬럼 구분 기준 x 좌표

        Returns:
            bool: 양쪽 컬럼에 모두 단어가 있으면 True
        """
        has_left = False
        has_right = False
        for w in line_words:
            mid_x = (float(w["x0"]) + float(w["x1"])) / 2
            if mid_x < split_x:
                has_left = True
            else:
                has_right = True
            if has_left and has_right:
                return True
        return False

    @staticmethod
    def _split_header_footer_words(
        words: list,
        page_height: float,
        header_pct: float = 0.08,
        footer_pct: float = 0.93,
    ) -> tuple[list, list, list]:
        """단어 목록을 y 좌표 기준으로 헤더·본문·푸터 세 구역으로 분리한다.

        페이지 상단 8%(기본)까지를 헤더, 하단 7%(기본)를 푸터,
        나머지를 본문으로 분류한다.

        Args:
            words:      분리할 단어 목록
            page_height: 페이지 높이(pt)
            header_pct: 헤더 높이 비율 (기본: 0.08 = 상단 8%)
            footer_pct: 푸터 시작 높이 비율 (기본: 0.93 = 하단 7%)

        Returns:
            tuple[list, list, list]: (헤더 단어, 본문 단어, 푸터 단어)
        """
        h = page_height
        header = [w for w in words if float(w.get("top", 0)) < h * header_pct]
        footer = [w for w in words if float(w.get("top", 0)) > h * footer_pct]
        main   = [w for w in words
                  if h * header_pct <= float(w.get("top", 0)) <= h * footer_pct]
        return header, main, footer

    @staticmethod
    def _split_words_by_column(words: list, gap_x: float) -> tuple[list, list]:
        """단어 목록을 gap_x를 기준으로 좌 컬럼과 우 컬럼으로 분리한다.

        단어 중심 x 좌표가 gap_x보다 작으면 좌 컬럼, 크거나 같으면 우 컬럼.

        Args:
            words: 분리할 단어 목록
            gap_x: 좌/우 컬럼 구분 기준 x 좌표(pt)

        Returns:
            tuple[list, list]: (좌 컬럼 단어 목록, 우 컬럼 단어 목록)
        """
        left = []
        right = []
        for w in words:
            mid_x = (float(w["x0"]) + float(w["x1"])) / 2
            if mid_x < gap_x:
                left.append(w)
            else:
                right.append(w)
        return left, right

    @classmethod
    def _find_body_start_y(
        cls,
        words: list,
        gap_x: float,
        page_height: float = 0.0,
    ) -> float | None:
        """양쪽 컬럼에 단어가 고르게 분포하는 본문 시작 y 좌표를 탐지한다.

        탐지 기준: 각 줄에서 좌·우 양쪽에 모두 단어가 2개 이상인 줄이
        연속 3줄 이상 이어지는 첫 번째 y 좌표를 본문 시작으로 판단한다.

        줄 간격이 GAP_THRESHOLD(30pt)를 초과하면 섹션 구분으로 보아 카운터를 리셋한다.
        이렇게 하면 제목(2줄)·저자(2줄) 같은 짧은 섹션은 본문으로 인식되지 않는다.

        Args:
            words:       분석할 단어 목록
            gap_x:       좌/우 컬럼 구분 기준 x 좌표
            page_height: 페이지 높이(pt). 미전달 시 words에서 추정.
                         상단 10% 이하 영역(런닝헤더)은 탐색에서 제외된다.

        Returns:
            float | None: 본문 시작 y 좌표(pt), 탐지 실패 시 None
        """
        lines = cls._group_words_by_line(words)
        if not lines:
            return None

        # page_height가 전달되지 않으면 단어들의 최대 bottom 좌표로 추정
        if not page_height and words:
            page_height = max(float(w.get("bottom", w.get("top", 0))) for w in words)
        min_y = page_height * 0.10  # 상단 10% = 런닝헤더 영역 → 탐색 제외

        WINDOW          = 3   # 연속으로 만족해야 할 줄 수
        MIN_EACH        = 2   # 각 줄에서 좌·우 각각 최소 단어 수
        GAP_THRESHOLD   = 30  # 이 pt를 초과하는 줄 간격은 섹션 구분으로 간주

        consecutive: int = 0          # 조건을 만족하는 연속 줄 수
        first_y: float | None = None  # 연속 구간의 첫 번째 y 좌표
        prev_yk: float | None = None  # 이전 줄의 y 좌표 (간격 계산용)

        for yk, line_words in lines:
            # 런닝헤더 영역은 건너뜀
            if yk < min_y:
                prev_yk = yk
                consecutive = 0
                first_y = None
                continue

            # 줄 간격이 GAP_THRESHOLD 초과 → 새 섹션 시작으로 간주하여 카운터 리셋
            if prev_yk is not None and (yk - prev_yk) > GAP_THRESHOLD:
                consecutive = 0
                first_y = None
            prev_yk = yk

            # 각 단어의 중심 x를 기준으로 좌·우 컬럼 단어 수 계산
            mid_xs = [(float(w["x0"]) + float(w["x1"])) / 2 for w in line_words]
            left_n  = sum(1 for x in mid_xs if x < gap_x)
            right_n = sum(1 for x in mid_xs if x >= gap_x)

            if left_n >= MIN_EACH and right_n >= MIN_EACH:
                # 양쪽 컬럼 조건 충족
                consecutive += 1
                if first_y is None:
                    first_y = yk
                if consecutive >= WINDOW:
                    # WINDOW개 연속 줄 조건을 모두 충족하면 첫 번째 y를 반환
                    return first_y
            else:
                # 조건 미충족 시 카운터 리셋
                consecutive = 0
                first_y = None

        return None

    @staticmethod
    def _join_columns(left_text: str, right_text: str) -> str:
        """좌 컬럼과 우 컬럼 텍스트를 올바른 순서와 구조로 연결한다.

        좌 컬럼의 마지막 단락이 문장 중간에서 끊긴 경우,
        우 컬럼의 첫 단락과 이어 붙인다 (2단 레이아웃에서 문장이 컬럼 경계를 넘는 경우).

        병합 조건:
          A. 우 첫 단락이 소문자로 시작 (원래 문장이 이어짐)
          B. 좌 마지막 단어가 문장 중단 단어이고, 우 컬럼 앞에 그림/표 캡션이 선행
             → 캡션을 건너뛰어 그 다음 단락과 병합
          C. 우 첫 단락이 (HBM3) 같은 대문자 약어 정의로 시작
        하이픈으로 끝나는 경우 하이픈을 제거하고 이어붙인다.

        Args:
            left_text:  좌 컬럼 전체 텍스트 (단락 구분: \\n\\n)
            right_text: 우 컬럼 전체 텍스트 (단락 구분: \\n\\n)

        Returns:
            str: 병합된 전체 텍스트
        """
        # 문장을 끝내는 구두점 집합
        _SENT_END = frozenset('.?!;')
        # 문장 중간에 올 수 있는 단어 (이 단어로 끝나면 다음 컬럼과 이어질 가능성 높음)
        _MID_WORD = frozenset({
            "a", "an", "the",
            "and", "or", "nor", "but", "yet", "so",
            "in", "of", "to", "at", "by", "for", "on", "as", "from",
            "with", "into", "onto", "upon", "over", "under",
            "is", "are", "was", "were", "be", "been", "being",
            "that", "which", "who", "whom", "whose", "when", "where",
            "this", "these", "those", "its", "their", "our", "your",
            "not", "also", "both", "either", "neither", "than", "such",
            "using", "shown", "given", "based", "called", "used",
            "mentioned", "described", "defined",
        })

        # 빈 단락 제거 후 좌·우 컬럼을 단락 목록으로 분리
        left_paras  = [p for p in left_text.strip().split('\n\n')  if p.strip()]
        right_paras = [p for p in right_text.strip().split('\n\n') if p.strip()]
        if not left_paras or not right_paras:
            return "\n\n".join(filter(None, [left_text.strip(), right_text.strip()]))

        last_left = left_paras[-1].rstrip()
        # 좌 컬럼 마지막 단락이 문장 종결 구두점으로 끝나면 이어붙이지 않음
        if not last_left or last_left[-1] in _SENT_END:
            return "\n\n".join(filter(None, [left_text.strip(), right_text.strip()]))

        # 우 컬럼 앞부분의 그림/표 캡션 단락들을 별도로 분리
        right_caps = []
        right_main = list(right_paras)
        while right_main and _COL_CAP_RE.match(right_main[0].strip()):
            right_caps.append(right_main.pop(0))
        # 우 컬럼이 캡션만으로 구성된 경우 캡션을 본문으로 처리
        if not right_main:
            right_main = right_caps
            right_caps = []

        first_right = right_main[0].lstrip()
        # 좌 마지막 단어를 소문자로 변환하고 구두점 제거
        last_word   = last_left.split()[-1].lower().rstrip(".,;:?!") if last_left.split() else ""

        # 세 가지 조건 중 하나라도 충족하면 병합
        should_merge = (
            first_right[:1].islower()                   # A: 우 컬럼 첫 글자가 소문자
            or (right_caps and last_word in _MID_WORD)  # B: 캡션 선행 + 좌 마지막 단어가 중간 단어
            or bool(_COL_ABBREV_RE.match(first_right))  # C: 우 컬럼이 대문자 약어 정의로 시작
        )

        if should_merge:
            # 하이픈으로 끝나면 하이픈 제거 후 공백 없이 연결, 아니면 공백으로 연결
            connector = "" if last_left.endswith('-') else " "
            ll = last_left[:-1] if last_left.endswith('-') else last_left
            merged      = ll + connector + first_right
            merged_left = '\n\n'.join(left_paras[:-1] + [merged])
            # 캡션과 우 컬럼 나머지 단락을 순서대로 이어붙임
            remainder   = '\n\n'.join(right_caps + right_main[1:])
            return "\n\n".join(filter(None, [merged_left, remainder]))

        return "\n\n".join(filter(None, [left_text.strip(), right_text.strip()]))

    @classmethod
    def _compose_multicol_text(
        cls,
        words: list,
        gap_x: float,
        body_start_y: float | None,
        page_height: float = 0.0,
    ) -> str:
        """다단 컬럼 페이지의 단어들을 헤더·본문·각주 순으로 올바르게 재구성한다.

        본문 시작 y(body_start_y)가 없으면 단순히 좌·우 컬럼을 합친다.
        본문 시작 y가 있으면:
          1) 본문 이전(헤더/제목) 영역을 처리 (2단 걸치는 줄은 헤더, 나머지는 좌·우 분리)
          2) 본문 영역을 좌·우 컬럼으로 분리하여 합침
          3) 좌 컬럼 하단 각주(35pt 이상 줄 간격이 있는 구간) 분리 후 마지막에 추가

        Args:
            words:        페이지의 모든 단어 목록
            gap_x:        컬럼 갭 중심 x 좌표(pt)
            body_start_y: 본문 시작 y 좌표(pt), None이면 헤더/본문 구분 없이 처리
            page_height:  페이지 높이(pt). 미전달 시 words에서 추정.

        Returns:
            str: 재구성된 페이지 전체 텍스트
        """
        # body_start_y가 없으면 전체를 바로 좌·우 분리하여 합침
        if body_start_y is None:
            left, right = cls._split_words_by_column(words, gap_x)
            return cls._join_columns(
                cls._words_to_text_with_paragraphs(left),
                cls._words_to_text_with_paragraphs(right),
            )

        header = []
        left_body = []
        right_body = []

        # 본문 이전(헤더/제목) 영역 처리
        pre_body_words = [w for w in words if float(w["top"]) < body_start_y]
        for _, line_words in cls._group_words_by_line(pre_body_words):
            if cls._line_spans_both_columns(line_words, gap_x):
                # 양쪽 컬럼에 걸치는 줄은 헤더(제목, 저자 등)로 분류
                header.extend(line_words)
                continue
            # 한 쪽에만 있는 줄은 해당 컬럼에 배정
            left_line, right_line = cls._split_words_by_column(line_words, gap_x)
            left_body.extend(left_line)
            right_body.extend(right_line)

        # 본문 영역 단어들을 좌·우로 분리
        body_words = [w for w in words if float(w["top"]) >= body_start_y]
        left_words, right_words = cls._split_words_by_column(body_words, gap_x)
        left_body.extend(left_words)
        right_body.extend(right_words)

        # ── 좌 컬럼 하단 각주 분리 ──────────────────────────────────────────
        # 페이지 하단 50% 구간에서 줄 간격 35pt 이상인 첫 지점을 각주 시작으로 판단.
        # 대상: 학술 논문의 소속·수신일 각주 등 컬럼 분리 이후 좌측에만 있는 하단 블록.
        if not page_height and words:
            page_height = max(float(w.get("bottom", w.get("top", 0))) for w in words)

        FOOTNOTE_GAP  = 35   # 이 pt 이상의 줄 간격이 있으면 각주 시작으로 판단
        FOOTNOTE_ZONE = 0.50 # 페이지 하단 50% 영역에서만 각주를 탐색

        note_start_y: float | None = None
        if page_height:
            zone_y = page_height * FOOTNOTE_ZONE
            prev_yk: float | None = None
            for yk, _ in cls._group_words_by_line(left_body):
                # 하단 50% 영역에서 큰 줄 간격이 발견되면 각주 시작으로 표시
                if yk > zone_y and prev_yk is not None and (yk - prev_yk) > FOOTNOTE_GAP:
                    note_start_y = yk
                    break
                prev_yk = yk

        # 각주 시작 y가 탐지되면 좌 컬럼 본문과 각주를 분리
        left_main = left_body
        left_note: list = []
        if note_start_y is not None:
            left_main = [w for w in left_body if float(w.get("top", 0)) < note_start_y]
            left_note  = [w for w in left_body if float(w.get("top", 0)) >= note_start_y]
        # ────────────────────────────────────────────────────────────────────

        header_text = cls._words_to_text_with_paragraphs(header)
        # 좌 컬럼 본문과 우 컬럼 본문을 문장 경계에 맞게 병합
        body_joined = cls._join_columns(
            cls._words_to_text_with_paragraphs(left_main),
            cls._words_to_text_with_paragraphs(right_body),
        )

        # 헤더 → 본문 → 각주 순으로 조합
        parts = [s for s in (header_text, body_joined) if s]
        if left_note:
            note_text = cls._words_to_text_with_paragraphs(left_note).strip()
            if note_text:
                parts.append(note_text)

        return "\n\n".join(parts)

    @classmethod
    def _words_to_text(cls, words: list) -> str:
        """단어 목록을 y → x 순서로 정렬하여 줄 단위 텍스트로 재구성한다.

        단락 구분 없이 단순히 줄(\n)로만 연결한다.
        단락 구분이 필요하면 _words_to_text_with_paragraphs()를 사용할 것.

        Args:
            words: pdfplumber에서 추출한 단어 목록

        Returns:
            str: 재구성된 텍스트
        """
        if not words:
            return ""
        lines = []
        for _, line_words in cls._group_words_by_line(words):
            lines.append(cls._join_line_words(line_words))
        return "\n".join(lines)

    @classmethod
    def _words_to_text_with_paragraphs(
        cls,
        words: list,
        line_height: int = 3,
    ) -> str:
        """단어 좌표 정보를 활용하여 줄 및 단락 구분이 있는 텍스트를 재구성한다.

        줄 간격(gap)을 분석하여 적절한 구분자를 선택한다:
          - 일반 줄 간격 → \\n (줄바꿈)
          - 단락 간격(중앙값의 1.35배 이상) → \\n\\n (단락 구분)
          - 넓은 단락 간격(중앙값의 2.1배 이상) → \\n\\n\\n (섹션 구분)

        Args:
            words:       pdfplumber에서 추출한 단어 목록
            line_height: 줄 그룹핑 반올림 단위 (기본: 3)

        Returns:
            str: 줄/단락 구분이 적용된 재구성 텍스트
        """
        if not words:
            return ""

        # 각 줄의 (시작 y, 끝 y, 텍스트) 튜플 목록 생성
        line_data: list[tuple[float, float, str]] = []
        for _, line_words in cls._group_words_by_line(words, line_height=line_height):
            text = cls._join_line_words(line_words).strip()
            if not text:
                continue
            sorted_line = sorted(line_words, key=lambda w: float(w.get("x0", 0)))
            line_data.append(
                (
                    float(sorted_line[0].get("top", 0)),    # 줄 시작 y (top)
                    max(float(w.get("bottom", w.get("top", 0))) for w in sorted_line),  # 줄 끝 y (bottom)
                    text,
                )
            )

        if not line_data:
            return ""
        if len(line_data) == 1:
            return line_data[0][2]

        # 연속된 줄 사이의 간격(gap) 목록 계산
        # gap = 다음 줄의 top - 현재 줄의 bottom (음수면 0으로 처리)
        gaps = [
            max(line_data[i + 1][0] - line_data[i][1], 0.0)
            for i in range(len(line_data) - 1)
        ]
        positive_gaps = [gap for gap in gaps if gap > 0]
        if not positive_gaps:
            # 모든 간격이 0 이하이면 단순 줄바꿈으로 연결
            return "\n".join(text for _, _, text in line_data)

        # 간격 목록의 중앙값을 기준 줄 간격으로 사용 (이상값에 강건)
        base_gap = sorted(positive_gaps)[len(positive_gaps) // 2]
        # 단락 구분 임계값: 기준 간격의 1.35배 또는 기준+1.0pt 중 큰 값
        paragraph_gap = max(base_gap * 1.35, base_gap + 1.0)
        # 넓은 단락(섹션) 구분 임계값: 기준 간격의 2.1배 또는 단락 임계값+3.0pt 중 큰 값
        wide_paragraph_gap = max(base_gap * 2.1, paragraph_gap + 3.0)

        parts: list[str] = [line_data[0][2]]
        for idx, gap in enumerate(gaps):
            # 간격 크기에 따라 구분자 선택
            if gap >= wide_paragraph_gap:
                separator = "\n\n\n"  # 섹션 구분
            elif gap >= paragraph_gap:
                separator = "\n\n"    # 단락 구분
            else:
                separator = "\n"      # 일반 줄바꿈
            parts.append(separator + line_data[idx + 1][2])

        return "".join(parts)

    # ------------------------------------------------------------------
    # 페이지별 적응형 OCR (메인 경로)
    # ------------------------------------------------------------------
    def _ocr_adaptive_pages(
        self,
        images: list,
        tess_cfg: dict,
        with_formula: bool = False,
    ) -> list[str]:
        """모든 페이지에 대해 레이아웃을 개별 감지하여 최적 OCR을 적용한다.

        각 페이지를 LayoutSplitter로 분석하여 단일/다단 컬럼을 자동 판별하고,
        각 컬럼에 Tesseract OCR을 적용한다.
        수식 포함 모드(with_formula=True)이면 수식 영역을 먼저 분리하여 처리한다.

        Args:
            images:       페이지별 PIL 이미지 목록
            tess_cfg:     Tesseract 설정 딕셔너리 (lang, psm, oem 등)
            with_formula: 수식 영역 분리 처리 여부 (기본: False)

        Returns:
            list[str]: 페이지별 OCR 결과 텍스트 목록
        """
        cfg_str = self._build_tess_config(tess_cfg)
        lang = tess_cfg["lang"]
        pages: list[str] = []

        for page_idx, img in enumerate(images):
            # 페이지 이미지를 컬럼 단위로 분리 (단일이면 1개, 2단이면 2개)
            cols = self.splitter._split_page(img)
            logger.debug(
                "페이지 %d: %s (%d컬럼)",
                page_idx + 1,
                "다단" if len(cols) > 1 else "단일",
                len(cols),
            )
            col_texts: list[str] = []
            for col in cols:
                if with_formula:
                    # 수식 영역을 분리하고 나머지 텍스트 영역만 OCR
                    text_img, regions = self.formula_handler.separate_formula(col)
                    text_img = self._enhance_for_ocr(text_img)
                    t = pytesseract.image_to_string(text_img, lang=lang, config=cfg_str)
                    if regions:
                        # 수식 영역 개수를 텍스트에 주석으로 추가
                        t += f"\n[수식 영역 {len(regions)}개]"
                else:
                    # 이미지를 이진화 전처리 후 Tesseract로 OCR
                    t = pytesseract.image_to_string(self._enhance_for_ocr(col), lang=lang, config=cfg_str)
                col_texts.append(t)
            # 컬럼별 텍스트를 줄바꿈으로 연결하여 한 페이지 텍스트 완성
            pages.append("\n".join(col_texts))

        return pages

    def _ocr_adaptive(
        self,
        images: list,
        tess_cfg: dict,
        with_formula: bool = False,
    ) -> tuple[str, float]:
        """적응형 OCR 실행 후 (전체 텍스트, 품질 점수) 튜플을 반환한다."""
        pages = self._ocr_adaptive_pages(images, tess_cfg, with_formula=with_formula)
        full_text = "\n\n".join(pages)
        return full_text, self._compute_quality(full_text)

    # ------------------------------------------------------------------
    # 단순 OCR (단일 컬럼, 언어별) — 하위 호환용
    # ------------------------------------------------------------------
    def _ocr_simple(self, images: list, tess_cfg: dict) -> tuple[str, float]:
        """단일 컬럼 Tesseract OCR을 수행한다. (적응형 OCR 도입 전 방식, 하위 호환용)

        Args:
            images:   페이지별 PIL 이미지 목록
            tess_cfg: Tesseract 설정 딕셔너리

        Returns:
            tuple[str, float]: (전체 텍스트, 품질 점수)
        """
        cfg_str = self._build_tess_config(tess_cfg)
        lang = tess_cfg["lang"]
        texts = [
            pytesseract.image_to_string(img, lang=lang, config=cfg_str)
            for img in images
        ]
        full_text = "\n\n".join(texts)
        return full_text, self._compute_quality(full_text)

    # ------------------------------------------------------------------
    # 다단 OCR
    # ------------------------------------------------------------------
    def _ocr_multi_column(self, images: list, tess_cfg: dict) -> tuple[str, float]:
        """다단 컬럼 레이아웃에 Tesseract OCR을 적용한다.

        LayoutSplitter로 각 페이지를 컬럼으로 분리한 뒤
        컬럼별로 OCR을 수행하고 결과를 순서대로 합친다.

        Args:
            images:   페이지별 PIL 이미지 목록
            tess_cfg: Tesseract 설정 딕셔너리

        Returns:
            tuple[str, float]: (전체 텍스트, 품질 점수)
        """
        cfg_str = self._build_tess_config(tess_cfg)
        lang = tess_cfg["lang"]
        all_texts = []
        for page_cols in self.splitter.split_columns(images):
            # 각 컬럼에 OCR을 적용하고 줄바꿈으로 연결
            col_texts = [
                pytesseract.image_to_string(col, lang=lang, config=cfg_str)
                for col in page_cols
            ]
            all_texts.append("\n".join(col_texts))
        full_text = "\n\n".join(all_texts)
        return full_text, self._compute_quality(full_text)

    # ------------------------------------------------------------------
    # 수식 포함 OCR
    # ------------------------------------------------------------------
    def _ocr_formula(self, images: list, tess_cfg: dict) -> tuple[str, float]:
        """수식 영역을 분리한 후 텍스트 영역만 Tesseract OCR을 적용한다.

        FormulaHandler로 각 페이지에서 수식 영역을 감지·제거하고,
        나머지 텍스트 영역에 OCR을 수행한다.
        수식이 있으면 텍스트에 수식 영역 개수 주석을 추가한다.

        Args:
            images:   페이지별 PIL 이미지 목록
            tess_cfg: Tesseract 설정 딕셔너리

        Returns:
            tuple[str, float]: (전체 텍스트, 품질 점수)
        """
        cfg_str = self._build_tess_config(tess_cfg)
        lang = tess_cfg["lang"]
        all_texts = []
        for img in images:
            # 수식 영역을 분리 (text_img: 수식 제거된 이미지, regions: 수식 영역 목록)
            text_img, regions = self.formula_handler.separate_formula(img)
            text = pytesseract.image_to_string(text_img, lang=lang, config=cfg_str)
            if regions:
                text += f"\n[수식 영역 {len(regions)}개]"
            all_texts.append(text)
        full_text = "\n\n".join(all_texts)
        return full_text, self._compute_quality(full_text)

    # ------------------------------------------------------------------
    # 다단 + 수식 복합
    # ------------------------------------------------------------------
    def _ocr_multi_formula(self, images: list, tess_cfg: dict) -> tuple[str, float]:
        """다단 컬럼과 수식 분리를 동시에 처리하는 복합 OCR을 수행한다.

        각 페이지를 컬럼으로 분리한 뒤, 각 컬럼에서 수식 영역을 제거하고
        텍스트 영역만 OCR을 적용한다.

        Args:
            images:   페이지별 PIL 이미지 목록
            tess_cfg: Tesseract 설정 딕셔너리

        Returns:
            tuple[str, float]: (전체 텍스트, 품질 점수)
        """
        cfg_str = self._build_tess_config(tess_cfg)
        lang = tess_cfg["lang"]
        all_texts = []
        for page_cols in self.splitter.split_columns(images):
            col_texts = []
            for col in page_cols:
                # 컬럼별 수식 영역 분리 후 OCR
                text_img, regions = self.formula_handler.separate_formula(col)
                text = pytesseract.image_to_string(text_img, lang=lang, config=cfg_str)
                if regions:
                    text += f"\n[수식 영역 {len(regions)}개]"
                col_texts.append(text)
            all_texts.append("\n".join(col_texts))
        full_text = "\n\n".join(all_texts)
        return full_text, self._compute_quality(full_text)

    # ------------------------------------------------------------------
    # 유틸
    # ------------------------------------------------------------------
    def _build_tess_config(self, tess_cfg: dict) -> str:
        """Tesseract 설정 딕셔너리를 커맨드라인 옵션 문자열로 변환한다.

        예: {"psm": 3, "oem": 1} → "--psm 3 --oem 1"

        Args:
            tess_cfg: psm, oem 등 Tesseract 옵션이 담긴 딕셔너리

        Returns:
            str: pytesseract에 전달할 config 문자열
        """
        parts = []
        if "psm" in tess_cfg:
            parts.append(f"--psm {tess_cfg['psm']}")
        if "oem" in tess_cfg:
            parts.append(f"--oem {tess_cfg['oem']}")
        return " ".join(parts)

    @staticmethod
    def _enhance_for_ocr(img_pil: PILImage.Image) -> PILImage.Image:
        """이미지를 적응형 이진화(binarization)로 전처리하여 OCR 인식률을 높인다.

        Gaussian 적응형 임계값 이진화 방식을 사용하여
        스캔 문서의 불균일한 조명이나 그림자에 강건하게 처리한다.
        한글처럼 획이 복잡한 문자의 인식률 개선에 효과적이다.

        Args:
            img_pil: 전처리할 PIL 이미지 (컬러 또는 그레이스케일)

        Returns:
            PILImage.Image: 이진화(흑백)된 PIL 이미지
        """
        # 컬러 이미지를 그레이스케일("L" 모드)로 변환
        g = np.array(img_pil.convert("L"))
        # 적응형 Gaussian 이진화:
        # - 블록 크기 31: 31×31 픽셀 영역의 평균으로 임계값 결정
        # - 상수 10: 평균에서 10을 빼 임계값을 조정 (노이즈 제거 효과)
        b = cv2.adaptiveThreshold(
            g, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 31, 10
        )
        return PILImage.fromarray(b)

    def _compute_quality(self, text: str) -> float:
        """추출된 텍스트의 품질 점수를 유효 문자 비율로 계산한다.

        전체 문자 중 _VALID_CHAR_RE 패턴에 매칭되는 문자(한글, 영문, 숫자, 구두점 등)의
        비율을 품질 점수로 사용한다. OCR 오인식이 많으면 특수문자·잡음이 늘어나
        점수가 낮아지는 원리를 이용한다.

        Args:
            text: 품질을 측정할 텍스트

        Returns:
            float: 품질 점수 0.0 ~ 1.0 (소수점 4자리 반올림)
        """
        if not text.strip():
            return 0.0
        # 유효 문자 패턴에 매칭되는 문자 수 계산
        valid = len(_VALID_CHAR_RE.findall(text))
        # 유효 문자 수 / 전체 문자 수 (0으로 나누기 방지를 위해 max 사용)
        return round(valid / max(len(text), 1), 4)

    def _get_page_images(self, pdf_path: str, dpi_scale: float = 4.17) -> list[PILImage.Image]:
        """pypdfium2 라이브러리를 사용하여 PDF 각 페이지를 이미지로 변환한다.

        poppler 설치 없이 PDF를 이미지로 변환할 수 있는 장점이 있다.
        dpi_scale=4.17은 약 300 DPI에 해당하며, OCR 품질과 처리 속도의 균형점이다.
        변환에 실패한 페이지는 A4 크기의 흰 빈 이미지로 대체한다.

        Args:
            pdf_path:  변환할 PDF 파일 경로
            dpi_scale: 렌더링 배율 (기본: 4.17 ≈ 300 DPI)

        Returns:
            list[PILImage.Image]: 페이지별 PIL 이미지 목록
        """
        images = []
        try:
            pdf = pdfium.PdfDocument(pdf_path)
        except Exception as e:
            logger.error("PDF 열기 실패 [%s]: %s", pdf_path, e)
            return images
        for i in range(len(pdf)):
            try:
                page = pdf[i]
                # scale 값으로 해상도 결정: 72 DPI * 4.17 ≈ 300 DPI
                bitmap = page.render(scale=dpi_scale)
                images.append(bitmap.to_pil())
            except Exception as e:
                logger.warning("페이지 %d 이미지 변환 실패 (빈 페이지 대체): %s", i, e)
                # 변환 실패 시 A4 72dpi×2 기준(1240×1754px) 흰 빈 이미지로 대체
                images.append(PILImage.new("RGB", (1240, 1754), color=(255, 255, 255)))
        return images


# 모듈 수준 싱글턴 인스턴스 — 여러 곳에서 import하여 재사용할 수 있도록 미리 생성
# 새로운 인스턴스가 필요하면 OcrEngine()으로 직접 생성하면 된다.
ocr_engine = OcrEngine()
