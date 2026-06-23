# -*- coding: utf-8 -*-
"""
Asset Pipeline — 표·그림 추출, OCR, DB 저장, 텍스트 참조 태그 삽입

디지털 PDF:
  pdfplumber 단어 좌표 기반으로 자산 위치를 파악하여
  - 표 셀 텍스트를 참조 태그로 교체
  - 그림 위치에 참조 태그 삽입
  - 분리된 문장 재연결

스캔 PDF:
  OCR 텍스트에 캡션 탐지 또는 위치 추정으로 참조 태그 삽입하고
  분리된 문장 재연결.

참조 태그 형식:  [TABLE_001], [FIGURE_001], ...
"""

import io
import logging
import os
import re
from collections import defaultdict
from pathlib import Path

import pdfplumber
import pytesseract
from PIL import Image as PILImage

from src.asset_db import AssetDatabase, asset_db
from src.config import config
from src.figure_extractor import extract_assets
from src.page_db import PageDatabase, page_db

# 이 모듈 전용 로거 생성 (로그 메시지 출처를 모듈명으로 구분)
logger = logging.getLogger(__name__)

# 문장이 끝났음을 나타내는 마침표 종류 모음 (집합 형태로 빠른 검사)
_SENT_END = frozenset(".?!;:。！？")

# 캡션 패턴 (표·그림 캡션 감지) — Fig / Fig. / Figure / 그림 모두 포함
# re.compile: 반복 사용할 정규식을 미리 컴파일해 두면 실행 속도가 빨라짐
_TABLE_CAPTION_RE = re.compile(r"^(?:표|Table)\s*\d+(?:[-\.]\d+)*(?![.-]\d)\s*[\.:]", re.IGNORECASE)
_FIGURE_CAPTION_RE = re.compile(r"^(?:그림|Figure|Fig\.?)\s*\d+(?:[-\.]\d+)*(?![.-]\d)\s*[\.:]", re.IGNORECASE)

# index만 포함된 캡션 라인 판정 (설명 텍스트 없이 번호만 존재)
# 예) "Fig 1", "Fig. 2", "Figure 3-1", "Table 4.", "그림 5"
_INDEX_ONLY_RE = re.compile(
    r"^(?:표|Table|그림|Figure|Fig\.?)\s*[\d][\d\-\.]*\.?\s*$",
    re.IGNORECASE,
)

# 이미 삽입된 참조 태그 감지 (중복 삽입 방지용)
# 예: [TABLE_001], [FIGURE_023] 같은 형태를 찾아냄
_REF_TAG_RE = re.compile(r"\[(?:TABLE|FIGURE)_\d{3}\]")


def _extract_simple_title(body_text: str) -> str:
    """본문 텍스트에서 문서 제목을 추정한다.

    첫 페이지 본문의 첫 번째 비어있지 않은 줄을 제목으로 간주한다.

    Args:
        body_text: 페이지 본문 텍스트 (줄바꿈 포함)

    Returns:
        첫 번째 비어있지 않은 줄 (최대 200자). 없으면 빈 문자열.
    """
    for line in body_text.splitlines():
        stripped = line.strip()
        if stripped:
            return stripped[:200]  # 너무 긴 제목은 200자로 잘라냄
    return ""


def _make_placeholder_png() -> bytes:
    """1×1 픽셀 흰색 PNG 이미지를 바이트로 생성한다.

    이미지 없이 캡션 텍스트만 있는 자산을 DB에 저장할 때
    image_data 컬럼을 채우기 위한 최소 크기의 더미 이미지다.

    Returns:
        PNG 포맷의 바이트 데이터
    """
    buf = io.BytesIO()  # 메모리 버퍼 (파일 대신 메모리에 저장)
    PILImage.new("RGB", (1, 1), (255, 255, 255)).save(buf, format="PNG")
    return buf.getvalue()


# 이미지 없는 캡션 전용 자산 저장 시 사용하는 1×1 흰색 PNG
# 모듈 로드 시 한 번만 생성하여 재사용 (상수처럼 활용)
_PLACEHOLDER_PNG: bytes = _make_placeholder_png()


class AssetPipeline:
    """표·그림 자산을 추출하여 별도 DB에 저장하고, 텍스트에 참조 태그를 삽입한다.

    PDF, DOCX, PPTX 등 다양한 문서 형식을 처리하며,
    디지털 PDF와 스캔 PDF를 각각 다른 방식으로 처리한다.

    Attributes:
        _adb: 자산(표·그림 이미지)을 저장하는 데이터베이스
        _pdb: 페이지 텍스트를 저장하는 데이터베이스
    """

    def __init__(self, adb: AssetDatabase = None, pdb: PageDatabase = None):
        """AssetPipeline 초기화.

        Args:
            adb: 자산 DB 인스턴스. None이면 전역 싱글턴 asset_db 사용.
            pdb: 페이지 DB 인스턴스. None이면 전역 싱글턴 page_db 사용.
        """
        # None이 전달되면 미리 만들어진 전역 DB 객체를 사용 (싱글턴 패턴)
        self._adb = adb or asset_db
        self._pdb = pdb or page_db

    # ──────────────────────────────────────────────────────────────────────
    # 공개 인터페이스
    # ──────────────────────────────────────────────────────────────────────

    def process_digital(self, file_path: str, doc_id: int) -> tuple[str, int]:
        """디지털 문서(PDF/DOCX/PPTX): 자산 추출·저장 + 참조 태그가 삽입된 전체 텍스트 반환.

        기존 DB 데이터를 삭제한 뒤 새로 추출하므로 재처리 시 중복이 없다.

        Args:
            file_path: 처리할 문서 파일 경로
            doc_id: 데이터베이스에서 이 문서를 식별하는 ID

        Returns:
            (text_with_refs, n_assets) 튜플.
            text_with_refs는 헤더·푸터를 제외한 본문 텍스트,
            n_assets는 저장된 자산(표·그림) 총 개수.
        """
        # 재처리 시 이전 자산/페이지 데이터 초기화
        self._adb.delete_assets_for_document(doc_id)
        self._pdb.delete_all(doc_id)

        # 파일 확장자로 문서 종류 판별
        suffix = Path(file_path).suffix.lower()
        if suffix == ".docx":
            return self._process_docx(file_path, doc_id)
        if suffix == ".pptx":
            return self._process_pptx(file_path, doc_id)

        # PDF 처리: 자산 추출 + 본문 구성
        pdf_path = file_path
        page_bodies, page_headers, page_footers, n_assets = \
            self._extract_digital_with_assets(pdf_path, doc_id)

        # 벡터 figure 등 이미지 미탐지 캡션을 텍스트 스캔으로 보완
        # (이미지 탐지기가 놓친 자산을 캡션 문자열 패턴으로 추가 발견)
        counts = self._adb.count_assets(doc_id)
        page_bodies, fig_seq, tbl_seq = self._inject_caption_refs(
            page_bodies, pdf_path, doc_id,
            fig_seq=counts.get("FIGURE", 0),
            tbl_seq=counts.get("TABLE", 0),
        )
        # 캡션 스캔으로 새로 발견된 자산 수를 합산
        n_assets += (fig_seq - counts.get("FIGURE", 0)) + (tbl_seq - counts.get("TABLE", 0))

        # 페이지별 헤더/본문/푸터를 page_contents DB에 저장
        for page_num, (header, body, footer) in enumerate(
            zip(page_headers, page_bodies, page_footers)
        ):
            self._pdb.save_page_content(
                doc_id, page_num,
                header_text=header, body_text=body, footer_text=footer,
            )

        # document_metadata: 제목은 page 0 body 첫 비어있지 않은 줄로 추정
        title = _extract_simple_title(page_bodies[0] if page_bodies else "")
        self._pdb.save_document_metadata(doc_id, title=title)

        # 헤더·푸터를 제외한 본문만 AnythingLLM 업로드용으로 반환
        full_text = "\n\n".join(page_bodies)
        logger.info("디지털 자산 처리 완료: %d건 (doc_id=%d)", n_assets, doc_id)
        return full_text, n_assets

    def process_scanned(
        self,
        pdf_path: str,
        doc_id: int,
        page_texts: list[str],
    ) -> tuple[str, int]:
        """스캔 PDF: 기존 OCR 텍스트에 참조 태그 삽입.

        스캔 PDF는 텍스트가 없으므로 별도 OCR 결과(page_texts)를 받아서
        이미지로 탐지한 자산의 참조 태그를 텍스트에 삽입한다.

        Args:
            pdf_path: 스캔 PDF 파일 경로
            doc_id: 데이터베이스 문서 ID
            page_texts: OCR 결과의 페이지별 텍스트 목록

        Returns:
            (text_with_refs, n_assets) 튜플.
        """
        # 기존 자산 데이터 초기화 (재처리 대비)
        self._adb.delete_assets_for_document(doc_id)

        # 모든 페이지에서 자산 추출 후 DB 저장
        all_assets, n_assets = self._extract_all_assets(pdf_path, doc_id)

        # 탐지된 자산이 있으면 OCR 텍스트에 참조 태그 삽입, 없으면 원본 유지
        modified = (
            self._inject_refs_scanned(page_texts, all_assets)
            if all_assets
            else list(page_texts)
        )

        # 이미지 미탐지 캡션을 텍스트 스캔으로 보완
        counts = self._adb.count_assets(doc_id)
        modified, fig_seq, tbl_seq = self._inject_caption_refs(
            modified, pdf_path, doc_id,
            fig_seq=counts.get("FIGURE", 0),
            tbl_seq=counts.get("TABLE", 0),
        )
        # 캡션 스캔으로 추가 발견된 자산 수 합산
        n_assets += (fig_seq - counts.get("FIGURE", 0)) + (tbl_seq - counts.get("TABLE", 0))

        # page_contents DB 저장 — 웹 UI가 EasyOCR 결과를 표시하려면 반드시 필요
        # (process_digital과 동일 패턴; 스캔 PDF는 header/footer 미분리이므로 body_text만)
        self._pdb.delete_all(doc_id)
        for page_num, body in enumerate(modified):
            self._pdb.save_page_content(doc_id, page_num, body_text=body)

        # 첫 페이지 첫 줄을 문서 제목으로 추정
        title = _extract_simple_title(modified[0] if modified else "")
        self._pdb.save_document_metadata(doc_id, title=title)

        logger.info("스캔 자산 처리 완료: %d건 (doc_id=%d)", n_assets, doc_id)
        return "\n\n".join(modified), n_assets

    # ──────────────────────────────────────────────────────────────────────
    # Office 문서 처리 (DOCX / PPTX)
    # ──────────────────────────────────────────────────────────────────────

    def _process_docx(self, file_path: str, doc_id: int) -> tuple[str, int]:
        """Word 문서(.docx): 본문 텍스트 추출 + 임베디드 이미지 DB 저장.

        문단 텍스트와 표 셀 텍스트를 결합하고,
        문서 내 모든 이미지 관계(rels)를 순회해 PNG로 저장한다.

        Args:
            file_path: DOCX 파일 경로
            doc_id: 데이터베이스 문서 ID

        Returns:
            (body_text, n_assets) 튜플.
        """
        # python-docx 라이브러리로 Word 파일 열기
        from docx import Document
        try:
            doc = Document(file_path)
        except Exception as e:
            logger.error("DOCX 열기 실패 (%s): %s", file_path, e)
            return "", 0

        # 본문 텍스트 조합 (문단 + 표)
        parts: list[str] = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                parts.append(t)
        for table in doc.tables:
            for row in table.rows:
                # 셀 내용을 "|"로 구분해 한 줄로 표현
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n\n".join(parts)

        # 이미지 추출 (문서 내 관계(rels) 목록에서 이미지 타입만 골라냄)
        n_assets = 0
        fig_seq = 0
        for rel in doc.part.rels.values():
            if "image" not in rel.reltype:  # 이미지 관계가 아니면 건너뜀
                continue
            try:
                image_blob = rel.target_part.blob  # 원본 이미지 바이트
                image_blob = self._to_png_safe(image_blob)  # PNG로 변환
                ref_tag = self._adb.save_asset(
                    document_id=doc_id,
                    file_path=file_path,
                    page_num=0,          # DOCX는 페이지 구분 없이 0번으로 통일
                    asset_type="FIGURE",
                    seq_in_doc=fig_seq,
                    bbox=None,           # DOCX는 정확한 좌표를 추출하기 어려움
                    image_data=image_blob,
                )
                fig_seq += 1
                n_assets += 1
                logger.debug("DOCX 이미지 저장: %s (doc_id=%d)", ref_tag, doc_id)
            except Exception as e:
                logger.warning("DOCX 이미지 추출 실패 (doc_id=%d): %s", doc_id, e)

        # 캡션 텍스트 기반 ref_tag 보완 (이미지 탐지에서 놓친 캡션 처리)
        counts = self._adb.count_assets(doc_id)
        page_texts = [text]
        page_texts, _, _ = self._inject_caption_refs(
            page_texts, file_path, doc_id,
            fig_seq=counts.get("FIGURE", 0),
            tbl_seq=counts.get("TABLE", 0),
        )
        body = page_texts[0]
        self._pdb.save_page_content(doc_id, 0, body_text=body)
        title = _extract_simple_title(body)
        self._pdb.save_document_metadata(doc_id, title=title)
        logger.info("DOCX 처리 완료: %d건 (doc_id=%d)", n_assets, doc_id)
        return body, n_assets

    def _process_pptx(self, file_path: str, doc_id: int) -> tuple[str, int]:
        """PowerPoint(.pptx): 슬라이드 텍스트 추출 + 임베디드 이미지 DB 저장.

        슬라이드를 페이지 단위로 처리하며, 각 슬라이드의 도형(shape)에서
        텍스트와 이미지를 추출한다.

        Args:
            file_path: PPTX 파일 경로
            doc_id: 데이터베이스 문서 ID

        Returns:
            (full_text, n_assets) 튜플.
        """
        # python-pptx 라이브러리로 PowerPoint 파일 열기
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        try:
            prs = Presentation(file_path)
        except Exception as e:
            logger.error("PPTX 열기 실패 (%s): %s", file_path, e)
            return "", 0

        page_texts: list[str] = []   # 슬라이드별 텍스트 저장
        n_assets = 0
        fig_seq = 0

        for slide_num, slide in enumerate(prs.slides):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                # 텍스트 프레임이 있는 도형: 텍스트 추출
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_texts.append(t)
                elif shape.has_table:
                    # 표 도형: 행별로 셀 텍스트를 "|"로 연결
                    for row in shape.table.rows:
                        cells = [c.text_frame.text.strip() for c in row.cells
                                 if c.text_frame.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))

                # 이미지 도형: PNG로 변환 후 DB 저장
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_blob = shape.image.blob
                        image_blob = self._to_png_safe(image_blob)
                        ref_tag = self._adb.save_asset(
                            document_id=doc_id,
                            file_path=file_path,
                            page_num=slide_num,   # 슬라이드 번호를 페이지 번호로 사용
                            asset_type="FIGURE",
                            seq_in_doc=fig_seq,
                            bbox=None,
                            image_data=image_blob,
                        )
                        fig_seq += 1
                        n_assets += 1
                        # 슬라이드 텍스트에 참조 태그도 추가 (위치 정보 포함)
                        slide_texts.append(ref_tag)
                        logger.debug("PPTX 이미지 저장: %s (슬라이드 %d)", ref_tag, slide_num)
                    except Exception as e:
                        logger.warning("PPTX 이미지 추출 실패 (슬라이드 %d): %s", slide_num, e)
            page_texts.append("\n\n".join(slide_texts))

        # 캡션 텍스트 기반 ref_tag 보완
        counts = self._adb.count_assets(doc_id)
        page_texts, _, _ = self._inject_caption_refs(
            page_texts, file_path, doc_id,
            fig_seq=counts.get("FIGURE", 0),
            tbl_seq=counts.get("TABLE", 0),
        )
        # 슬라이드별로 DB에 저장
        for slide_num, slide_body in enumerate(page_texts):
            self._pdb.save_page_content(doc_id, slide_num, body_text=slide_body)
        title = _extract_simple_title(page_texts[0] if page_texts else "")
        self._pdb.save_document_metadata(doc_id, title=title)
        full_text = "\n\n".join(page_texts)
        logger.info("PPTX 처리 완료: %d건 (doc_id=%d)", n_assets, doc_id)
        return full_text, n_assets

    @staticmethod
    def _to_png_safe(image_blob: bytes) -> bytes:
        """이미지 바이트를 PIL로 열어 PNG로 변환한다. 변환 실패 시 원본 반환.

        DOCX/PPTX 내부 이미지는 다양한 포맷(JPEG, EMF, WMF 등)일 수 있으므로
        일관된 PNG 형식으로 통일한다.

        Args:
            image_blob: 원본 이미지 바이트

        Returns:
            PNG 포맷 바이트. 변환 실패 시 원본 image_blob 반환.
        """
        try:
            buf = io.BytesIO()
            # RGBA로 변환하면 투명도(알파 채널)가 있는 이미지도 처리 가능
            PILImage.open(io.BytesIO(image_blob)).convert("RGBA").save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            # 알 수 없는 포맷이면 원본 그대로 반환 (예: EMF, WMF 등)
            return image_blob

    # ──────────────────────────────────────────────────────────────────────
    # 디지털 PDF 처리
    # ──────────────────────────────────────────────────────────────────────

    def _extract_digital_with_assets(
        self,
        pdf_path: str,
        doc_id: int,
    ) -> tuple[list[str], list[str], list[str], int]:
        """디지털 PDF에서 페이지별 자산 추출 + 참조 태그 포함 텍스트 구성.

        각 페이지를 순회하며:
        1. 자산(표·그림) 탐지 및 PNG 렌더링
        2. DB 저장 및 참조 태그 할당
        3. 단어 좌표 기반으로 본문 텍스트 구성 (단단/다단 분기)

        Args:
            pdf_path: PDF 파일 경로
            doc_id: 데이터베이스 문서 ID

        Returns:
            (page_bodies, page_headers, page_footers, n_assets) 튜플.
            각 리스트는 페이지 순서대로 정렬됨.
        """
        from src.ocr_engine import OcrEngine  # 순환 import 방지

        page_bodies:  list[str] = []
        page_headers: list[str] = []
        page_footers: list[str] = []
        n_assets = 0
        table_seq = 0    # TABLE 자산의 문서 내 순번 (001, 002, ...)
        figure_seq = 0   # FIGURE 자산의 문서 내 순번

        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # 워터마크 제거 후 페이지 객체 반환
                filtered = OcrEngine.filter_watermarks(page)
                page_height = float(page.height)
                page_width  = float(page.width)

                # ① 자산 탐지 + PNG 렌더링 (표·그림 영역을 이미지로 추출)
                detected = extract_assets(pdf_path, page_num)

                # 전체 단어 추출 후 헤더/본문/푸터 분리
                # pdfplumber는 각 단어의 x, y 좌표를 딕셔너리로 반환
                all_words = filtered.extract_words()
                header_words, main_words, footer_words = OcrEngine._split_header_footer_words(
                    all_words, page_height
                )
                header_text = OcrEngine._words_to_text(header_words) if header_words else ""
                footer_text = OcrEngine._words_to_text(footer_words) if footer_words else ""

                # 자산이 없는 페이지: 텍스트만 구성
                if not detected:
                    if main_words:
                        # 다단 여부 판단: 컬럼 사이 공백(gap_x)이 있으면 다단
                        gap_x = OcrEngine._find_column_gap(main_words, page_width)
                        if gap_x is None:
                            # 단단(단일 컬럼): 일반 텍스트 구성
                            body_text = OcrEngine._words_to_text_with_paragraphs(main_words)
                        else:
                            # 다단: 캡션 처리를 포함한 다단 텍스트 구성
                            body_start_y = OcrEngine._find_body_start_y(main_words, gap_x)
                            body_text = self._compose_multicol_text_with_captions(main_words, gap_x, body_start_y)
                    else:
                        # 단어 추출 실패 시 전체 텍스트로 폴백
                        body_text = filtered.extract_text() or ""
                    page_bodies.append(body_text)
                    page_headers.append(header_text)
                    page_footers.append(footer_text)
                    continue

                # ② DB 저장 + ref_tag 할당
                asset_refs: list[dict] = []
                for asset in detected:
                    # 자산 이미지에서 텍스트 OCR 추출 (표는 psm 6, 그림은 psm 3)
                    ocr_text = self._ocr_image(asset.image_data, asset.asset_type)
                    if asset.asset_type == "TABLE":
                        seq = table_seq
                        table_seq += 1
                    else:
                        seq = figure_seq
                        figure_seq += 1

                    # 자산을 DB에 저장하고 참조 태그(예: [TABLE_001]) 반환
                    ref_tag = self._adb.save_asset(
                        document_id=doc_id,
                        file_path=pdf_path,
                        page_num=page_num,
                        asset_type=asset.asset_type,
                        seq_in_doc=seq,
                        bbox=asset.bbox,           # (x0, top, x1, bottom) 좌표
                        image_data=asset.image_data,
                        ocr_text=ocr_text,
                        caption=asset.caption,
                    )
                    # 참조 태그와 위치 정보를 묶어 보관 (본문 삽입 시 사용)
                    asset_refs.append({
                        "ref_tag":    ref_tag,
                        "bbox":       asset.bbox,
                        "asset_type": asset.asset_type,
                    })
                    n_assets += 1
                    logger.debug("자산 저장: %s (p%d)", ref_tag, page_num)

                # ③ 본문 단어 + 참조 태그 삽입 (다단 여부 판단)
                gap_x = OcrEngine._find_column_gap(main_words, page_width) if main_words else None
                if gap_x is None:
                    # 단단 페이지: 자산 위치 기준으로 텍스트 분할 후 태그 삽입
                    body_text = self._build_digital_page_text(main_words, asset_refs, page_height)
                else:
                    # 다단 페이지: 컬럼별로 텍스트 구성 후 합침
                    body_start_y = OcrEngine._find_body_start_y(main_words, gap_x)
                    body_text = self._build_digital_page_text_multicol(
                        main_words, asset_refs, page_height, gap_x, body_start_y
                    )
                page_bodies.append(body_text)
                page_headers.append(header_text)
                page_footers.append(footer_text)

        return page_bodies, page_headers, page_footers, n_assets

    def _build_digital_page_text(
        self,
        words: list,
        asset_refs: list[dict],
        page_height: float,
    ) -> str:
        """단어 목록과 자산 목록으로 참조 태그가 포함된 페이지 텍스트를 구성한다.

        자산의 bbox(경계 상자) 기준으로 텍스트를 위/아래 구역으로 나누고,
        자산 위치에 참조 태그를 삽입한다.

        Args:
            words: pdfplumber에서 추출한 단어 딕셔너리 목록
            asset_refs: 자산 참조 정보 목록 (ref_tag, bbox, asset_type 포함)
            page_height: 페이지 전체 높이 (픽셀)

        Returns:
            참조 태그가 삽입된 페이지 텍스트 문자열.
        """
        from src.ocr_engine import OcrEngine

        # 자산이 없으면 단어만으로 텍스트 구성
        if not asset_refs:
            return OcrEngine._words_to_text_with_paragraphs(words) if words else ""

        # 단어가 없으면 자산 태그만 세로로 나열 (y 좌표 순)
        if not words:
            return "\n".join(
                a["ref_tag"] for a in sorted(asset_refs, key=lambda a: a["bbox"][1])
            )

        # y 좌표(bbox[1] = top) 기준으로 자산을 위에서 아래 순으로 정렬
        sorted_assets = sorted(asset_refs, key=lambda a: a["bbox"][1])

        # 자산 bbox 내부에 있는 단어 제거 (표 셀 텍스트 등이 본문에 중복 포함되지 않도록)
        clean_words = [
            w for w in words
            if not any(self._word_in_bbox(w, a["bbox"]) for a in sorted_assets)
        ]

        # 구역 분리: 자산 위치를 기준으로 텍스트를 위/아래로 나눔
        # zones: [("text", "텍스트 내용"), ("asset", "[TABLE_001]"), ...]
        zones: list[tuple[str, str]] = []
        current_y = 0.0   # 현재까지 처리한 y 좌표 (아래로 내려가면서 증가)

        for asset in sorted_assets:
            top    = asset["bbox"][1]    # 자산 상단 y 좌표
            bottom = asset["bbox"][3]    # 자산 하단 y 좌표

            # 이전 자산 하단부터 현재 자산 상단까지의 단어 (자산 위쪽 텍스트)
            above = [w for w in clean_words
                     if current_y - 2 <= float(w.get("top", 0)) < top - 2]
            if above:
                zones.append(("text", OcrEngine._words_to_text_with_paragraphs(above)))

            # 자산 위치에 참조 태그 삽입
            zones.append(("asset", asset["ref_tag"]))
            current_y = bottom + 2  # 다음 탐색 시작 y 좌표 (자산 하단 + 여유값)

        # 마지막 자산 아래쪽 나머지 텍스트
        remaining = [w for w in clean_words if float(w.get("top", 0)) >= current_y - 2]
        if remaining:
            zones.append(("text", OcrEngine._words_to_text_with_paragraphs(remaining)))

        # 구역들을 합치면서 자산으로 분리된 문장 재연결
        return self._join_zones(zones)

    def _build_digital_page_text_multicol(
        self,
        words: list,
        asset_refs: list[dict],
        page_height: float,
        gap_x: float,
        body_start_y: float | None,
    ) -> str:
        """다단 페이지: 컬럼별로 자산 참조 태그를 삽입한 뒤 좌→우 순서로 연결한다.

        자산을 좌측 컬럼, 우측 컬럼, 또는 양 컬럼에 걸친(스패닝) 세 가지로 분류한다.

        Args:
            words: 페이지 전체 단어 목록
            asset_refs: 자산 참조 정보 목록
            page_height: 페이지 높이
            gap_x: 좌우 컬럼 사이 공백의 x 좌표
            body_start_y: 본문 시작 y 좌표 (None이면 전체가 본문)

        Returns:
            참조 태그가 삽입된 다단 텍스트 문자열.
        """
        from src.ocr_engine import OcrEngine

        # y 좌표 기준 정렬
        sorted_assets = sorted(asset_refs, key=lambda a: a["bbox"][1])

        # 자산 bbox 내부 단어 제거 (중복 방지)
        clean_words = [
            w for w in words
            if not any(self._word_in_bbox(w, a["bbox"]) for a in sorted_assets)
        ]

        # 헤더(전체폭)와 본문 분리
        if body_start_y is not None:
            header_words = [w for w in clean_words if float(w.get("top", 0)) < body_start_y]
            body_words   = [w for w in clean_words if float(w.get("top", 0)) >= body_start_y]
        else:
            header_words = []
            body_words   = clean_words

        header_text = OcrEngine._words_to_text_with_paragraphs(header_words) if header_words else ""

        # 본문 단어를 컬럼별로 분리 (gap_x 기준으로 좌/우 분류)
        left_words, right_words = OcrEngine._split_words_by_column(body_words, gap_x)

        # 자산을 컬럼별 분류 (좌·우·스패닝)
        left_assets:  list[dict] = []   # 좌측 컬럼에만 속하는 자산
        right_assets: list[dict] = []   # 우측 컬럼에만 속하는 자산
        span_assets:  list[dict] = []   # 양쪽 컬럼에 걸친 자산
        for asset in sorted_assets:
            ax0, _, ax1, _ = asset["bbox"]
            if ax1 <= gap_x + 5:          # 자산 오른쪽 끝이 gap보다 왼쪽 → 좌측
                left_assets.append(asset)
            elif ax0 >= gap_x - 5:        # 자산 왼쪽 끝이 gap보다 오른쪽 → 우측
                right_assets.append(asset)
            else:                          # 그 사이 → 스패닝
                span_assets.append(asset)

        # 각 컬럼별로 텍스트 구성 (단단 처리 함수 재사용)
        left_text  = self._build_digital_page_text(left_words,  left_assets,  page_height)
        right_text = self._build_digital_page_text(right_words, right_assets, page_height)

        # 스패닝 자산: 좌측과 우측 텍스트 사이에 삽입
        if span_assets:
            span_part = "\n".join(a["ref_tag"] for a in span_assets)
            body_text = "\n\n".join(s for s in [left_text, span_part, right_text] if s)
        else:
            # 캡션/런닝헤더를 고려하며 컬럼 연결
            body_text = self._join_columns_strip_captions(left_text, right_text)

        # 헤더 텍스트와 본문을 합쳐 최종 결과 반환
        return "\n\n".join(s for s in [header_text, body_text] if s)

    def _join_columns_strip_captions(self, left_text: str, right_text: str) -> str:
        """캡션·런닝헤더를 임시 제거한 뒤 컬럼을 연결하고 제거된 항목을 재삽입한다.

        캡션이 좌 컬럼 끝 또는 우 컬럼 머리에 있으면 _join_columns가 연결을 막음.
        - 좌 컬럼 꼬리: caption 단락 제거
        - 우 컬럼 머리: caption 및 런닝헤더("…" 종결) 단락 제거
        제거 후 연결 → 연결된 첫 단락 뒤에 제거 항목 복원.

        Args:
            left_text: 좌측 컬럼 텍스트
            right_text: 우측 컬럼 텍스트

        Returns:
            두 컬럼이 자연스럽게 연결된 텍스트.
        """
        from src.ocr_engine import OcrEngine

        # 텍스트를 단락 단위로 분리 (빈 단락은 제외)
        left_paras  = [p for p in left_text.strip().split('\n\n')  if p.strip()]
        right_paras = [p for p in right_text.strip().split('\n\n') if p.strip()]

        # 한쪽이 비어있으면 단순 연결
        if not left_paras or not right_paras:
            return OcrEngine._join_columns(left_text, right_text)

        # 좌 컬럼 꼬리의 캡션 단락을 제거 (나중에 복원)
        left_tail: list[str] = []
        while left_paras:
            s = left_paras[-1].strip()
            if _TABLE_CAPTION_RE.match(s) or _FIGURE_CAPTION_RE.match(s):
                left_tail.insert(0, left_paras.pop())  # 맨 앞에 넣어 순서 유지
            else:
                break

        # 우 컬럼 머리의 캡션 + 런닝헤더 + ref_tag 제거 (나중에 복원)
        right_head: list[str] = []
        while right_paras:
            s = right_paras[0].strip()
            if (_TABLE_CAPTION_RE.match(s) or _FIGURE_CAPTION_RE.match(s)
                    or s.endswith('…')          # "…" 로 끝나는 런닝헤더
                    or _REF_TAG_RE.match(s)):   # 비트맵 탐지 ref_tag
                right_head.append(right_paras.pop(0))
            else:
                break

        # 제거 대상이 없으면 단순 연결
        if not left_tail and not right_head:
            return OcrEngine._join_columns(left_text, right_text)

        n_left = len(left_paras)   # 좌 컬럼 단락 수 (join 후에도 기준점으로 사용)
        clean_left  = '\n\n'.join(left_paras)
        clean_right = '\n\n'.join(right_paras)
        # 캡션을 제거한 상태에서 컬럼 연결 (문장 재연결 포함)
        joined = OcrEngine._join_columns(clean_left, clean_right)
        joined_paras = [p for p in joined.split('\n\n') if p.strip()]
        if not joined_paras:
            return '\n\n'.join(left_tail + right_head)

        # joined_paras[:n_left] = 좌 컬럼 단락들(마지막은 우 컬럼 첫 단락과 합쳐질 수 있음)
        # joined_paras[n_left:] = 우 컬럼 나머지 단락들
        # 좌 컬럼 단락들 + 제거했던 캡션들 + 우 컬럼 나머지 단락들 순으로 재조합
        result_parts = joined_paras[:n_left] + left_tail + right_head + joined_paras[n_left:]
        return '\n\n'.join(result_parts)

    def _compose_multicol_text_with_captions(
        self,
        words: list,
        gap_x: float,
        body_start_y: float | None,
    ) -> str:
        """다단 텍스트를 구성한다. OcrEngine._compose_multicol_text와 동일하나 캡션 처리를 포함한다.

        OcrEngine의 기본 다단 구성 함수 대신 이 함수를 사용하면
        컬럼 연결 시 캡션·런닝헤더를 올바르게 처리한다.

        Args:
            words: 페이지 단어 목록
            gap_x: 컬럼 구분 x 좌표
            body_start_y: 본문 시작 y 좌표

        Returns:
            다단 텍스트 문자열.
        """
        from src.ocr_engine import OcrEngine

        # body_start_y가 없으면 단순히 좌/우 분리 후 연결
        if body_start_y is None:
            left, right = OcrEngine._split_words_by_column(words, gap_x)
            return self._join_columns_strip_captions(
                OcrEngine._words_to_text_with_paragraphs(left),
                OcrEngine._words_to_text_with_paragraphs(right),
            )

        # 본문 시작 전 단어들: 전체폭 줄은 헤더로, 나머지는 컬럼에 배분
        header: list = []
        left_body: list = []
        right_body: list = []

        pre_body_words = [w for w in words if float(w["top"]) < body_start_y]
        for _, line_words in OcrEngine._group_words_by_line(pre_body_words):
            if OcrEngine._line_spans_both_columns(line_words, gap_x):
                # 양쪽 컬럼에 걸쳐 있는 줄 → 헤더로 분류
                header.extend(line_words)
                continue
            left_line, right_line = OcrEngine._split_words_by_column(line_words, gap_x)
            left_body.extend(left_line)
            right_body.extend(right_line)

        # 본문 시작 후 단어들: 컬럼별로 분리
        body_words = [w for w in words if float(w["top"]) >= body_start_y]
        left_words, right_words = OcrEngine._split_words_by_column(body_words, gap_x)
        left_body.extend(left_words)
        right_body.extend(right_words)

        # 헤더 텍스트 구성
        header_text = OcrEngine._words_to_text_with_paragraphs(header)
        # 캡션 처리를 포함한 컬럼 연결
        body_joined = self._join_columns_strip_captions(
            OcrEngine._words_to_text_with_paragraphs(left_body),
            OcrEngine._words_to_text_with_paragraphs(right_body),
        )
        # 헤더와 본문을 합쳐 반환
        return "\n\n".join(s for s in (header_text, body_joined) if s)

    # ──────────────────────────────────────────────────────────────────────
    # 스캔 PDF 처리
    # ──────────────────────────────────────────────────────────────────────

    def _extract_all_assets(
        self,
        pdf_path: str,
        doc_id: int,
    ) -> tuple[list[dict], int]:
        """스캔 PDF의 모든 페이지에서 자산을 추출하여 DB에 저장한다.

        각 페이지의 탐지된 자산을 순차적으로 처리하고,
        나중에 텍스트에 삽입할 수 있도록 페이지 번호·bbox 정보와 함께 반환한다.

        Args:
            pdf_path: PDF 파일 경로
            doc_id: 데이터베이스 문서 ID

        Returns:
            (asset_list, n_assets) 튜플.
            asset_list: 자산별 정보 딕셔너리 목록 (ref_tag, page_num, bbox 등 포함)
            n_assets: 총 저장된 자산 수
        """
        all_assets: list[dict] = []
        n_assets = 0
        table_seq = 0
        figure_seq = 0

        # 먼저 전체 페이지 수와 각 페이지 높이를 읽어둠
        try:
            with pdfplumber.open(pdf_path) as pdf:
                n_pages = len(pdf.pages)
                page_heights = [float(p.height) for p in pdf.pages]
        except Exception as e:
            logger.error("페이지 정보 조회 실패 (%s): %s", pdf_path, e)
            return [], 0

        for page_num in range(n_pages):
            detected = extract_assets(pdf_path, page_num)
            for asset in detected:
                # 자산 이미지에서 텍스트 OCR 추출
                ocr_text = self._ocr_image(asset.image_data, asset.asset_type)
                if asset.asset_type == "TABLE":
                    seq = table_seq
                    table_seq += 1
                else:
                    seq = figure_seq
                    figure_seq += 1

                # DB에 저장하고 참조 태그 반환
                ref_tag = self._adb.save_asset(
                    document_id=doc_id,
                    file_path=pdf_path,
                    page_num=page_num,
                    asset_type=asset.asset_type,
                    seq_in_doc=seq,
                    bbox=asset.bbox,
                    image_data=asset.image_data,
                    ocr_text=ocr_text,
                    caption=asset.caption,
                )
                # 텍스트 삽입 시 필요한 위치 정보 포함하여 보관
                all_assets.append({
                    "ref_tag":     ref_tag,
                    "page_num":    page_num,
                    "bbox":        asset.bbox,
                    "asset_type":  asset.asset_type,
                    "page_height": page_heights[page_num] if page_num < len(page_heights) else None,
                })
                n_assets += 1
                logger.debug("자산 저장: %s (p%d)", ref_tag, page_num)

        return all_assets, n_assets

    def _inject_refs_scanned(
        self,
        page_texts: list[str],
        assets: list[dict],
    ) -> list[str]:
        """스캔 PDF OCR 텍스트에 페이지별로 참조 태그를 삽입한다.

        자산을 페이지별로 그룹화한 뒤, 각 페이지 텍스트에
        자산 위치 기반으로 참조 태그를 삽입한다.

        Args:
            page_texts: OCR로 생성된 페이지별 텍스트 목록
            assets: _extract_all_assets에서 반환된 자산 정보 목록

        Returns:
            참조 태그가 삽입된 페이지 텍스트 목록.
        """
        # defaultdict: 키가 없어도 자동으로 빈 리스트 생성 (일반 dict는 KeyError 발생)
        page_asset_map: dict[int, list[dict]] = defaultdict(list)
        # 페이지 번호 → y 좌표 순으로 정렬하여 맵에 추가
        for a in sorted(assets, key=lambda x: (x["page_num"], x["bbox"][1])):
            page_asset_map[a["page_num"]].append(a)

        result = list(page_texts)  # 원본을 변경하지 않도록 복사
        for page_num, page_assets in page_asset_map.items():
            if page_num >= len(result):
                continue  # 페이지 번호가 텍스트 목록 범위를 벗어나면 건너뜀
            text = result[page_num]
            if not text.strip():
                # 텍스트가 없으면 자산 태그만 나열
                result[page_num] = "\n".join(a["ref_tag"] for a in page_assets)
                continue
            # 텍스트가 있으면 자산 위치 기반으로 태그 삽입
            result[page_num] = self._inject_into_page_text(text, page_assets)

        return result

    def _inject_into_page_text(
        self,
        text: str,
        sorted_assets: list[dict],
    ) -> str:
        """단일 페이지 텍스트에 자산 참조 태그를 삽입한다.

        캡션 라인을 먼저 탐색하고, 없으면 y 좌표 비율로 위치를 추정하여 삽입한다.

        Args:
            text: 페이지 OCR 텍스트
            sorted_assets: y 좌표 순으로 정렬된 자산 정보 목록

        Returns:
            참조 태그가 삽입된 텍스트 문자열.
        """
        lines = text.split("\n")

        for asset in sorted_assets:
            ref_tag      = asset["ref_tag"]
            asset_type   = asset["asset_type"]
            page_height  = asset.get("page_height")

            # ① 캡션 라인 탐색 (예: "Table 1.", "그림 2." 등의 패턴 탐색)
            caption_idx = self._find_caption_line(lines, asset_type)
            if caption_idx is not None:
                # 캡션 근처에 태그 삽입
                lines = self._insert_after_line(lines, caption_idx, ref_tag)
                continue

            # ② 위치 추정 (bbox y 좌표 기반)
            # 자산이 페이지 내 어느 위치에 있는지 비율로 계산하여 줄 번호 추정
            if page_height and page_height > 0:
                rel_y  = asset["bbox"][1] / page_height  # 0.0 ~ 1.0 사이의 상대적 위치
                target = max(0, min(len(lines) - 1, int(rel_y * len(lines))))
            else:
                target = len(lines) // 2  # 페이지 높이를 모르면 중간에 삽입

            # 문장 경계를 찾아 태그 삽입
            lines = self._insert_at_boundary(lines, target, ref_tag)

        return "\n".join(lines)

    def _find_caption_line(self, lines: list[str], asset_type: str) -> int | None:
        """캡션 패턴과 일치하는 라인 인덱스를 반환한다 (처음 발견된 것).

        Args:
            lines: 페이지 텍스트 줄 목록
            asset_type: "TABLE" 또는 "FIGURE"

        Returns:
            캡션 줄의 인덱스 (0부터 시작). 없으면 None.
        """
        # 자산 타입에 따라 적절한 정규식 패턴 선택
        pattern = _TABLE_CAPTION_RE if asset_type == "TABLE" else _FIGURE_CAPTION_RE
        for i, line in enumerate(lines):
            if pattern.match(line.strip()):
                return i
        return None

    def _insert_after_line(
        self, lines: list[str], idx: int, ref_tag: str
    ) -> list[str]:
        """캡션 라인(idx) 근처에 ref_tag를 삽입한다.

        우선순위:
        1. index만 있는 라인("Fig 1", "Table 2" 등) → 같은 줄 끝에 인라인 삽입
        2. 다음 줄이 소문자로 시작(분리 문장) → 같은 줄 끝에 인라인 삽입
        3. 그 외 → 다음 줄로 삽입

        Args:
            lines: 페이지 텍스트 줄 목록
            idx: 캡션 라인의 인덱스
            ref_tag: 삽입할 참조 태그 문자열 (예: "[FIGURE_001]")

        Returns:
            참조 태그가 삽입된 줄 목록.
        """
        new_lines = list(lines)
        before    = new_lines[idx].rstrip()

        # ① index만 포함된 라인이면 같은 줄 끝에 태그 추가 (예: "Fig 1" → "Fig 1 [FIGURE_001]")
        if _INDEX_ONLY_RE.match(before.strip()):
            new_lines[idx] = before + " " + ref_tag
            return new_lines

        next_idx = idx + 1
        if next_idx < len(new_lines):
            after = new_lines[next_idx].lstrip()
            if after and after[0].islower():
                # ② 분리 문장 재연결: 캡션 끝 + ref_tag (다음 줄은 그대로 유지)
                # 다음 줄이 소문자 시작 = 앞 문장에서 잘린 것으로 간주
                new_lines[idx] = before + " " + ref_tag
                return new_lines

        # ③ 다음 줄로 삽입 (일반적인 경우)
        new_lines.insert(idx + 1, ref_tag)
        return new_lines

    def _insert_at_boundary(
        self,
        lines: list[str],
        near: int,
        ref_tag: str,
    ) -> list[str]:
        """near 근처에서 문단 경계를 찾아 ref_tag를 삽입한다. 분리 문장이면 재연결.

        캡션을 찾지 못했을 때 y 좌표 추정으로 삽입 위치를 결정하는 함수.
        빈 줄 또는 문장 종결 기호를 우선 탐색하여 자연스러운 위치에 삽입한다.

        Args:
            lines: 페이지 텍스트 줄 목록
            near: 삽입 목표 줄 번호 (y 좌표 추정값)
            ref_tag: 삽입할 참조 태그

        Returns:
            참조 태그가 삽입된 줄 목록.
        """
        n = len(lines)
        if n == 0:
            return [ref_tag]

        # near ±6 범위에서 문장 종결 또는 빈 라인 탐색 (가장 가까운 경계 우선)
        best_idx = near
        found = False
        for delta in range(min(n, 7)):
            for sign in ([0] if delta == 0 else [1, -1]):
                idx = near + sign * delta
                if not (0 <= idx < n):
                    continue
                line = lines[idx].rstrip()
                if not line:                       # 빈 라인 → 단락 경계
                    best_idx = idx + 1
                    found = True
                    break
                if line[-1:] in _SENT_END:         # 마침표 등 문장 종결 기호
                    best_idx = idx + 1
                    found = True
                    break
            if found:
                break

        # 범위 초과 방지
        best_idx = max(0, min(n, best_idx))

        # 삽입 위치 전후 줄을 확인하여 분리 문장인지 검사
        before_line = lines[best_idx - 1].rstrip() if best_idx > 0 else ""
        after_line  = lines[best_idx].lstrip()     if best_idx < n else ""

        # 앞줄이 문장 종결 기호 없이 끝나고 뒷줄이 소문자로 시작 → 분리 문장
        if (before_line
                and before_line[-1:] not in _SENT_END
                and after_line
                and after_line[0:1].islower()):
            # 분리 문장 사이에 ref_tag를 끼워 넣어 문장 흐름 유지
            new_lines = list(lines)
            new_lines[best_idx - 1] = before_line + " " + ref_tag + " " + after_line
            if best_idx < len(new_lines):
                del new_lines[best_idx]  # 뒷줄은 앞줄에 합쳐졌으므로 삭제
            return new_lines

        # 일반 경계에 새 줄로 삽입
        new_lines = list(lines)
        new_lines.insert(best_idx, ref_tag)
        return new_lines

    # ──────────────────────────────────────────────────────────────────────
    # 캡션 텍스트 기반 보완
    # ──────────────────────────────────────────────────────────────────────

    def _inject_caption_refs(
        self,
        page_texts: list[str],
        pdf_path: str,
        doc_id: int,
        fig_seq: int,
        tbl_seq: int,
    ) -> tuple[list[str], int, int]:
        """텍스트 캡션 스캔으로 참조 태그가 누락된 Figure/Table을 보완한다.

        이미지 탐지기가 놓친 자산(주로 벡터 그래픽, 텍스트 기반 표)을
        캡션 패턴 문자열로 찾아내어 더미 이미지와 함께 DB에 저장한다.
        이미 ref_tag가 있는 캡션 줄은 건너뛴다.

        Args:
            page_texts: 페이지별 텍스트 목록
            pdf_path: 원본 PDF 경로 (DB 저장 시 file_path로 사용)
            doc_id: 데이터베이스 문서 ID
            fig_seq: 다음에 사용할 FIGURE 순번 (기존 자산 수에서 시작)
            tbl_seq: 다음에 사용할 TABLE 순번

        Returns:
            (modified_page_texts, next_fig_seq, next_tbl_seq) 튜플.
        """
        result: list[str] = []
        for page_num, text in enumerate(page_texts):
            lines = text.split("\n")
            i = 0
            while i < len(lines):
                stripped = lines[i].strip()

                # 현재 줄에 이미 ref_tag 존재 → 이중 삽입 방지
                if _REF_TAG_RE.search(lines[i]):
                    i += 1
                    continue
                # 다음 줄에 ref_tag 존재 → 캡션 바로 뒤에 이미 태그가 있는 경우
                if i + 1 < len(lines) and _REF_TAG_RE.search(lines[i + 1]):
                    i += 1
                    continue
                # 바로 이전 비어있지 않은 줄이 ref_tag → 비트맵 탐지와 중복 방지
                prev_ne = i - 1
                while prev_ne >= 0 and not lines[prev_ne].strip():
                    prev_ne -= 1  # 빈 줄 건너뛰기
                if prev_ne >= 0 and _REF_TAG_RE.search(lines[prev_ne]):
                    i += 1
                    continue

                # 캡션 패턴 매칭: 표 또는 그림 캡션인지 판별
                if _TABLE_CAPTION_RE.match(stripped):
                    asset_type, seq = "TABLE", tbl_seq
                    tbl_seq += 1
                elif _FIGURE_CAPTION_RE.match(stripped):
                    asset_type, seq = "FIGURE", fig_seq
                    fig_seq += 1
                else:
                    i += 1
                    continue  # 캡션이 아니면 건너뜀

                try:
                    # 이미지 없이 캡션만으로 DB에 저장 (더미 PNG 사용)
                    ref_tag = self._adb.save_asset(
                        document_id=doc_id,
                        file_path=pdf_path,
                        page_num=page_num,
                        asset_type=asset_type,
                        seq_in_doc=seq,
                        bbox=None,
                        image_data=_PLACEHOLDER_PNG,  # 1×1 흰색 더미 이미지
                        ocr_text=None,
                        caption=stripped,             # 캡션 원문 보존
                    )
                    # 캡션 라인을 ref_tag으로 교체 (캡션 텍스트는 DB caption 컬럼에 보존)
                    lines[i] = ref_tag
                    logger.debug("캡션 ref_tag 교체: %s p%d '%s'",
                                 ref_tag, page_num, stripped[:50])

                    # 캡션이 문장 사이에 끼어 있으면 앞뒤 문장 재연결
                    # 조건: 앞줄이 문장 종결 기호 없음 + 뒷줄이 소문자 시작
                    prev_idx = i - 1
                    while prev_idx >= 0 and not lines[prev_idx].strip():
                        prev_idx -= 1  # 빈 줄 건너뛰기
                    next_idx = i + 1
                    while next_idx < len(lines) and not lines[next_idx].strip():
                        next_idx += 1  # 빈 줄 건너뛰기
                    if 0 <= prev_idx and next_idx < len(lines):
                        ptext = lines[prev_idx].rstrip()
                        ntext = lines[next_idx].lstrip()
                        if (ptext and ptext[-1] not in _SENT_END
                                and ntext and ntext[0].islower()):
                            # 앞문장 끝 + 뒷문장 앞 + ref_tag를 하나의 줄로 병합
                            lines[prev_idx] = ptext + " " + ntext + " " + ref_tag
                            # 병합된 줄들(next_idx ~ i)을 역순으로 삭제
                            for k in range(next_idx, i - 1, -1):
                                del lines[k]
                            i = prev_idx  # 합쳐진 줄 위치로 인덱스 이동
                except Exception as exc:
                    logger.warning("캡션 ref_tag 보완 실패 (p%d '%s'): %s",
                                   page_num, stripped[:40], exc)
                i += 1

            result.append("\n".join(lines))

        return result, fig_seq, tbl_seq

    # ──────────────────────────────────────────────────────────────────────
    # 공통 유틸리티
    # ──────────────────────────────────────────────────────────────────────

    def _join_zones(self, zones: list[tuple[str, str]]) -> str:
        """텍스트·자산 구역들을 병합한다. 자산으로 분리된 문장은 재연결한다.

        zones 리스트의 각 항목은 ("text", 내용) 또는 ("asset", ref_tag) 형태다.
        자산 구역이 문장 중간에 있으면 앞뒤 텍스트를 자연스럽게 이어붙인다.

        Args:
            zones: (타입, 내용) 튜플의 목록

        Returns:
            병합된 텍스트 문자열.
        """
        if not zones:
            return ""

        # 자산 구역별로 인라인 삽입 여부를 미리 계산
        # 인라인 조건: 앞 텍스트가 종결 기호 없이 끝나고 뒤 텍스트가 소문자로 시작
        inline_flags = [False] * len(zones)
        for i, (zone_type, _) in enumerate(zones):
            if zone_type != "asset":
                continue

            # 이 자산 앞에 있는 가장 가까운 텍스트 구역 찾기
            prev_text = ""
            for j in range(i - 1, -1, -1):
                if zones[j][0] == "text" and zones[j][1].strip():
                    prev_text = zones[j][1]
                    break

            # 이 자산 뒤에 있는 가장 가까운 텍스트 구역 찾기
            next_text = ""
            for j in range(i + 1, len(zones)):
                if zones[j][0] == "text" and zones[j][1].strip():
                    next_text = zones[j][1]
                    break

            # 앞 텍스트의 마지막 줄과 뒤 텍스트의 첫 줄을 가져옴
            prev_last  = prev_text.rstrip().rsplit("\n", 1)[-1]  if prev_text.strip() else ""
            next_first = next_text.lstrip().split("\n",  1)[0]   if next_text.strip() else ""

            # 인라인 조건: 앞줄이 종결 기호 없이 끝나고 뒷줄이 소문자로 시작
            inline_flags[i] = (
                prev_last
                and prev_last[-1:]  not in _SENT_END
                and next_first
                and next_first[0:1].islower()
            )

        # 구역들을 순서대로 처리하여 최종 텍스트 구성
        parts: list[str] = []

        for i, (zone_type, content) in enumerate(zones):
            if zone_type == "text":
                # 직전 자산이 인라인이었으면 현재 텍스트를 앞 부분에 이어붙임
                prev_was_inline = (
                    i > 0
                    and zones[i - 1][0] == "asset"
                    and inline_flags[i - 1]
                )
                if prev_was_inline and parts:
                    # 공백 하나로 연결 (lstrip으로 앞 공백 제거)
                    parts[-1] = parts[-1] + " " + content.lstrip()
                else:
                    parts.append(content)
            else:
                # 자산 구역: 인라인이면 앞 텍스트 끝에 붙이고, 아니면 새 항목 추가
                if inline_flags[i] and parts:
                    parts[-1] = parts[-1].rstrip() + " " + content
                else:
                    parts.append(content)

        # 각 부분을 "\n\n"으로 구분하여 연결 (단락 구분)
        return "\n\n".join(p.strip() for p in parts if p.strip())

    @staticmethod
    def _word_in_bbox(word: dict, bbox: tuple, margin: float = 2.0) -> bool:
        """단어의 중심점이 bbox(경계 상자) 내부에 있는지 확인한다.

        margin을 두어 경계에 걸친 단어도 포함한다.

        Args:
            word: pdfplumber 단어 딕셔너리 (x0, x1, top, bottom 키 포함)
            bbox: (x0, top, x1, bottom) 형식의 자산 경계 상자
            margin: bbox 경계에서 허용하는 여유 픽셀 (기본 2.0)

        Returns:
            단어 중심점이 bbox 내부에 있으면 True.
        """
        x0, top, x1, bottom = bbox
        wx0    = float(word.get("x0", 0))
        wx1    = float(word.get("x1", 0))
        w_top  = float(word.get("top", 0))
        w_bot  = float(word.get("bottom", w_top))
        # 단어의 가로·세로 중심점 계산
        wcx    = (wx0 + wx1) / 2
        wcy    = (w_top + w_bot) / 2
        # 중심점이 margin을 포함한 bbox 범위 안에 있는지 확인
        return (x0 - margin <= wcx <= x1 + margin and
                top - margin <= wcy <= bottom + margin)

    def _ocr_image(self, image_data: bytes, asset_type: str) -> str:
        """Tesseract OCR로 자산 이미지에서 텍스트를 추출한다.

        자산 타입에 따라 OCR 페이지 분할 모드(psm)를 다르게 설정한다.
        - TABLE: psm 6 (균일한 텍스트 블록으로 간주)
        - FIGURE: psm 3 (자동 분할)

        Args:
            image_data: OCR할 이미지 바이트 (PNG 형식)
            asset_type: "TABLE" 또는 "FIGURE"

        Returns:
            추출된 텍스트 문자열. 실패 시 빈 문자열.
        """
        if not image_data:
            return ""
        try:
            # Tesseract 실행 경로 설정 (config에 지정된 경우)
            if config.ocr_tesseract_cmd:
                pytesseract.pytesseract.tesseract_cmd = config.ocr_tesseract_cmd
            # Tesseract 언어 데이터 경로 설정
            if config.ocr_tessdata_dir:
                os.environ["TESSDATA_PREFIX"] = config.ocr_tessdata_dir

            img  = PILImage.open(io.BytesIO(image_data))
            # 표는 psm 6 (균일 블록), 그림은 psm 3 (자동 페이지 분할)
            cfg  = "--psm 6" if asset_type == "TABLE" else "--psm 3"
            lang = getattr(config, "ocr_language", "kor+eng")  # 기본값: 한국어+영어
            return pytesseract.image_to_string(img, lang=lang, config=cfg).strip()
        except Exception as e:
            logger.warning("자산 이미지 OCR 실패 (%s): %s", asset_type, e)
            return ""


# 모듈 전체에서 공유하는 싱글턴 인스턴스
# 매번 새로 생성하지 않고 이 인스턴스를 import해서 사용
asset_pipeline = AssetPipeline()
